# ============================================
# 备课助手 Web 版 - 用户登录版 4.0
# 启动: python web_app.py
# 访问: http://localhost:6000
# ============================================

import re
import os
import sys
import yaml
import csv
import io
import datetime
import random
import sqlite3
import threading
import queue
from pathlib import Path
from functools import wraps
from flask import Flask, render_template, request, jsonify, send_file, Response, redirect, url_for, make_response
from openai import OpenAI
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import time
import json
import uuid
import hashlib

# 导入数据库模块
import database
import supabase_client
import template_filler
from exam_module import register_exam_routes
from visual_module import register_visual_routes

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys._MEIPASS)
    return Path(__file__).parent

def get_data_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys.executable).parent
    return Path(__file__).parent

BASE_DIR = get_base_dir()
DATA_DIR = get_data_dir()

# 读配置
CONFIG_PATH = BASE_DIR / "config.yaml"
cfg = {}
if CONFIG_PATH.exists():
    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

SECRET_KEY = cfg.get("secret_key") or os.environ.get("SECRET_KEY")
if not SECRET_KEY:
    raise ValueError("必须配置 secret_key（在 config.yaml 中设置或设置环境变量 SECRET_KEY）")

template_dir = str(BASE_DIR / "templates")
app = Flask(__name__, template_folder=template_dir)
app.secret_key = SECRET_KEY
app.config['TEMPLATES_AUTO_RELOAD'] = True

API_KEY = os.environ.get("DEEPSEEK_API_KEY") or os.environ.get("OPENAI_API_KEY") or cfg.get("api_key") or ""
BASE_URL = cfg.get("base_url") or os.environ.get("DEEPSEEK_BASE_URL") or os.environ.get("OPENAI_BASE_URL") or "https://api.deepseek.com/v1"
MODEL = cfg.get("model", "deepseek-chat")
TEMP = cfg.get("temperature", 0.7)
MAX_TOKENS = cfg.get("max_tokens", 8192)
MAX_HISTORY = cfg.get("max_history", 6)
OUTPUT_DIR = cfg.get("output_dir", ".")
if not os.path.isabs(OUTPUT_DIR):
    OUTPUT_DIR = os.path.abspath(os.path.join(str(DATA_DIR), OUTPUT_DIR))
else:
    OUTPUT_DIR = os.path.abspath(OUTPUT_DIR)

os.makedirs(OUTPUT_DIR, exist_ok=True)

HISTORY_DIR = os.path.join(OUTPUT_DIR, "历史记录")
os.makedirs(HISTORY_DIR, exist_ok=True)

VISUAL_OUTPUT_DIR = DATA_DIR / "output"
os.makedirs(str(VISUAL_OUTPUT_DIR), exist_ok=True)

SUPABASE_URL = cfg.get("supabase_url") or os.environ.get("SUPABASE_URL")
SUPABASE_ANON_KEY = cfg.get("supabase_anon_key") or os.environ.get("SUPABASE_ANON_KEY")
SUPABASE_SERVICE_KEY = cfg.get("supabase_service_role_key") or os.environ.get("SUPABASE_SERVICE_ROLE_KEY")

if SUPABASE_URL and SUPABASE_ANON_KEY:
    supabase_client.init_supabase(SUPABASE_URL, SUPABASE_ANON_KEY, SUPABASE_SERVICE_KEY)
    if not SUPABASE_SERVICE_KEY:
        print("Warning: supabase_service_role_key not configured. Admin features will not work.")
else:
    print("Warning: Supabase not configured. User authentication may not work.")

QWEN_IMAGE_API_KEY = cfg.get("qwen_image_api_key") or os.environ.get("QWEN_IMAGE_API_KEY") or ""
QWEN_IMAGE_MODEL = cfg.get("qwen_image_model", "qwen-image-plus")

IMAGE_SEARCH_PROVIDER = cfg.get("image_search_provider", "unsplash")
IMAGE_SEARCH_API_KEY = cfg.get("image_search_api_key") or os.environ.get("IMAGE_SEARCH_API_KEY") or ""

ALLOWED_UPLOAD_EXT = {'.md', '.txt', '.docx', '.xlsx', '.csv', '.pdf'}
MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # 20MB
MAX_TEXT_CHARS = 30000  # 单文件最多送3万字符给AI

import hmac
import secrets

# 管理员账号与会话已迁移到数据库（database.admin_users / admin_sessions 表），
# 不再使用进程内存字典 ADMIN_CREDENTIALS 和进程级常量 ADMIN_TOKEN_VALUE。
# 详见 database.py 中的 create_admin_user / verify_admin_password / create_admin_session 等函数。

_CSRF_TOKENS = {}
_CSRF_LOCK = threading.Lock()
_CSRF_EXEMPT = {
    '/login', '/register', '/api/auth/session', '/api/auth/get-csrf-token',
    '/api/chat', '/api/chat/stream', '/api/exam/generate', '/api/exam/export/docx', '/api/exam/export/pdf',
    '/api/visual/generate', '/api/save_blob',
    '/api/lesson_plan/template_upload', '/api/lesson_plan/template_fill', '/api/lesson_plan/template_confirm', '/api/lesson_plan/template_generate',
    '/api/upload', '/api/download', '/api/avatar',
    '/api/profile/avatar', '/api/profile/username',
    '/logout',
    '/admin/login', '/admin/logout',
    '/admin/api/admin/change_password'
}

def _generate_csrf_token(user_id):
    token = secrets.token_hex(32)
    with _CSRF_LOCK:
        _CSRF_TOKENS[token] = {'user_id': user_id, 'created_at': time.time()}
    return token

def _validate_csrf_token(token, user_id):
    if not token:
        return False
    with _CSRF_LOCK:
        entry = _CSRF_TOKENS.get(token)
        if not entry:
            return False
        if entry['user_id'] != user_id:
            return False
        if time.time() - entry['created_at'] > 3600:
            del _CSRF_TOKENS[token]
            return False
        del _CSRF_TOKENS[token]
        return True

def _format_ai_error(e):
    """识别常见 AI API 异常，返回对用户友好的错误提示。

    覆盖：余额不足/配额超限、API Key 无效、超时、限流、连接错误等。
    不向客户端暴露原始异常细节（避免泄露 key/路径），但保留可读的归类提示。
    """
    # openai SDK 异常通常带 status_code 属性
    status_code = getattr(e, 'status_code', None)
    raw = str(e) or ''
    low = raw.lower()

    # 余额不足（智谱 code 1113 / 通用 quota / insufficient balance）
    if ('余额不足' in raw or 'insufficient' in low or 'quota' in low
            or '1113' in raw or 'resource' in low and 'package' in low):
        return 'AI 服务余额不足或配额已用完，请联系管理员充值后重试'

    # API Key 无效 / 鉴权失败
    if status_code == 401 or 'api key' in low or 'invalid key' in low \
            or 'authentication' in low or 'unauthorized' in low:
        return 'API Key 无效或未配置，请联系管理员'

    # 超时
    if 'timeout' in low or 'timed out' in low or status_code == 504 \
            or 'APITimeoutError' in type(e).__name__:
        return 'AI 请求超时，请稍后重试'

    # 限流（非余额类）
    if status_code == 429 or 'rate limit' in low:
        return '请求过于频繁，请稍后重试'

    # 连接错误
    if 'connection' in low or 'conn' in low or 'network' in low:
        return '网络连接异常，请检查网络后重试'

    return 'AI 请求失败，请稍后重试'

def _check_csrf():
    if request.path in _CSRF_EXEMPT:
        return True
    if request.method != 'POST':
        return True
    user = get_current_user()
    if not user:
        return True
    token = request.headers.get('X-CSRF-Token') or request.form.get('_csrf_token')
    if not _validate_csrf_token(token, user['id']):
        return False
    return True


@app.before_request
def csrf_protect():
    if not _check_csrf():
        print(f"[CSRF] Blocked request: path={request.path}, method={request.method}, user={get_current_user() is not None}")
        return jsonify({"error": "请求无效"}), 403

_env_admin_user = os.environ.get('ADMIN_USERNAME', 'LSW')
_env_admin_pass = os.environ.get('ADMIN_PASSWORD', 'LSWYYDS')
# 启动时确保至少有一个超管存在：若 admin_users 表为空则用环境变量创建初始超管。
# 注意：登录账号 _env_admin_user 仅用于首次初始化，运行期不再以此覆盖任何状态。
try:
    _init_super_result = database.init_default_super_admin(_env_admin_user, _env_admin_pass)
    if _init_super_result.get("created"):
        print("=" * 60)
        print("已创建初始超级管理员账号：")
        print(f"登录账号：{_init_super_result['admin']['username']}")
        print("请尽快登录后台修改密码，并删除或修改默认密码。")
        print("=" * 60)
    elif not _init_super_result.get("message", "").startswith("已存在超管"):
        print(f"[startup] 超管初始化: {_init_super_result.get('message')}")
    if not os.environ.get('ADMIN_PASSWORD'):
        print("=" * 60)
        print(f"环境变量默认管理员账号（仅首次启动用于初始化）：")
        print(f"用户名：{_env_admin_user}")
        print(f"密码：{_env_admin_pass}")
        print("如需更换默认账号，请设置环境变量 ADMIN_USERNAME / ADMIN_PASSWORD 后重启。")
        print("运行期账号管理请在后台 /admin 的管理员账号管理界面操作。")
        print("=" * 60)
except Exception as _init_err:
    print(f"[startup] 初始化超管失败: {_init_err}")

TEMPLATE_CACHE = {}
TEMPLATE_CACHE_TTL = 30 * 60  # 30分钟
TEMPLATE_CACHE_LOCK = threading.Lock()


def cleanup_template_cache():
    now = time.time()
    with TEMPLATE_CACHE_LOCK:
        expired = [k for k, v in TEMPLATE_CACHE.items() if now - v.get('timestamp', 0) > TEMPLATE_CACHE_TTL]
        for k in expired:
            del TEMPLATE_CACHE[k]


def safe_get_json():
    try:
        data = request.get_json(silent=True)
        if data is None:
            data = request.form if request.form else None
        return data
    except Exception:
        return None


def get_user_output_dir(user_id=None):
    if user_id:
        user_dir = os.path.join(HISTORY_DIR, f"user_{user_id}")
    else:
        user_dir = os.path.join(HISTORY_DIR, "public")
    os.makedirs(user_dir, exist_ok=True)
    return user_dir

SYSTEM_PROMPT = """你是专业的「AI 备课助手」，服务对象是中小学 / 职业院校的一线教师。你的任务是根据教师的需求，按规范格式生成高质量的教学资源，并**严格用指定标签包裹可保存的文件内容**，方便系统自动解析并下载。

---

## 🎯 核心能力（6 大模块）
1. **PPT 课件生成**：输出结构化 PPT 大纲，系统会自动解析为 `.pptx` 课件
2. **PPT 讲解纲要**：为每页幻灯片生成详细的教师讲解稿（逐页讲稿 + 互动设计）
3. **教案生成**：45 分钟标准格式教案（三维目标 / 教学过程 / 板书 / 反思）
4. **习题生成**：基础 / 提高 / 拓展 三级难度分层习题，附答案解析
5. **作业分析**：作业批改数据分析 + 错因归因 + 辅导建议 + 教学调整策略
6. **复习提纲**：期末 / 单元复习资料，知识脉络 + 核心考点 + 典型例题

---

## 🚫 最重要的输出规则（不遵守会导致文件无法下载）
**必须用 `[文件:xxx.md]  …内容…  [/文件]` 或 `[PPT:xxx.pptx]  …大纲…  [/PPT]` 标签把生成结果包裹起来**，系统会自动保存为可下载的文件。
- 如果用户要的是「教案 / 习题 / 作业分析 / 讲解纲要 / 复习提纲」 → 用 `[文件:xxx.md]`
- 如果用户要的是「PPT 课件 / 幻灯片」 → 用 `[PPT:xxx.pptx]`
- 若用户同时要多份内容（如「同时生成教案和PPT」），可同时出现多个标签块
- 纯文字的简短说明 / 寒暄可以写在标签外，但**核心生成内容必须全部放在标签块内**

✅ 正确示例：
```
好的，为您生成《背影》的教案如下：
[文件:初中语文《背影》教案.md]
# 初中语文《背影》教案

## 一、三维教学目标
- **知识与技能**：……
- **过程与方法**：……
- **情感态度与价值观**：……

## 二、教学重难点
……（以下内容全部写在标签内）……

## 六、教学反思栏
……
[/文件]
```

❌ 错误示例（严禁！会导致用户无法下载文件）：
- 直接输出 Markdown，不写任何标签包裹
- 标签名写错，如 `[教案:xxx]`、`[下载:xxx]`
- 文件内容一半在标签里，一半在标签外

---

## 📖 模块 1：教案生成 标准格式（`[文件:xxx教案.md]`）
结构必须完整，标题层级严格按下面顺序：

```
# 《课题名》教案（学段 + 学科）

## 一、基本信息
- 学科：xxx
- 学段：xxx（小学/初中/高中/中职）
- 年级：xxx
- 课时：1课时（45分钟）
- 课型：新授课 / 复习课 / 讲评课 / 实验课 / 活动课

## 二、三维教学目标 / 核心素养目标
- 知识与技能：……
- 过程与方法：……
- 情感态度与价值观：……

## 三、教学重难点
- **教学重点**：……
- **教学难点**：……

## 四、学情分析与教学方法
- 学情分析：……
- 教学方法：讲授法 / 讨论法 / 实验法 / 情境教学法 / 任务驱动法

## 五、教学准备
- 教具准备：多媒体课件、实验器材、学案……
- 学生准备：……

## 六、详细教学过程（精确到分钟）
### 1. 导入新课（约 5 分钟）
……教师活动……
……学生活动……

### 2. 新知讲授（约 15 分钟）
……分步骤讲解……

### 3. 课堂练习 / 小组讨论（约 12 分钟）
……

### 4. 课堂小结（约 5 分钟）
……由学生总结 + 教师补充……

### 5. 作业布置（约 3 分钟）
- 基础作业：……
- 拓展作业：……

### 6. 板书设计（约 5 分钟同步板书）
主板书 | 副板书
-------|-------
……    | ……

## 七、板书设计
- 主板书：……（结构化呈现）
- 副板书：……

## 八、分层作业设计
- ✅ **基础层（必做）**：……
- 📈 **提高层（选做）**：……
- 🚀 **拓展层（挑战）**：……

## 九、教学反思
- 本节课亮点：……
- 待改进点：……
- 课堂生成问题记录（留白教师填写）：
```

---

## 📝 模块 2：习题生成 标准格式（`[文件:xxx习题.md]`）
```
# 《课题 / 知识点》分层练习题

## 一、基本信息
- 学段学科：xxx
- 考查知识点：xxx
- 题目总数：xx 道
- 建议用时：xx 分钟
- 难度分布：基础 60% / 提高 30% / 拓展 10%

---

## 二、基础巩固题（约 60%）
### 一、选择题（每题 3 分，共 xx 分）
1. 题目内容……
   A. xxx   B. xxx   C. xxx   D. xxx
2. ……

### 二、填空题（每题 2 分，共 xx 分）
1. _________________
2. _________________

---

## 三、能力提高题（约 30%）
### 三、解答题 / 计算题（每题 xx 分，共 xx 分）
1. 题目内容……

---

## 四、拓展探究题（约 10%）
### 四、综合应用题 / 探究题
1. 题目内容……

---

## 五、参考答案与详细解析
1. 选择题 1：答案 B
   - 解析：……
   - 考点：……
   - 易错分析：……
2. ……
```

---

## 📊 模块 3：作业分析报告（`[文件:xxx作业分析.md]`）
```
# 《作业名称》学情分析报告

## 一、整体概况
- 学科：xxx
- 班级：xxx（共 xx 人）
- 应交：xx 份，实交：xx 份，上交率：xx%
- 平均分：xx 分 ｜ 最高分：xx ｜ 最低分：xx ｜ 及格率：xx% ｜ 优秀率：xx%

### 分数段分布
| 分数段    | 人数 | 占比  |
|-----------|------|-------|
| 90-100 分 | x 人 | xx%  |
| 80-89 分  | x 人 | xx%  |
| 70-79 分  | x 人 | xx%  |
| 60-69 分  | x 人 | xx%  |
| 60 分以下 | x 人 | xx%  |

## 二、知识点掌握情况雷达
| 知识点         | 平均得分率 | 掌握等级 |
|----------------|------------|----------|
| 知识点 A       | xx%        | 熟练 / 一般 / 薄弱 |
| 知识点 B       | xx%        | ……     |

## 三、高频错题 TOP 5 + 归因分析
| 排名 | 题号 | 知识点 | 得分率 | 主要错解 | 错因归类 |
|------|------|--------|--------|----------|----------|
| 1    | 第x题 | xxx | xx% | …… | ⚠️ 概念不清 / 审题错误 / 计算失误 / 方法未掌握 |
| 2    | …… | …… | …… | …… | …… |

### 典型错解展示（第 x 题）
> 错误解法示例：……
> ✅ 正确解法：……
> 🔍 错因剖析：……

## 四、分层辅导建议
### 🔴 学困生（xx 人，60 分以下）
- 薄弱点：……
- 辅导策略：① …… ② …… ③ ……
- 补充练习题：……

### 🟡 中等生（xx 人，60-84 分）
- 提升点：……
- 辅导策略：……
- 强化练习：……

### 🟢 优等生（xx 人，85 分以上）
- 拓展方向：……
- 挑战性题目：……

## 五、下节课教学调整建议
- 需要补讲的知识点：……（建议用时 xx 分钟）
- 课堂讲评顺序建议：先讲第 x/x/x 题（得分率最低）
- 教学方法调整：……
- 与家长沟通要点（家校共育）：……

## 六、讲评课时分配建议
| 环节            | 时间  | 内容说明 |
|-----------------|-------|----------|
| 整体情况通报    | 3 分钟 | 成绩分布 + 表扬优秀进步学生 |
| 高频错题精讲    | 20 分钟 | TOP 5 错因分析 + 变式训练 |
| 小组互助订正    | 10 分钟 | 学困生结对，同伴讲解 |
| 针对性补充练习  | 8 分钟  | 相似题型，当堂巩固 |
| 总结 + 二次过关 | 4 分钟  | 关键方法总结，小测过关 |
```

---

## 🎤 模块 4：PPT 讲解纲要（逐页讲稿，`[文件:xxx讲解纲要.md]`）
```
# 《课件名》PPT 逐页讲解纲要

## 使用说明
- 总页数：xx 页
- 建议总时长：45 分钟
- 适用对象：xxx 年级学生

---

## 第 1 页：封面（约 1 分钟）
🎤 **教师讲解词**：
同学们好，今天我们一起来学习《……》。在上课之前请大家看屏幕上的这张图片 / 这个问题，有没有同学能说一说……
❓ **互动提问**（可选）：……
⭐ **语气提示**：语速稍慢，吸引注意力

---

## 第 2 页：教学目标（约 2 分钟）
🎤 **教师讲解词**：
本节课我们要达成三个学习目标：第一，……；第二，……；第三，……
⭐ **强调**：第 2 条是重点，请用红色笔在学案上画出来
⏰ 本页用时：2 分钟

---

## 第 3 页：xxx（按 PPT 每页依次写）
🎤 **教师讲解词**：……
❓ **互动提问**：……（预留学生作答的留白，可写学生可能的回答）
⚠️ **易错提醒**：……
⏰ 本页用时：xx 分钟

---
（每页 PPT 都按上面格式写一节）
---

## 最后一页：课堂小结 + 作业布置（约 3 分钟）
🎤 **教师讲解词**：
好，本节课我们学习了……，主要内容可以用三句话记住：①……②……③……
下课后请大家完成：……（作业内容）
下课，同学们再见！
```

---

## 🤖 模块 5：复习提纲（`[文件:xxx复习提纲.md]`）
```
# 《课程名》期末 / 单元复习提纲

## 📚 一、复习内容概览
- 章节范围：第 x 章 — 第 x 章
- 建议复习用时：xx 课时
- 重要程度：★★★★★（5 星为必考）

## 🧠 二、核心知识体系（思维导图式）
### 模块一：xxx
- 核心概念 1：……
  - 定义：……
  - 关键词：……
- 核心概念 2：……

### 模块二：xxx
……

## 📐 三、重点公式 / 定理 / 结论速记
| 名称 | 公式 / 结论 | 适用条件 | 考频 |
|------|-------------|----------|------|
| xxx  | ……         | ……       | ⭐⭐⭐ |

## 📝 四、典型例题精讲
### 例题 1（★★ 基础题）
> 题目：……
> 解题步骤：
> ① ……
> ② ……
> 💡 点拨：……

### 例题 2（★★★★ 高频考点）
……

## ⚠️ 五、易错点警示 TOP 10
1. ❌ 错误理解：…… → ✅ 正确：……
2. ❌ ……

## 🎯 六、考点预测 + 分值分布
| 考点         | 预测分值 | 题型       | 难度 |
|--------------|----------|------------|------|
| xxx          | 8-12 分  | 选择 + 解答 | ★★★ |

## ✅ 七、自我检测（附答案）
……（10 道精选小题，覆盖全部考点）
```

---

## 📊 模块 6：PPT 课件生成（必须用 `[PPT:xxx.pptx]` … `[/PPT]` 包裹）
**结构严格遵守以下格式，系统会自动解析成 `.pptx` 文件：**

```
[PPT:荷塘月色_语文课件.pptx]
# 《荷塘月色》语文课件
朱自清散文 · 高中语文必修上册

## 第1页：情境导入
- 展示荷花池夜色图片
- 提问：同学们记忆中关于"月"和"荷"的诗句有哪些？
- 引出作者：朱自清
- 板书课题

## 第2页：学习目标
- 📖 知识与技能：掌握重点词语，理解关键语句含义
- 🧠 过程与方法：通过朗读体会文章的语言美和意境美
- ❤️ 情感态度：体会作者情感，感受文学作品的感染力

## 第3页：作者介绍
- 朱自清（1898-1948），原名自华，字佩弦
- 现代著名散文家、诗人、学者
- 代表作品：《背影》《春》《匆匆》《荷塘月色》
- 散文风格：语言洗练，文笔清丽，情感真挚

## 第4页：写作背景
……

## 第N页：课堂小结 + 作业布置
……
[/PPT]
```

**PPT 格式要求（务必遵守）：**
1. 第一行 `# 主标题` 是封面主标题，第二行（无特殊符号）是副标题
2. 每页幻灯片用 `## 第N页：页面标题` 开头
3. 每页下面的内容用 `- 要点` 列（最多 7 条，每条简短）
4. 页数控制在 8-20 页之间，结构完整
5. **绝对不要**把多页内容塞进一个 `##` 里

---

## 🎨 其它输出规范
- 全文中文表达，专业简洁
- Markdown 层级清晰：# 一级 → ## 二级 → ### 三级
- 重点用 **粗体**、表格整理数据、列表梳理条目
- 编号必须连续，避免乱码字符

## 📎 上传文件智能识别规则（最重要！）
用户消息中如果出现 `【用户上传文件：文件名（类型）】` 标记，说明用户上传了该文件的全文内容，你必须按以下逻辑处理：

### 情况 A：文件类型 = "教案/教学设计模板"（文件名或内容含「教案」「教学设计」「学案」「模板」「____」「[填空]」「__」等占位符）
- 你的任务 = **按模板中所有的空白 / 占位符 / 待填写项，自动填充完整的教学内容**
- 规则：
  1. 严格**保留模板原有的章节结构、编号、表格结构、填空占位符的位置**
  2. 模板里写着「____」「___」「...」「待填写」「填写」「留白」「XXX」「[     ]」的位置，全部填上具体、详实、符合学科逻辑的内容
  3. 若用户在上传文件之外还写了需求（如"请按人教版高二语文"），需优先按该需求填充内容
  4. 没有的学科/年级信息时，按模板里能推断的学段和学科默认选一个最通用的
  5. **生成结果必须用 `[文件:xxx_已填充.md]` 完整包裹起来**，把整个模板+填充后的全部内容放进去，方便用户下载

### 情况 B：文件类型 = "学生作业/成绩/批改记录"（文件名或内容含「作业」「成绩」「得分」「分数」「错题」「批改」「考勤」「班级」「学生」「姓名」等）
- 你的任务 = **按「模块 5 作业分析」标准格式生成一份完整的学情分析报告**
- 报告内容必须包含 6 部分：
  1. **整体概况**（班级人数、均分、及格率、优秀率、最高分/最低分、分数段分布）
  2. **错题归因分析**（哪题错最多？错因归类：概念不清 / 计算失误 / 审题错误 / 知识点没掌握 / 书写规范）
  3. **知识点掌握雷达**（按知识点列出 掌握率 = 做对人数/总人数）
  4. **学生分层画像**（优生/中等/待进步各占多少 %，每一层的典型问题）
  5. **针对性辅导建议**（全班教学调整 + 分层布置作业 + 个别辅导名单）
  6. **后续教学调整策略**（下一节课怎么改、是否加小测、是否安排讲评课）
- **生成结果必须用 `[文件:xxx_作业分析报告.md]` 完整包裹**

### 情况 C：其他类型文件（如参考资料、课文原文、讲义等）
- 先明确告诉用户："已读取到您上传的《xxx》（共 x 字 / x 行）"
- 再结合用户的具体问题给出对应的回答；如果用户没有明确问题，自动判断：内容像课文原文 → 建议生成《xxx》教案/PPT/习题；内容像复习资料 → 建议生成复习提纲

最后再次强调：**所有生成内容必须放在对应的标签块中，否则用户点击"下载"时将没有任何文件！**"""

# 存会话（内存缓存 + 数据库持久化）
# 用 user_id:sesssion_id 作为 key，确保用户隔离
sessions: dict[str, list[dict]] = {}
# 会话字典并发访问锁，避免多请求/多线程同时修改导致历史错乱
sessions_lock = threading.Lock()

def _session_key(user_id, session_id):
    return f"{user_id}:{session_id}"

def get_current_user():
    # 管理员：从服务端 admin_sessions 表反查会话，cookie 仅携带 token
    admin_session_token = request.cookies.get('admin_session')
    if admin_session_token:
        session = database.get_admin_session(admin_session_token)
        if session and session.get('is_active'):
            admin = database.get_admin_by_id(session['admin_id'])
            if admin and admin['is_active']:
                display_name = admin['display_name'] or admin['username']
                return {
                    'id': admin['id'],
                    'username': display_name,
                    'account': admin['username'],
                    'email': 'admin@example.com',
                    'is_admin': True,
                    'is_super': admin['is_super'],
                    'avatar_url': admin['avatar_url']
                }

    token = request.cookies.get('session_token')
    if token:
        user = database.get_user_by_token(token)
        if user:
            user['is_admin'] = False
        return user
    return None

def trim_history(messages):
    if not messages or len(messages) <= MAX_HISTORY + 1:
        return list(messages)
    system_msg = messages[0]
    recent = messages[-(MAX_HISTORY):]
    return [system_msg] + list(recent)

def markdown_to_word(markdown_text: str, title: str = "教案") -> BytesIO:
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)

    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title)
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(31, 78, 121)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    lines = markdown_text.split('\n')
    in_code_block = False
    code_content = []

    for line in lines:
        if line.strip().startswith('```'):
            if in_code_block:
                if code_content:
                    code_para = doc.add_paragraph()
                    code_run = code_para.add_run('\n'.join(code_content))
                    code_run.font.name = 'Courier New'
                    code_run.font.size = Pt(9)
                    code_run.font.color.rgb = RGBColor(128, 128, 128)
                    code_para.paragraph_format.left_indent = Pt(20)
                code_content = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_content.append(line)
            continue

        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('#### '):
            doc.add_heading(line[5:], level=4)
        elif '**' in line:
            count = line.count('**')
            if count % 2 != 0:
                # 奇数个 **，无法配对，按普通文本处理
                doc.add_paragraph(line)
            else:
                parts = line.split('**')
                para = doc.add_paragraph()
                for i, part in enumerate(parts):
                    if i % 2 == 1:
                        run = para.add_run(part)
                        run.bold = True
                    else:
                        para.add_run(part)
        elif line.strip().startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.strip().startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif re.match(r'^\d+\.\s', line.strip()):
            doc.add_paragraph(re.sub(r'^\d+\.\s', '', line.strip()), style='List Number')
        elif '|' in line and line.strip():
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if len(cells) >= 2:
                if not all(re.match(r'^-+$', cell) for cell in cells):
                    para = doc.add_paragraph()
                    para.add_run(line)
        elif line.strip().startswith('>'):
            para = doc.add_paragraph(line[1:].strip())
            para.paragraph_format.left_indent = Pt(20)
            para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
        elif line.strip() == '---':
            doc.add_paragraph('_' * 50)
        elif line.strip():
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def _strip_md(text):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'^#{1,6}\s*', '', text)
    text = text.replace('•', '').replace('◦', '').replace('▪', '')
    text = re.sub(r'^\s*[-*]\s+', '', text)
    text = re.sub(r'^\s*\d+[.、)]\s*', '', text)
    return text.strip()


def _detect_header_level(s):
    m = re.match(r'^(#{1,6})\s+(.+)$', s)
    if m:
        return len(m.group(1)), m.group(2).strip()
    return 0, s


def _is_bullet(s):
    return bool(re.match(r'^\s*[-*•◦▪]\s+', s))


def _is_numbered(s):
    return bool(re.match(r'^\s*\d+[.、)]\s*', s))


def _indent_level(raw):
    stripped = raw.lstrip(' \t')
    n = len(raw) - len(stripped)
    if n == 0:
        return 0
    if n <= 3:
        return 1
    if n <= 6:
        return 2
    return 3


def parse_ppt_outline(markdown_text):
    lines = markdown_text.split('\n')
    title_main = ""
    title_sub = ""
    raw_sections = []

    header_levels = set()
    for raw in lines:
        s = raw.strip()
        if not s:
            continue
        lvl, _ = _detect_header_level(s)
        if lvl > 0:
            header_levels.add(lvl)

    if header_levels:
        if 1 in header_levels:
            cover_level = 1
            slide_levels = {l for l in header_levels if l > 1}
        else:
            cover_level = 0
            slide_levels = header_levels
    else:
        cover_level = 0
        slide_levels = set()

    min_slide = min(slide_levels) if slide_levels else 0

    cur_title = None
    cur_items = []

    def flush():
        nonlocal cur_title, cur_items
        if cur_title is not None:
            raw_sections.append((cur_title, cur_items))
        cur_title = None
        cur_items = []

    for raw in lines:
        if not raw.strip():
            continue
        s = raw.strip()
        lvl, content = _detect_header_level(s)

        if lvl == cover_level and cover_level > 0:
            flush()
            title_main = content
        elif lvl in slide_levels:
            flush()
            t = re.sub(r'^第\s*\d+\s*页\s*[：:]\s*', '', content)
            cur_title = t
        elif lvl > 0 and min_slide > 0:
            rel = max(0, lvl - min_slide)
            cur_items.append((min(rel, 3), content))
        elif _is_bullet(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*[-*•◦▪]\s+', '', s)
            cur_items.append((ind, content))
        elif _is_numbered(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*\d+[.、)]\s*', '', s)
            cur_items.append((ind, content))
        else:
            if title_main and not title_sub and not raw_sections and cur_title is None:
                title_sub = s
            else:
                cur_items.append((_indent_level(raw), s))
    flush()

    merged = []
    for stitle, items in raw_sections:
        if not items and merged:
            prev_t, prev_items = merged[-1]
            prev_items.append((0, stitle))
            merged[-1] = (prev_t, prev_items)
        elif not items and not merged:
            # 首个 section 无内容且无前置可合并项，跳过避免出现空白页
            continue
        else:
            merged.append((stitle, items))
    raw_sections = merged

    if not raw_sections and title_main:
        all_items = []
        for raw in lines:
            s = raw.strip()
            if not s or _detect_header_level(s)[0] == cover_level:
                continue
            if _is_bullet(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*[-*•◦▪]\s+', '', s)))
            elif _is_numbered(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*\d+[.、)]\s*', '', s)))
            elif _detect_header_level(s)[0] > 0:
                _, content = _detect_header_level(s)
                all_items.append((0, content))
            else:
                all_items.append((_indent_level(raw), s))
        if all_items:
            raw_sections = [("内容概览", all_items)]

    return title_main, title_sub, raw_sections


def _split_overflow(items, max_items=7):
    if len(items) <= max_items:
        return [items]
    result = []
    cur = []
    for lvl, text in items:
        # 在顶层（lvl==0）边界分页，或当子级项目累计已达上限时强制分页
        if len(cur) >= max_items and (lvl == 0 or len(cur) >= max_items + 3):
            result.append(cur)
            cur = []
        cur.append((lvl, text))
    if cur:
        result.append(cur)
    return result


class PptTheme:
    PRIMARY = PptRGBColor(0x1E, 0x40, 0xAF)
    DARK = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT_SEC = PptRGBColor(0x4B, 0x55, 0x63)
    MUTED = PptRGBColor(0x9C, 0xA3, 0xAF)
    WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = PptRGBColor(0x3B, 0x82, 0xF6)
    FONT = 'Microsoft YaHei'


class PptBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.sw = self.prs.slide_width
        self.sh = self.prs.slide_height
        self.m = Inches(0.8)

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _tbox(self, slide, left, top, w, h, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def _run(self, para, text, size, color, bold=False):
        run = para.add_run()
        run.text = text
        run.font.size = PptPt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = PptTheme.FONT
        return run

    def _rect(self, slide, left, top, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _page_num(self, slide, num, total):
        tf = self._tbox(slide, self.sw - Inches(1.6), self.sh - Inches(0.45),
                        Inches(1.4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        self._run(p, f"{num} / {total}", 11, PptTheme.MUTED)

    def _footer(self, slide):
        tf = self._tbox(slide, self.m, self.sh - Inches(0.45),
                        Inches(4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        self._run(p, "AI 备课助手", 10, PptTheme.MUTED)

    def _font_size(self, items):
        n = len(items)
        chars = sum(len(_strip_md(t)) for _, t in items)
        max_lvl = max((lvl for lvl, _ in items), default=0)
        if n <= 4 and chars <= 120:
            base = 30
        elif n <= 6 and chars <= 250:
            base = 26
        elif n <= 8 and chars <= 400:
            base = 22
        elif n <= 12 and chars <= 600:
            base = 18
        else:
            base = 15
        if max_lvl >= 2 and base > 22:
            base = 22
        return base

    def cover(self, title, subtitle, speaker='', date_str=''):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(2.8), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.1), self.sw - self.m * 2,
                        Inches(2), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        t = _strip_md(title) or "课件"
        size = 46 if len(t) <= 18 else (36 if len(t) <= 36 else 28)
        self._run(p, t, size, PptTheme.WHITE, bold=True)

        if subtitle:
            tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                            Inches(0.8), MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            self._run(p, _strip_md(subtitle), 22, PptTheme.MUTED)

        footer = '  |  '.join(filter(None, [speaker, date_str]))
        if footer:
            tf = self._tbox(slide, self.m, self.sh - Inches(0.8),
                            self.sw - self.m * 2, Inches(0.4), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            self._run(p, footer, 14, PptTheme.MUTED)

    def content(self, title, items, page_num, total, dark=False):
        slide = self._blank()

        bg = PptTheme.DARK if dark else PptTheme.WHITE
        tc = PptTheme.WHITE if dark else PptTheme.TEXT
        sc = PptTheme.MUTED if dark else PptTheme.TEXT_SEC
        bc = PptTheme.ACCENT if dark else PptTheme.PRIMARY

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        tf = self._tbox(slide, self.m, Inches(0.5), self.sw - self.m * 2,
                        Inches(1), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        clean_t = _strip_md(title)
        ts = 36 if len(clean_t) <= 20 else (28 if len(clean_t) <= 40 else 22)
        self._run(p, clean_t, ts, tc, bold=True)

        self._rect(slide, self.m, Inches(1.5), Inches(2), Inches(0.06), bc)

        base = self._font_size(items)

        BULLETS = ["●", "○", "▪", "·"]
        INDENT_EMU = [0, 457200, 914400, 1371600]
        SIZE_DEC = [0, 3, 5, 6]

        valid_items = [(min(lvl, 3), _strip_md(text)) for lvl, text in items
                       if _strip_md(text)]
        n = len(valid_items)
        est_line_h = (base * 1.3 + 8) / 72.0
        est_height = Inches(min(n * est_line_h + 0.4, 4.6))

        avail_top = Inches(1.85)
        avail_bottom = self.sh - Inches(0.7)
        avail_h = avail_bottom - avail_top
        if est_height < avail_h:
            content_top = avail_top + (avail_h - est_height) / 2
            content_h = est_height
            anchor = MSO_ANCHOR.MIDDLE
        else:
            content_top = avail_top
            content_h = avail_h
            anchor = MSO_ANCHOR.TOP

        tf = self._tbox(slide, self.m, content_top, self.sw - self.m * 2,
                        content_h, anchor)

        first = True
        for lvl, clean in valid_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PptPt(6)
            p.line_spacing = 1.3

            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(INDENT_EMU[lvl]))
            pPr.set('indent', str(-228600))

            bchar = BULLETS[lvl]
            bsize = base - SIZE_DEC[lvl]
            tsize = base - SIZE_DEC[lvl]
            tcolor = tc if lvl == 0 else sc
            self._run(p, f"{bchar}  ", bsize, bc, bold=(lvl == 0))
            self._run(p, clean, tsize, tcolor, bold=False)

        self._page_num(slide, page_num, total)
        if not dark:
            self._footer(slide)

    def ending(self):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(3.0), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.3), self.sw - self.m * 2,
                        Inches(1.5), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "谢谢观看", 48, PptTheme.WHITE, bold=True)

        tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                        Inches(0.6), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "AI 备课助手 · 自动生成", 16, PptTheme.MUTED)

    def build(self, title, subtitle, content_slides, speaker='', date_str=''):
        self.cover(title, subtitle, speaker, date_str)

        expanded = []
        for stitle, items in content_slides:
            chunks = _split_overflow(items, max_items=7)
            for ci, chunk in enumerate(chunks):
                ct = stitle if ci == 0 else f"{stitle}（续{ci}）"
                expanded.append((ct, chunk))

        # 总页数 = 封面(1) + 内容页数 + 结尾页(1)
        total = len(expanded) + 2
        for i, (stitle, items) in enumerate(expanded):
            is_summary = any(k in stitle for k in
                             ['总结', '小结', '结语', '谢谢', '感谢', 'CTA', '行动'])
            self.content(stitle, items, i + 2, total, dark=is_summary)

        if not expanded:
            self.content("内容概览", [(0, "暂无具体内容")], 2, 3, dark=False)

        self.ending()

        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf


def _parse_kv(items):
    kv = {}
    for lvl, text in items:
        clean = _strip_md(text).strip()
        if not clean:
            continue
        # 优先匹配中文冒号 / 英文冒号（强键值分隔符）
        for sep in ['：', ':']:
            if sep in clean:
                parts = clean.split(sep, 1)
                if len(parts) == 2:
                    key = parts[0].strip()
                    val = parts[1].strip()
                    if key and val:
                        kv[key] = val
                    break
        else:
            # 仅当 key 为纯中文/字母标签（不含数字+空格组合，避免 "日期 2024" 被误当 key）
            m = re.match(r'^([\u4e00-\u9fa5A-Za-z]{1,8})\s*[—-]\s*(.+)$', clean)
            if m:
                key = m.group(1).strip()
                val = m.group(2).strip()
                if key and val:
                    kv[key] = val
    return kv


def markdown_to_ppt(markdown_text, title="课件"):
    title_main, title_sub, slides = parse_ppt_outline(markdown_text)
    if not title_main:
        title_main = title
    if not slides:
        # 兜底：将所有非空行作为单个 section 的要点，并通过 _split_overflow 自动分页
        all_items = [(0, t) for t in markdown_text.split('\n') if t.strip()]
        slides = [(title_main or title, all_items)]

    cover_data = {}
    content_slides = []
    for stitle, items in slides:
        if stitle == '封面' or stitle.startswith('封面'):
            cover_data = _parse_kv(items)
        else:
            content_slides.append((stitle, items))

    if cover_data:
        c_title = _strip_md(cover_data.get('主标题', title_main))
        c_sub = _strip_md(cover_data.get('副标题', title_sub))
        speaker = _strip_md(cover_data.get('演讲人', ''))
        date_str = _strip_md(cover_data.get('日期', time.strftime('%Y年%m月%d日')))
    else:
        c_title = _strip_md(title_main) if title_main else title
        c_sub = _strip_md(title_sub) if title_sub else ''
        speaker = ''
        date_str = time.strftime('%Y年%m月%d日')

    builder = PptBuilder()
    return builder.build(c_title, c_sub, content_slides, speaker, date_str)


def generate_ppt_files(reply_text, user_id=None):
    base_dir = get_user_output_dir(user_id)
    ppt_blocks = re.findall(r"\[PPT:\s*([^\]]+)\](.*?)\[/PPT\]", reply_text, re.DOTALL)
    saved_ppts = []
    failed_ppts = []
    for filename, outline in ppt_blocks:
        filename = filename.strip()
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"
        safe_name = _sanitize_filename(filename)
        if not safe_name.lower().endswith(".pptx"):
            safe_name += ".pptx"
        filepath = os.path.join(base_dir, safe_name)
        if not filepath.startswith(os.path.abspath(base_dir)):
            continue
        ppt_buffer = None
        try:
            title = os.path.splitext(safe_name)[0]
            ppt_buffer = markdown_to_ppt(outline, title=title)
            with open(filepath, "wb") as f:
                f.write(ppt_buffer.getvalue())
            saved_ppts.append(safe_name)
        except Exception as e:
            print(f"[PPT] 生成失败 {safe_name}")
            failed_ppts.append({"filename": safe_name, "error": "生成失败"})
        finally:
            if ppt_buffer is not None:
                try:
                    ppt_buffer.close()
                except Exception:
                    pass
            try:
                if os.path.exists(filepath) and os.path.getsize(filepath) == 0:
                    os.remove(filepath)
            except Exception:
                pass
    display = reply_text
    if ppt_blocks:
        display = re.sub(r"\[PPT:\s*[^\]]+\](.*?)\[/PPT\]", "", display, flags=re.DOTALL)
    if failed_ppts:
        notice_lines = ["\n\n> ⚠️ 以下 PPT 生成失败："]
        for fp in failed_ppts:
            notice_lines.append(f"> - {fp['filename']}")
        display = display + "".join(notice_lines)
    return saved_ppts, display


def auto_wrap_and_save_fallback(user_message: str, ai_reply: str, user_id):
    """兜底：AI 忘记写 [文件:] / [PPT:] 标签时，自动检测内容类型并保存文件，保证一定能下载。
    返回 (extra_files, extra_ppts, display)"""
    import os as _os, re as _re, time as _time
    if not ai_reply or not ai_reply.strip():
        return [], [], ai_reply
    has_file = bool(_re.search(r"\[文件:\s*[^\]]+\]", ai_reply))
    has_ppt = bool(_re.search(r"\[PPT:\s*[^\]]+\]", ai_reply))
    if has_file or has_ppt:
        return [], [], ai_reply

    # 修复 L1：已删除未使用的 text / lower 死代码（下方仅使用 reply_text / lower_reply）

    safe_title = _re.sub(r'[<>:"/\\|?*]', '', (user_message or "AI生成内容").strip())[:20] or "AI生成内容"
    ts = _time.strftime("%Y%m%d_%H%M%S")

    # 仅根据 AI 回复内容判断文件类型，避免用户消息干扰
    reply_text = (ai_reply or "")
    lower_reply = reply_text.lower()

    if ("ppt" in lower_reply) or ("课件" in reply_text) or ("幻灯片" in reply_text):
        filename = f"课件_{safe_title}_{ts}.pptx"
        wrapped = f"[PPT:{filename}]\n{ai_reply.strip()}\n[/PPT]"
        saved_ppts, display = generate_ppt_files(wrapped, user_id)
        return [], saved_ppts, display

    if "作业" in reply_text and ("分析" in reply_text or "批改" in reply_text or "错因" in reply_text or "讲评" in reply_text):
        fn = f"作业分析_{safe_title}_{ts}.md"
    elif "教案" in reply_text or "教学设计" in reply_text or "说课稿" in reply_text or "教学过程" in reply_text:
        fn = f"教案_{safe_title}_{ts}.md"
    elif "习题" in reply_text or "练习" in reply_text or "试卷" in reply_text or "测试题" in reply_text or ("题" in reply_text and "答案" in reply_text):
        fn = f"习题_{safe_title}_{ts}.md"
    elif ("讲解" in reply_text and ("纲要" in reply_text or "讲稿" in reply_text or "逐页" in reply_text)) or "说课稿" in reply_text:
        fn = f"讲解纲要_{safe_title}_{ts}.md"
    elif "复习" in reply_text or "提纲" in reply_text or "知识点总结" in reply_text or "知识体系" in reply_text:
        fn = f"复习提纲_{safe_title}_{ts}.md"
    else:
        fn = f"备课资料_{safe_title}_{ts}.md"

    base_dir = get_user_output_dir(user_id)
    filepath = _os.path.join(base_dir, fn)
    content = ai_reply.strip()
    if content.startswith("```"):
        content = _re.sub(r"^```\w*\s*\n", "", content)
        content = _re.sub(r"\n```\s*$", "", content)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    return [fn], [], ai_reply


def _sanitize_filename(filename):
    filename = str(filename).strip()
    filename = re.sub(r'[<>:"/\\|?*]', '', filename)
    filename = re.sub(r'\.\.+', '.', filename)
    filename = filename.strip('.')
    win_reserved = {'CON', 'PRN', 'AUX', 'NUL', 'COM1', 'COM2', 'COM3', 'COM4', 'COM5', 'COM6', 'COM7', 'COM8', 'COM9', 'LPT1', 'LPT2', 'LPT3', 'LPT4', 'LPT5', 'LPT6', 'LPT7', 'LPT8', 'LPT9'}
    if filename.upper() in win_reserved:
        filename = filename + "_file"
    if len(filename) > 200:
        filename = filename[:200]
    return filename or "unnamed"


def extract_and_save_all(user_message, reply_text, user_id):
    """统一处理 AI 回复：解析文件块 + 生成 PPT + 兜底保存，返回(saved_files, saved_ppts, display_for_chat)"""
    user_output_dir = get_user_output_dir(user_id)

    saved_files = []
    files = re.findall(r"\[文件:\s*([^\]]+)\](.*?)\[/文件\]", reply_text, re.DOTALL)
    for filename, content in files:
        safe_name = _sanitize_filename(filename)
        filepath = os.path.join(user_output_dir, safe_name)
        if not filepath.startswith(os.path.abspath(user_output_dir)):
            continue
        clean = content.strip()
        if clean.startswith("```"):
            clean = re.sub(r"^```\w*\s*\n", "", clean)
            clean = re.sub(r"\n```\s*$", "", clean)
        try:
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(clean)
            saved_files.append(safe_name)
        except Exception:
            pass

    display = reply_text
    if files:
        display = re.sub(r"\[文件:.*?\[/文件\]", "", reply_text, flags=re.DOTALL).strip()

    saved_ppts, display = generate_ppt_files(display, user_id)

    if not saved_files and not saved_ppts:
        extra_fs, extra_ppts, display = auto_wrap_and_save_fallback(
            user_message, reply_text, user_id
        )
        saved_files.extend(extra_fs)
        saved_ppts.extend(extra_ppts)

    return saved_files, saved_ppts, display


@app.route("/health")
def health_check():
    return {"status": "ok"}

@app.route("/")
def index():
    user = get_current_user()
    if not user:
        return redirect(url_for('login'))
    return render_template("index.html", model=MODEL, username=user['username'])

_login_attempts = {}
_login_lock = threading.Lock()

def _get_client_ip():
    # 修复 S7：使用 ipaddress 模块完整判断私有/回环地址，仅信任非私有的 X-Forwarded-For 首个 IP，
    # 防止客户端伪造 X-Forwarded-For 绕过限流
    import ipaddress
    if request:
        forwarded = request.headers.get('X-Forwarded-For')
        if forwarded:
            first_ip = forwarded.split(',')[0].strip()
            try:
                ip_obj = ipaddress.ip_address(first_ip)
                if not ip_obj.is_private:
                    return first_ip
            except ValueError:
                pass
        return request.remote_addr or '127.0.0.1'
    return '127.0.0.1'

def _check_login_rate_limit(ip: str, username: str) -> bool:
    now = time.time()
    ip_key = f"ip:{ip}"
    user_key = f"user:{username}"
    with _login_lock:
        if ip_key not in _login_attempts:
            _login_attempts[ip_key] = []
        if user_key not in _login_attempts:
            _login_attempts[user_key] = []
        
        ip_attempts = [t for t in _login_attempts[ip_key] if now - t < 900]
        user_attempts = [t for t in _login_attempts[user_key] if now - t < 900]
        
        _login_attempts[ip_key] = ip_attempts
        _login_attempts[user_key] = user_attempts
        
        if len(_login_attempts) > 1000:
            expired_keys = [k for k, v in _login_attempts.items() if not v]
            for k in expired_keys:
                del _login_attempts[k]
        
        if len(ip_attempts) >= 15 or len(user_attempts) >= 5:
            return False
        return True

def _record_login_failure(ip: str, username: str):
    now = time.time()
    ip_key = f"ip:{ip}"
    user_key = f"user:{username}"
    with _login_lock:
        if ip_key not in _login_attempts:
            _login_attempts[ip_key] = []
        if user_key not in _login_attempts:
            _login_attempts[user_key] = []
        _login_attempts[ip_key].append(now)
        _login_attempts[user_key].append(now)

def _clear_login_failures(ip: str, username: str):
    ip_key = f"ip:{ip}"
    user_key = f"user:{username}"
    with _login_lock:
        _login_attempts.pop(ip_key, None)
        _login_attempts.pop(user_key, None)

def _set_admin_auth_cookies(response, admin_session_token):
    """下发管理员会话 cookie。仅携带服务端会话 token，不再下发 admin_user 之类的身份信息。"""
    is_secure = os.environ.get('FLASK_ENV') == 'production'
    response.set_cookie('admin_session', admin_session_token, httponly=True, secure=is_secure,
                        samesite='Lax', max_age=database.ADMIN_SESSION_TTL)

def _clear_auth_cookies(response):
    # 兼容旧 cookie（admin_token/admin_user）一并清除，确保旧会话失效
    for cookie_name in ['session_token', 'admin_session', 'admin_token', 'admin_user']:
        response.delete_cookie(cookie_name, samesite='Lax')


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        data = request.get_json(silent=True)
        if data is None:
            data = request.form
        if not data:
            return jsonify({"error": "无效的请求格式"}), 400
        username = data.get("username")
        password = data.get("password")
        ip = _get_client_ip()

        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400

        if not _check_login_rate_limit(ip, username):
            return jsonify({"error": "登录尝试过于频繁，请5分钟后再试"}), 429

        # 管理员：通过数据库 admin_users 表验证，登录成功创建服务端会话
        admin_result = database.verify_admin_password(username, password)
        if admin_result.get("success"):
            admin = admin_result["admin"]
            _clear_login_failures(ip, username)
            session_token = database.create_admin_session(
                admin["id"], ip_address=ip,
                user_agent=request.headers.get('User-Agent', '')[:200]
            )
            if not session_token:
                return jsonify({"error": "登录失败，会话创建异常"}), 500
            database.record_login_attempt(ip, username, success=True)
            database.log_admin_action(admin["username"], "login",
                                     details=f"管理员登录（前台入口）display={admin['display_name']}",
                                     ip_address=ip)
            response = jsonify({"status": "success", "username": admin["display_name"], "is_admin": True})
            _set_admin_auth_cookies(response, session_token)
            return response

        # 管理员账号被识别但验证失败（密码错/被禁用），直接走失败流程
        if admin_result.get("reason") in ("wrong_password", "inactive", "corrupt"):
            _record_login_failure(ip, username)
            database.record_login_attempt(ip, username, success=False)
            return jsonify({"error": admin_result.get("error", "用户名或密码错误")}), 401

        login_email = username
        if '@' not in username:
            local_email = database.get_email_by_username(username)
            if local_email:
                login_email = local_email
            else:
                email_result = supabase_client.get_email_by_username(username)
                if not email_result["success"]:
                    _record_login_failure(ip, username)
                    return jsonify({"error": "用户名或密码错误"}), 401
                login_email = email_result["email"]

        user = database.authenticate_user(login_email, password)
        if user:
            _clear_login_failures(ip, username)
            access_token = user.get('access_token') or database.generate_session_token()
            csrf_token = _generate_csrf_token(user['id'])
            response = jsonify({"status": "success", "username": user['username'], "is_admin": False, "csrf_token": csrf_token})
            # 普通用户：下发 session_token cookie
            is_secure = os.environ.get('FLASK_ENV') == 'production'
            response.set_cookie('session_token', access_token, httponly=True, secure=is_secure,
                                samesite='Lax', max_age=86400*7)
            return response
        else:
            _record_login_failure(ip, username)
            return jsonify({"error": "用户名或密码错误"}), 401

    return render_template("login.html")

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        data = safe_get_json()
        if not data:
            return jsonify({"error": "无效的请求格式"}), 400
        username = data.get("username")
        password = data.get("password")
        email = data.get("email", "")
        
        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400
        
        if len(password) < 8:
            return jsonify({"error": "密码长度至少8位"}), 400
        
        if not re.search(r'[A-Za-z]', password) or not re.search(r'[0-9]', password):
            return jsonify({"error": "密码需包含字母和数字"}), 400

        if not email or not re.match(r'^[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}$', email):
            return jsonify({"error": "请输入有效的邮箱地址"}), 400

        if len(username) < 2 or len(username) > 20:
            return jsonify({"error": "用户名长度需在2-20位之间"}), 400
        
        if not re.match(r'^[a-zA-Z0-9_\u4e00-\u9fa5]+$', username):
            return jsonify({"error": "用户名只允许字母、数字、下划线和中文"}), 400
        
        success = database.register_user(username, password, email)
        if success:
            return jsonify({"status": "success"})
        else:
            return jsonify({"error": "注册失败，请稍后重试"}), 409
    
    return render_template("register.html")

@app.route("/logout", methods=["GET", "POST"])
def logout():
    user_token = request.cookies.get('session_token')
    if user_token:
        database.clear_session_token(user_token)
        try:
            supabase_client.sign_out_by_token(user_token)
        except Exception:
            pass
    
    admin_token = request.cookies.get('admin_session')
    if admin_token:
        database.delete_admin_session(admin_token)
    
    if request.method == 'GET':
        response = make_response(redirect(url_for('login')))
    else:
        response = make_response(jsonify({"status": "success"}))
    _clear_auth_cookies(response)
    return response

@app.route("/profile")
def profile():
    user = get_current_user()
    if not user:
        return redirect(url_for('login'))

    all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]

    # 获取用户使用统计
    usage_stats = {"total_sessions": 0, "total_messages": 0, "total_files": 0}
    try:
        stats = database.get_user_sessions_with_stats(user['id'])
        usage_stats["total_sessions"] = len(stats)
        total_msgs = sum(s.get("total_count") or 0 for s in stats)
        usage_stats["total_messages"] = total_msgs
        total_files = 0
        for s in stats:
            files = s.get("files") or []
            total_files += len(files)
        usage_stats["total_files"] = total_files
    except Exception as e:
        print(f"[profile] 获取使用统计失败: {e}")

    return render_template("profile.html", user=user, all_modules=all_modules, usage_stats=usage_stats)

@app.route("/api/profile", methods=["GET"])
def get_profile():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401
    
    if user.get('is_admin'):
        avatar_url = database.get_admin_avatar(user['id'])
        return jsonify({
            "success": True,
            "user": {
                "username": user['username'],
                "email": user['email'],
                "avatar_url": avatar_url
            }
        })
    
    result = supabase_client.get_user_by_id(user['id'])
    if result["success"]:
        return jsonify({"success": True, "user": result["user"]})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500

@app.route("/api/profile/username", methods=["POST"])
def update_username():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    new_username = (data.get("username") or "").strip()
    # 校验：与注册时保持一致 —— 长度 2-20，仅允许字母/数字/下划线/中文
    if not new_username or len(new_username) < 2:
        return jsonify({"error": "用户名至少2个字符"}), 400
    if len(new_username) > 20:
        return jsonify({"error": "用户名长度不超过20位"}), 400
    if not re.match(r'^[a-zA-Z0-9_\u4e00-\u9fa5]+$', new_username):
        return jsonify({"error": "用户名只允许字母、数字、下划线和中文"}), 400

    # 显示名与当前一致则直接返回成功
    if new_username == user.get('username'):
        return jsonify({"success": True, "username": new_username})

    if user.get('is_admin'):
        # 管理员：只改显示名，不改登录账号（admin_users.username 永远不变）
        if not database.update_admin_display_name(user['id'], new_username):
            return jsonify({"error": "保存失败"}), 500
        return jsonify({"success": True, "username": new_username})

    exist_result = supabase_client.username_exists(new_username)
    if exist_result["success"]:
        if exist_result["exists"] and exist_result.get("user_id") != user['id']:
            return jsonify({"success": False, "error": "该用户名已被使用"}), 409

    result = supabase_client.update_username(user['id'], new_username)
    if result["success"]:
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        try:
            cursor.execute('UPDATE users SET username = ? WHERE user_id = ?', (new_username, user['id']))
            conn.commit()
        finally:
            conn.close()
        return jsonify({"success": True, "username": new_username})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500

@app.route("/api/profile/password", methods=["POST"])
def update_password():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401
    
    if user.get('is_admin'):
        return jsonify({"success": False, "error": "管理员账号不支持修改密码"}), 400
    
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400
    
    new_password = data.get("password")
    confirm_password = data.get("confirm_password")
    email = data.get("email")
    
    if not new_password or not confirm_password:
        return jsonify({"error": "请输入密码"}), 400
    
    if new_password != confirm_password:
        return jsonify({"error": "两次输入的密码不一致"}), 400
    
    if len(new_password) < 6:
        return jsonify({"error": "密码长度至少6位"}), 400
    
    if email != user.get('email'):
        return jsonify({"error": "邮箱验证失败"}), 400
    
    result = supabase_client.send_password_reset_email(email)
    if result["success"]:
        return jsonify({"success": True, "message": "密码重置邮件已发送，请查看邮箱完成密码修改"})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500

@app.route("/api/profile/avatar", methods=["POST"])
def upload_avatar():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401
    
    if 'avatar' not in request.files:
        return jsonify({"error": "请选择图片"}), 400
    
    file = request.files['avatar']
    if file.filename == '':
        return jsonify({"error": "请选择图片"}), 400
    
    allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
    if '.' not in file.filename or file.filename.rsplit('.', 1)[1].lower() not in allowed_extensions:
        return jsonify({"error": "只支持 PNG、JPG、JPEG、GIF、WEBP 格式"}), 400
    
    file_data = file.read()
    if len(file_data) > 5 * 1024 * 1024:
        return jsonify({"error": "图片大小不能超过5MB"}), 400
    
    valid_signatures = {
        'png': b'\x89PNG\r\n\x1a\n',
        'jpg': b'\xff\xd8\xff',
        'jpeg': b'\xff\xd8\xff',
        'gif': b'GIF8',
        'webp': b'RIFF....WEBP',
    }
    ext = file.filename.rsplit('.', 1)[1].lower()
    sig = valid_signatures.get(ext)
    if sig:
        if ext == 'webp':
            if not file_data.startswith(b'RIFF') or len(file_data) < 12 or file_data[8:12] != b'WEBP':
                return jsonify({"error": "无效的图片文件"}), 400
        else:
            if not file_data.startswith(sig):
                return jsonify({"error": "无效的图片文件"}), 400
    
    import uuid
    import os
    filename = f"avatar_{user['id']}_{uuid.uuid4().hex[:8]}.{ext}"
    upload_dir = os.path.join(DATA_DIR, 'uploads', 'avatars')
    os.makedirs(upload_dir, exist_ok=True)
    filepath = os.path.join(upload_dir, filename)

    # 获取旧头像 URL，便于新头像写入成功后清理旧文件
    old_avatar_url = None
    if user.get('is_admin'):
        old_avatar_url = database.get_admin_avatar(user['id'])
    else:
        try:
            user_lookup = supabase_client.get_user_by_id(user['id'])
            if user_lookup.get("success"):
                old_avatar_url = user_lookup["user"].get("avatar_url")
        except Exception:
            pass

    with open(filepath, "wb") as f:
        f.write(file_data)

    avatar_url = f"/api/avatar/{filename}"

    def _cleanup_old_avatar_file(old_url):
        """从旧的 /api/avatar/<filename> URL 提取文件名并安全删除"""
        if not old_url or '/api/avatar/' not in old_url:
            return
        old_filename = old_url.rsplit('/', 1)[-1]
        if not _is_safe_filename(old_filename):
            return
        old_filepath = os.path.realpath(os.path.join(upload_dir, old_filename))
        try:
            if os.path.isfile(old_filepath) and old_filepath.startswith(os.path.realpath(upload_dir) + os.sep):
                os.remove(old_filepath)
        except OSError:
            pass

    if user.get('is_admin'):
        if not database.update_admin_avatar(user['id'], avatar_url):
            # 数据库写入失败，删除已落盘的文件避免悬空
            try:
                os.remove(filepath)
            except OSError:
                pass
            return jsonify({"error": "保存失败"}), 500
        _cleanup_old_avatar_file(old_avatar_url)
        return jsonify({"success": True, "avatar_url": avatar_url})

    result = supabase_client.update_user_avatar(user['id'], avatar_url)

    if result["success"]:
        _cleanup_old_avatar_file(old_avatar_url)
        return jsonify({"success": True, "avatar_url": avatar_url})
    else:
        try:
            os.remove(filepath)
        except Exception:
            pass
        return jsonify({"success": False, "error": "上传失败，请稍后重试"}), 500

@app.route("/api/auth/session", methods=["POST"])
def set_auth_session():
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400
    
    access_token = data.get("access_token")
    if not access_token:
        return jsonify({"error": "缺少 access_token"}), 400
    
    result = supabase_client.get_user(access_token)
    if not result["success"]:
        return jsonify({"error": "无效的 token"}), 401
    
    response = jsonify({"status": "success", "username": result["user"]["username"] or result["user"]["email"].split("@")[0], "csrf_token": _generate_csrf_token(result["user"]["id"])})
    is_secure = os.environ.get('FLASK_ENV') == 'production'
    response.set_cookie('session_token', access_token, httponly=True, secure=is_secure,
                        samesite='Lax', max_age=86400*7)
    return response

@app.route("/api/auth/get-csrf-token", methods=["GET"])
def get_csrf_token():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401
    csrf_token = _generate_csrf_token(user['id'])
    return jsonify({"success": True, "csrf_token": csrf_token})


@app.route("/api/auth/get-email", methods=["POST"])
def get_email_by_username():
    # 修复 S3：限制仅管理员可调用，避免通过该接口枚举用户邮箱
    if not is_admin_logged_in():
        return jsonify({"error": "仅管理员可执行此操作"}), 403

    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400
    
    username = data.get("username")
    if not username:
        return jsonify({"error": "请提供用户名"}), 400
    
    local_email = database.get_email_by_username(username)
    if local_email:
        return jsonify({"success": True, "email": local_email})
    
    result = supabase_client.get_email_by_username(username)
    if result["success"]:
        return jsonify({"success": True, "email": result["email"]})
    else:
        return jsonify({"success": False, "error": result["error"]}), 401

_CHAT_RATE_LIMIT = {}
_CHAT_RATE_LOCK = threading.Lock()


def _check_chat_rate_limit(ip: str) -> bool:
    now = time.time()
    with _CHAT_RATE_LOCK:
        if ip not in _CHAT_RATE_LIMIT:
            _CHAT_RATE_LIMIT[ip] = []
        attempts = _CHAT_RATE_LIMIT[ip]
        attempts = [t for t in attempts if now - t < 60]
        _CHAT_RATE_LIMIT[ip] = attempts
        if len(attempts) >= 20:
            return False
        _CHAT_RATE_LIMIT[ip].append(now)
        # 修复 M3：定期清理空IP键，避免限流字典无界增长导致内存泄漏
        if len(_CHAT_RATE_LIMIT) > 1000:
            expired_keys = [k for k, v in _CHAT_RATE_LIMIT.items() if not v]
            for k in expired_keys:
                del _CHAT_RATE_LIMIT[k]
    return True


def _prepare_chat_message(data):
    user_message = data.get("message", "").strip()
    attachments = data.get("attachments") or []

    if isinstance(attachments, list) and len(attachments) <= 10:
        blocks = []
        for att in attachments:
            if not isinstance(att, dict):
                continue
            name = str(att.get("name", "未命名文件"))[:120]
            kind = str(att.get("kind", "资料"))[:30]
            full_text = str(att.get("text") or "")  # 修复 L10：att.get("text", "") 在值为 None 时仍返回 None，需用 or 兜底
            text = full_text[:MAX_TEXT_CHARS]
            original_len = len(full_text)
            truncated_notice = ""
            if original_len > MAX_TEXT_CHARS:
                truncated_notice = f"\n（⚠️ 原文共约 {original_len} 字，已截断仅显示前 {MAX_TEXT_CHARS} 字，如需更多请告知）"
            blocks.append(
                f"【用户上传文件：{name}（{kind}）】\n"
                f"文件内容预览（节选，共约 {len(text)} 字）：\n"
                f"```\n{text}\n```\n"
                f"【上传文件结束】{truncated_notice}"
            )
        if blocks:
            if user_message:
                user_message = user_message + "\n\n---\n\n" + "\n\n".join(blocks)
            else:
                user_message = "\n\n".join(blocks)

    if not user_message:
        return None, "消息不能为空"

    if len(user_message) > 24000:
        user_message = user_message[:24000] + "\n\n[内容过长已截断]"

    return user_message, None


@app.route("/api/chat", methods=["POST"])
def chat():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    ip = request.remote_addr or 'unknown'
    if request.headers.get('X-Forwarded-For'):
        ip = request.headers.get('X-Forwarded-For').split(',')[0].strip()

    if not _check_chat_rate_limit(ip):
        return jsonify({"error": "请求过于频繁，请1分钟后再试"}), 429

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "无效请求"}), 400

    session_id = str(data.get("session_id", "default"))[:64]
    user_message, err = _prepare_chat_message(data)
    if err:
        return jsonify({"error": err}), 400

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return jsonify({"error": "参数格式错误"}), 400

    key = _session_key(user['id'], session_id)
    with sessions_lock:
        if key not in sessions:
            sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]
        sessions[key].append({"role": "user", "content": user_message})

    try:
        client = OpenAI(api_key=API_KEY, base_url=BASE_URL, timeout=60.0)
        with sessions_lock:
            msgs = trim_history(sessions[key])
        response = client.chat.completions.create(
            model=_model,
            messages=msgs,
            temperature=_temp,
            max_tokens=_max_tokens,
        )
    except Exception as e:
        print(f"[Chat] API request failed: {e}")
        with sessions_lock:
            msgs = sessions.get(key, [])
            if msgs and msgs[-1].get("role") == "user" and msgs[-1].get("content") == user_message:
                msgs.pop()
        return jsonify({"error": "请求失败，请稍后重试"}), 500

    if not response.choices or not response.choices[0].message.content:
        with sessions_lock:
            msgs = sessions.get(key, [])
            if msgs and msgs[-1].get("role") == "user" and msgs[-1].get("content") == user_message:
                msgs.pop()
        return jsonify({"error": "请求失败，请稍后重试"}), 500

    reply = response.choices[0].message.content
    usage = response.usage
    # 修复 M1：response.usage 可能为 None，访问前需判空，避免 AttributeError
    prompt_tokens = usage.prompt_tokens if usage else 0
    completion_tokens = usage.completion_tokens if usage else 0

    with sessions_lock:
        sessions[key].append({"role": "assistant", "content": reply})
        first_user_msg = next((m["content"] for m in sessions[key] if m["role"] == "user"), "")
    title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
    database.save_user_chat_session(user['id'], session_id, title)

    saved_files, saved_ppts, display = extract_and_save_all(user_message, reply, user['id'])

    return jsonify({
        "reply": display,
        "raw_reply": reply,
        "model": response.model,
        "usage": {
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
        },
        "saved_files": saved_files,
        "saved_ppts": saved_ppts,
    })

def extract_keywords(message):
    keywords = []
    subjects = ["物理", "数学", "语文", "英语", "化学", "生物", "历史", "地理", "政治", "信息技术", "通用技术", "体育", "音乐", "美术", "科学", "道德与法治", "思想品德"]
    for subj in subjects:
        if subj in message:
            keywords.append(subj)
            break
    levels = ["小学", "初中", "高中", "大学", "职业院校", "中职", "高职", "一年级", "二年级", "三年级", "四年级", "五年级", "六年级", "初一", "初二", "初三", "高一", "高二", "高三"]
    for level in levels:
        if level in message:
            keywords.append(level)
            break
    types = ["教案", "课件", "习题", "试卷", "复习", "提纲", "实验", "实训", "说课稿", "导学案", "任务单", "教学设计"]
    for t in types:
        if t in message:
            keywords.append(t)
            break
    if not keywords:
        keywords = ["教学", "备课"]
    return keywords

knowledge_base = {
    "物理": ["牛顿运动定律", "能量守恒", "电磁感应", "光学", "热力学", "力学", "相对论", "量子物理", "波动", "磁场"],
    "数学": ["函数", "方程", "几何", "概率统计", "数列", "三角函数", "向量", "导数", "积分", "不等式"],
    "语文": ["文言文", "现代文", "诗词鉴赏", "写作", "阅读", "修辞手法", "表现手法", "作文", "记叙文", "议论文"],
    "英语": ["语法", "词汇", "阅读", "写作", "听力", "口语", "时态", "句型", "翻译", "完形填空"],
    "化学": ["元素周期表", "化学反应", "有机化学", "无机化学", "化学平衡", "电化学", "化学实验", "化学键"],
    "生物": ["细胞", "遗传", "生态", "光合作用", "呼吸作用", "进化论", "微生物", "人体生理"],
    "历史": ["中国古代史", "中国近代史", "世界史", "历史事件", "历史人物", "朝代", "战争", "改革"],
    "地理": ["自然地理", "人文地理", "气候", "地形", "地图", "环境保护", "区域地理", "人口城市"],
    "政治": ["经济生活", "政治生活", "文化生活", "哲学生活", "法律", "道德", "国情"],
    "default": ["教学目标", "教学重难点", "教学过程", "教学方法", "板书设计", "课后作业", "教学反思", "学情分析"]
}

def generate_thinking_phases(user_message):
    keywords = extract_keywords(user_message)
    subject = keywords[0] if keywords else "default"
    kbs = knowledge_base.get(subject, knowledge_base["default"])
    selected_kb = random.sample(kbs, min(4, len(kbs)))

    phases = [
        {"title": "📋 需求分析阶段", "items": [
            f"解析用户需求：{user_message[:25]}{'...' if len(user_message) > 25 else ''}",
            f"识别学科与学段：{'、'.join(keywords)}",
            "明确输出类型与格式要求",
        ]},
        {"title": "🔍 知识检索阶段", "items": [
            "在知识库中检索相关教学资源...",
            f"匹配知识点：{selected_kb[0]}、{selected_kb[1]}",
            "查找课程标准与教学大纲要求",
            f"筛选参考资料：{selected_kb[2]}相关案例",
        ]},
        {"title": "🧠 内容规划阶段", "items": [
            "梳理知识体系与逻辑结构",
            "设计三维教学目标",
            "确定教学重点与难点",
            "规划教学环节与时间分配",
            "选择教学方法与策略",
        ]},
        {"title": "✍️ 内容生成阶段", "items": [
            "组织语言，生成详细内容...",
            "优化表述，确保专业准确",
            "调整结构，保证逻辑清晰",
            "检查内容完整性与合理性",
        ]},
    ]
    return phases

def thinking_to_text(phases, phase_idx, item_idx):
    lines = []
    for i, phase in enumerate(phases):
        if i > phase_idx:
            break
        lines.append(f"【{phase['title']}】")
        for j, item in enumerate(phase['items']):
            if i == phase_idx and j > item_idx:
                break
            if i < phase_idx or (i == phase_idx and j < item_idx):
                icon = "✓"
            elif i == phase_idx and j == item_idx:
                icon = "⏳"
            else:
                icon = "○"
            lines.append(f"  {icon} {item}")
        lines.append("")
    return "\n".join(lines)

@app.route("/api/chat/stream", methods=["POST"])
def chat_stream():
    user = get_current_user()
    if not user:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '请先登录'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    ip = request.remote_addr or 'unknown'
    if request.headers.get('X-Forwarded-For'):
        ip = request.headers.get('X-Forwarded-For').split(',')[0].strip()

    if not _check_chat_rate_limit(ip):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '请求过于频繁，请1分钟后再试'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '无效请求'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    session_id = str(data.get("session_id", "default"))[:64]
    user_message, err = _prepare_chat_message(data)
    if err:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': err}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '参数格式错误'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    key = _session_key(user['id'], session_id)
    with sessions_lock:
        if key not in sessions:
            sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]
        sessions[key].append({"role": "user", "content": user_message})

    def generate():
        thinking_phases = generate_thinking_phases(user_message)
        thinking_done = False
        ai_done = False
        ai_error = None
        full_reply = ""
        ai_queue = queue.Queue()
        first_content_received = False
        stop_event = threading.Event()
        start_time = time.time()
        MAX_STREAM_DURATION = 300

        def ai_worker():
            nonlocal full_reply, ai_done, ai_error
            try:
                client = OpenAI(api_key=API_KEY, base_url=BASE_URL, timeout=60.0)
                with sessions_lock:
                    msgs = trim_history(sessions[key])
                stream = client.chat.completions.create(
                    model=_model,
                    messages=msgs,
                    temperature=_temp,
                    max_tokens=_max_tokens,
                    stream=True,
                )
                for chunk in stream:
                    if stop_event.is_set() or (time.time() - start_time) > MAX_STREAM_DURATION:
                        try:
                            stream.close()
                        except Exception:
                            pass
                        break
                    if chunk.choices and chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        full_reply += content
                        ai_queue.put(("content", content))
                ai_queue.put(("done", None))
                ai_done = True
            except Exception as e:
                print(f"[Chat] Stream request failed: {type(e).__name__}: {e}")
                ai_error = _format_ai_error(e)
                ai_queue.put(("error", ai_error))
                ai_done = True

        ai_thread = threading.Thread(target=ai_worker, daemon=True)
        ai_thread.start()

        total_thinking_time = 0
        min_thinking_time = 1.2
        max_thinking_time = 3.0

        for phase_idx, phase in enumerate(thinking_phases):
            if thinking_done:
                break
            for item_idx, item in enumerate(phase["items"]):
                if thinking_done:
                    break

                thinking_text = thinking_to_text(thinking_phases, phase_idx, item_idx)
                yield f"data: {json.dumps({'type': 'thinking', 'content': thinking_text}, ensure_ascii=False)}\n\n"

                sleep_time = random.uniform(0.25, 0.4)
                if total_thinking_time + sleep_time > max_thinking_time and phase_idx >= 2:
                    sleep_time = 0.1
                total_thinking_time += sleep_time
                time.sleep(sleep_time)

                if total_thinking_time >= min_thinking_time:
                    while not ai_queue.empty():
                        msg_type, msg_data = ai_queue.get()
                        if msg_type == "content":
                            # 修复 B1：首个 content 分片必须立即 yield，否则会丢失
                            yield f"data: {json.dumps({'type': 'content', 'content': msg_data}, ensure_ascii=False)}\n\n"
                            first_content_received = True
                            break
                        elif msg_type == "done":
                            first_content_received = True
                            ai_done = True
                            break
                        elif msg_type == "error":
                            ai_error = msg_data
                            first_content_received = True
                            ai_done = True
                            break

                if first_content_received and total_thinking_time >= min_thinking_time:
                    thinking_done = True
                    final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
                    yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
                    break

        if not thinking_done:
            final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
            yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
            thinking_done = True

        try:
            while not ai_done or not ai_queue.empty():
                try:
                    msg_type, msg_data = ai_queue.get(timeout=0.1)
                    if msg_type == "content":
                        yield f"data: {json.dumps({'type': 'content', 'content': msg_data}, ensure_ascii=False)}\n\n"
                    elif msg_type == "done":
                        ai_done = True
                    elif msg_type == "error":
                        ai_error = msg_data
                        ai_done = True
                except queue.Empty:
                    if not ai_done:
                        if (time.time() - start_time) > MAX_STREAM_DURATION:
                            ai_error = "请求超时，请稍后重试"
                            ai_done = True
                        else:
                            time.sleep(0.05)
        except GeneratorExit:
            stop_event.set()
            ai_thread.join(timeout=5.0)
            if full_reply:
                with sessions_lock:
                    sessions[key].append({"role": "assistant", "content": full_reply})
                database.save_user_chat_session(user['id'], session_id,
                    (next((m["content"] for m in sessions[key] if m["role"] == "user"), "")[:30]) + "...")
            return

        if ai_error:
            if full_reply:
                with sessions_lock:
                    sessions[key].append({"role": "assistant", "content": full_reply})
                first_user_msg = next((m["content"] for m in sessions[key] if m["role"] == "user"), "")
                title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
                database.save_user_chat_session(user['id'], session_id, title)
                # 修复 B2：删除重复 yield full_reply（内容已在流中下发，避免重复发送）
            yield f"data: {json.dumps({'type': 'error', 'content': ai_error}, ensure_ascii=False)}\n\n"
            if not full_reply:
                with sessions_lock:
                    msgs = sessions.get(key, [])
                    if msgs and msgs[-1].get("role") == "user" and msgs[-1].get("content") == user_message:
                        msgs.pop()
            return

        if full_reply:
            with sessions_lock:
                sessions[key].append({"role": "assistant", "content": full_reply})
                first_user_msg = next((m["content"] for m in sessions[key] if m["role"] == "user"), "")
            title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
            database.save_user_chat_session(user['id'], session_id, title)

            saved_files, saved_ppts, display = extract_and_save_all(user_message, full_reply, user['id'])

            final_thinking_text = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
            yield f"data: {json.dumps({'type': 'done', 'content': display, 'raw_content': full_reply, 'saved_files': saved_files, 'saved_ppts': saved_ppts, 'thinking': final_thinking_text}, ensure_ascii=False)}\n\n"

    return Response(generate(), mimetype="text/event-stream", headers={
        'Cache-Control': 'no-cache',
        'X-Accel-Buffering': 'no',
        'Connection': 'keep-alive',
    })

@app.route("/api/clear", methods=["POST"])
def clear():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json() or {}
    session_id = data.get("session_id", "default")
    key = _session_key(user['id'], session_id)
    with sessions_lock:
        sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]
    return jsonify({"status": "ok"})

@app.route("/api/sessions", methods=["GET"])
def list_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    db_sessions = database.get_user_sessions(user['id'])
    result = []
    for s in db_sessions:
        result.append({
            "id": s['session_id'],
            "title": s['title'],
            "message_count": 0,
            "created_at": s['created_at']
        })
    return jsonify({"sessions": result})

@app.route("/api/session/<session_id>", methods=["GET"])
def get_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    key = _session_key(user['id'], session_id)
    with sessions_lock:
        if key in sessions:
            msgs = sessions[key]
        else:
            msgs = []
    if msgs:
        user_msgs = [m for m in msgs if m["role"] == "user"]
        title = "新对话"
        if user_msgs:
            first_user = user_msgs[0]["content"]
            title = first_user[:20] + "..." if len(first_user) > 20 else first_user
        return jsonify({
            "id": session_id,
            "title": title,
            "messages": msgs
        })
    
    db_sessions = database.get_user_sessions(user['id'])
    db_session = next((s for s in db_sessions if s['session_id'] == session_id), None)
    if db_session:
        return jsonify({
            "id": session_id,
            "title": db_session['title'],
            "messages": []
        })
    
    return jsonify({"error": "会话不存在"}), 404

@app.route("/api/session/<session_id>", methods=["DELETE"])
def delete_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    database.delete_user_session(user['id'], session_id)
    key = _session_key(user['id'], session_id)
    with sessions_lock:
        if key in sessions:
            del sessions[key]
    return jsonify({"status": "ok"})

@app.route("/api/sessions/batch", methods=["DELETE"])
def delete_sessions_batch():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json() or {}
    ids = data.get("ids", [])
    deleted = 0
    for sid in ids:
        database.delete_user_session(user['id'], sid)
        key = _session_key(user['id'], sid)
        with sessions_lock:
            if key in sessions:
                del sessions[key]
        deleted += 1
    return jsonify({"status": "ok", "deleted": deleted})

@app.route("/api/sessions/all", methods=["DELETE"])
def delete_all_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    prefix = f"{user['id']}:"
    with sessions_lock:
        to_delete = [key for key in sessions if key.startswith(prefix)]
        for key in to_delete:
            del sessions[key]
    db_deleted = database.delete_all_user_sessions(user['id'])
    return jsonify({"status": "ok", "memory_deleted": len(to_delete), "db_deleted": db_deleted})

@app.route("/api/search_history", methods=["GET"])
def get_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    history = database.get_search_history(user['id'], limit=50)
    return jsonify({"history": history})

@app.route("/api/search_history", methods=["DELETE"])
def clear_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    deleted = database.clear_search_history(user['id'])
    return jsonify({"status": "ok", "deleted": deleted})

@app.route("/api/config", methods=["GET"])
def get_config():
    return jsonify({
        "model": MODEL,
        "temperature": TEMP,
        "max_tokens": MAX_TOKENS,
    })

MIME_MAP = {
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.md': 'text/markdown; charset=utf-8',
    '.txt': 'text/plain; charset=utf-8',
    '.csv': 'text/csv; charset=utf-8',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.pdf': 'application/pdf',
}


def _read_text_bytes(b: bytes) -> str:
    for enc in ('utf-8-sig', 'utf-8', 'gbk', 'gb18030', 'utf-16'):
        try:
            return b.decode(enc)
        except UnicodeDecodeError:
            continue
    return b.decode('utf-8', errors='ignore')


def _guess_file_kind(name: str, text: str) -> str:
    low = (name + "\n" + text[:1500]).lower()
    template_keys = ['教案', '教学设计', '学案', '模板', '____', '___', '待填写', '留白', '[    ]', '[___]', '填写说明', '教学目标', '教学重难点']
    homework_keys = ['作业', '成绩', '得分', '分数', '错题', '批改', '班级', '学生', '姓名', '学号', '排名', '及格', '均分', '考勤', '总分']
    t_hit = sum(1 for k in template_keys if k in low)
    h_hit = sum(1 for k in homework_keys if k in low)
    if t_hit >= 2 and t_hit >= h_hit:
        return "教案/教学设计模板"
    if h_hit >= 2:
        return "学生作业/成绩/批改记录"
    if '.xlsx' in name.lower() or '.csv' in name.lower():
        return "数据表"
    if '.pdf' in name.lower():
        return "PDF文档"
    return "参考资料"


def _table_to_markdown(rows):
    rows = list(rows)
    if not rows:
        return ""
    rows = [[("" if v is None else str(v)).replace('|', '\\|').replace('\n', ' ') for v in r] for r in rows]
    max_cols = max(len(r) for r in rows)
    rows = [r + [''] * (max_cols - len(r)) for r in rows]
    header = rows[0]
    sep = ['---'] * max_cols
    body = rows[1:]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join(sep) + " |"]
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


@app.route("/api/upload", methods=["POST"])
def api_upload():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    if 'file' not in request.files:
        return jsonify({"error": "未收到上传文件字段 file"}), 400

    f = request.files['file']
    if not f or not f.filename:
        return jsonify({"error": "文件名为空"}), 400

    filename = os.path.basename(f.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXT:
        return jsonify({"error": f"不支持的格式 {ext}，支持：md / txt / docx / xlsx / csv / pdf"}), 400

    # 修复 M4：先校验 Content-Length，避免把超大文件整体读入内存后才报错
    content_length = request.content_length or 0
    if content_length > MAX_UPLOAD_BYTES:
        return jsonify({"error": "文件过大，超出上传上限"}), 413

    raw = f.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return jsonify({"error": f"文件过大 {len(raw)} bytes，上限 {MAX_UPLOAD_BYTES} bytes（20MB）"}), 400

    text = ""
    try:
        if ext in ('.txt', '.md'):
            text = _read_text_bytes(raw)
        elif ext == '.docx':
            bio = BytesIO(raw)
            doc = Document(bio)
            parts = []
            for p in doc.paragraphs:
                if p.text:
                    parts.append(p.text)
            for t in doc.tables:
                rows = [[c.text for c in row.cells] for row in t.rows]
                md = _table_to_markdown(rows)
                if md:
                    parts.append("\n[表格]\n" + md + "\n")
            text = "\n".join(parts)
        elif ext == '.csv':
            s = _read_text_bytes(raw)
            import csv as _csv
            reader = _csv.reader(s.splitlines())
            rows = list(reader)[:200]
            text = _table_to_markdown(rows) or (s[:MAX_TEXT_CHARS])
        elif ext == '.xlsx':
            import openpyxl
            bio = BytesIO(raw)
            # 修复 M5：使用 with 语句保证 workbook 在异常或正常结束后均被关闭，避免资源泄漏
            with openpyxl.load_workbook(bio, read_only=True, data_only=True) as wb:
                sheets_text = []
                for ws in wb.worksheets:
                    rows = []
                    for idx, row in enumerate(ws.iter_rows(values_only=True)):
                        if idx > 60:
                            break
                        rows.append(list(row)[:16])
                    md = _table_to_markdown(rows)
                    if md:
                        sheets_text.append(f"# 工作表：{ws.title}\n{md}")
                text = "\n\n".join(sheets_text)
        elif ext == '.pdf':
            try:
                import PyPDF2
                bio = BytesIO(raw)
                reader = PyPDF2.PdfReader(bio)
                pages = []
                for i, page in enumerate(reader.pages[:30]):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    if t:
                        pages.append(f"\n---第 {i+1} 页---\n{t}")
                text = "\n".join(pages)
            except ImportError:
                return jsonify({"error": "PDF 需要先安装 PyPDF2：pip install PyPDF2"}), 400
    except Exception as ee:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"解析文件失败: {ee}"}), 500

    if not text or not text.strip():
        return jsonify({"error": "文件内容为空或无法识别文本"}), 400

    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n[原文过长已截断，仅保留前部分内容用于分析]"

    kind = _guess_file_kind(filename, text)
    preview = text[:600]

    return jsonify({
        "ok": True,
        "name": filename,
        "kind": kind,
        "chars": len(text),
        "preview": preview,
        "full_text": text,
    })


def _is_safe_filename(filename):
    if not filename or not isinstance(filename, str):
        return False
    base = os.path.basename(filename)
    if base != filename or base in ('.', '..') or not base:
        return False
    if any(c in filename for c in ('\\', '/', '..')):
        return False
    if filename.startswith('.'):
        return False
    return True


def _find_user_file(filename):
    if not _is_safe_filename(filename):
        return None
    user = get_current_user()
    
    if user:
        user_dir = os.path.realpath(get_user_output_dir(user['id']))
        fp = os.path.realpath(os.path.join(user_dir, filename))
        if os.path.isfile(fp) and fp.startswith(user_dir + os.sep):
            return fp
    
    out_real = os.path.realpath(OUTPUT_DIR)
    fp = os.path.realpath(os.path.join(out_real, filename))
    if os.path.isfile(fp) and fp.startswith(out_real + os.sep):
        return fp
    
    vis_real = os.path.realpath(str(VISUAL_OUTPUT_DIR))
    fp = os.path.realpath(os.path.join(vis_real, filename))
    if os.path.isfile(fp) and fp.startswith(vis_real + os.sep):
        return fp
    
    avatars_dir = os.path.realpath(os.path.join(DATA_DIR, 'uploads', 'avatars'))
    fp = os.path.realpath(os.path.join(avatars_dir, filename))
    if os.path.isfile(fp) and fp.startswith(avatars_dir + os.sep):
        return fp
    
    return None


@app.route("/api/download/<path:filename>")
def download(filename):
    # 修复 S1：download 路由必须强制登录认证，避免匿名用户遍历下载文件
    user = get_current_user()
    if not user:
        return jsonify({"error": "未登录"}), 401
    from urllib.parse import unquote, quote as _qd
    filename = unquote(filename)
    if not _is_safe_filename(filename):
        return jsonify({"error": "文件不存在"}), 404
    filepath = _find_user_file(filename)
    if not filepath:
        return jsonify({"error": "文件不存在"}), 404

    ext = os.path.splitext(filename)[1].lower()
    mimetype = MIME_MAP.get(ext, 'application/octet-stream')
    try:
        safe_ascii = filename.encode('ascii', errors='ignore').decode('ascii') or ('download' + ext)
    except Exception:
        safe_ascii = 'download' + ext
    try:
        resp = send_file(filepath, as_attachment=True, download_name=safe_ascii, mimetype=mimetype)
        resp.headers['X-Filename-Encoded'] = _qd(filename)
        resp.headers['X-Filename'] = safe_ascii
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as ee:
        print(f"[download] 发送文件异常: {ee}")
        import traceback; traceback.print_exc()
        return jsonify({"error": "文件下载失败"}), 500


@app.route("/api/avatar/<filename>")
def serve_avatar(filename):
    if not _is_safe_filename(filename):
        return jsonify({"error": "无效的文件"}), 400
    
    avatars_dir = os.path.realpath(os.path.join(DATA_DIR, 'uploads', 'avatars'))
    filepath = os.path.realpath(os.path.join(avatars_dir, filename))
    
    if not os.path.isfile(filepath) or not filepath.startswith(avatars_dir + os.sep):
        return jsonify({"error": "文件不存在"}), 404
    
    ext = os.path.splitext(filename)[1].lower()
    mimetype = MIME_MAP.get(ext, 'image/png')
    
    try:
        return send_file(filepath, mimetype=mimetype)
    except Exception as ee:
        print(f"[avatar] 发送文件异常: {ee}")
        return jsonify({"error": "图片加载失败"}), 500


@app.route("/api/save_blob", methods=["POST"])
def save_blob():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401
    
    import base64
    
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400
    
    filename = data.get('filename', 'download')
    blob_data = data.get('data')
    
    if not blob_data:
        return jsonify({"error": "缺少文件数据"}), 400
    
    safe_filename = _sanitize_filename(filename) if filename else 'download'
    if not safe_filename:
        safe_filename = 'download'
    # 修复 M11：在文件名前加用户ID前缀，避免不同用户同名文件相互覆盖/越权下载
    user_prefix = f"{user['id']}_"
    safe_filename = f"{user_prefix}{safe_filename}"
    filepath = os.path.join(str(VISUAL_OUTPUT_DIR), safe_filename)

    real_visual_dir = os.path.realpath(str(VISUAL_OUTPUT_DIR))
    real_filepath = os.path.realpath(filepath)
    if not real_filepath.startswith(real_visual_dir + os.sep):
        return jsonify({"error": "无效的文件名"}), 400

    try:
        decoded_data = base64.b64decode(blob_data)
        if len(decoded_data) > 10 * 1024 * 1024:
            return jsonify({"error": "文件大小不能超过10MB"}), 400
        with open(filepath, 'wb') as f:
            f.write(decoded_data)

        download_url = f'/api/download/{safe_filename}'
        return jsonify({"success": True, "url": download_url})
    except Exception as e:
        print(f"[save_blob] 保存文件异常: {e}")
        return jsonify({"error": "文件保存失败"}), 500


@app.route("/api/export/word", methods=["POST"])
def export_to_word():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    markdown_content = data.get("content", "").strip()
    title = str(data.get("title", "教案")).strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])

    try:
        word_buffer = markdown_to_word(markdown_content, title)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:20] or "untitled"
        filename = f"教案_{safe_title}_{timestamp}.docx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(word_buffer.getvalue())

        from urllib.parse import quote as _q
        from unicodedata import normalize as _norm
        try:
            _ascii_fn = filename.encode('ascii', errors='ignore').decode('ascii') or 'document.docx'
        except Exception:
            _ascii_fn = 'document.docx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn,
                         mimetype=MIME_MAP['.docx'])
        resp.headers['X-Filename-Encoded'] = _q(filename)
        resp.headers['X-Filename'] = _ascii_fn
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": "导出失败"}), 500

_PPT_RATE_LIMIT = {}
_PPT_RATE_LOCK = threading.Lock()


def _check_ppt_rate_limit(ip: str) -> bool:
    now = time.time()
    with _PPT_RATE_LOCK:
        if ip not in _PPT_RATE_LIMIT:
            _PPT_RATE_LIMIT[ip] = []
        attempts = _PPT_RATE_LIMIT[ip]
        attempts = [t for t in attempts if now - t < 300]
        _PPT_RATE_LIMIT[ip] = attempts
        if len(attempts) >= 10:
            return False
        _PPT_RATE_LIMIT[ip].append(now)
        # 修复 M3：定期清理空IP键，避免限流字典无界增长导致内存泄漏
        if len(_PPT_RATE_LIMIT) > 1000:
            expired_keys = [k for k, v in _PPT_RATE_LIMIT.items() if not v]
            for k in expired_keys:
                del _PPT_RATE_LIMIT[k]
    return True


def _cleanup_user_expired_files(user_output_dir, max_age_hours=24):
    try:
        if not os.path.exists(user_output_dir):
            return
        now = time.time()
        cutoff = now - max_age_hours * 3600
        for filename in os.listdir(user_output_dir):
            filepath = os.path.join(user_output_dir, filename)
            if os.path.isfile(filepath):
                try:
                    mtime = os.path.getmtime(filepath)
                    if mtime < cutoff and (filename.endswith('.docx') or filename.endswith('.pptx') or filename.endswith('.md')):
                        os.remove(filepath)
                except Exception:
                    pass
    except Exception:
        pass


@app.route("/api/export/ppt", methods=["POST"])
def export_to_ppt():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    ip = request.remote_addr or 'unknown'
    if request.headers.get('X-Forwarded-For'):
        ip = request.headers.get('X-Forwarded-For').split(',')[0].strip()

    if not _check_ppt_rate_limit(ip):
        return jsonify({"error": "请求过于频繁，请5分钟后再试"}), 429

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "无效请求"}), 400

    markdown_content = data.get("content", "")
    if not isinstance(markdown_content, str):
        return jsonify({"error": "内容格式错误"}), 400
    markdown_content = markdown_content.strip()

    title = data.get("title", "课件")
    if not isinstance(title, str):
        title = "课件"
    title = title.strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])
    _cleanup_user_expired_files(user_output_dir)

    ppt_buffer = None
    try:
        ppt_buffer = markdown_to_ppt(markdown_content, title)
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        safe_title = _sanitize_filename(title)[:20] or "untitled"
        filename = f"课件_{safe_title}_{timestamp}.pptx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(ppt_buffer.getvalue())

        from urllib.parse import quote as _q2
        try:
            _ascii_fn2 = filename.encode('ascii', errors='ignore').decode('ascii') or 'slides.pptx'
        except Exception:
            _ascii_fn2 = 'slides.pptx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn2,
                         mimetype=MIME_MAP['.pptx'])
        resp.headers['X-Filename-Encoded'] = _q2(filename)
        resp.headers['X-Filename'] = _ascii_fn2
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"[PPT] 导出失败: {e}", flush=True)
        return jsonify({"error": f"导出 PPT 失败: {e}"}), 500
    finally:
        if ppt_buffer is not None:
            try:
                ppt_buffer.close()
            except Exception:
                pass


@app.route("/api/lesson_plan/template_upload", methods=["POST"])
def api_lesson_plan_template_upload():
    cleanup_template_cache()
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    if 'file' not in request.files:
        return jsonify({"error": "未收到上传文件字段 file"}), 400

    f = request.files['file']
    if not f or not f.filename:
        return jsonify({"error": "文件名为空"}), 400

    filename = os.path.basename(f.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext != '.docx':
        return jsonify({"error": "仅支持 .docx 格式的教案模板"}), 400

    raw = f.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return jsonify({"error": f"文件过大 {len(raw)} bytes，上限 {MAX_UPLOAD_BYTES} bytes（20MB）"}), 400

    try:
        structure_map, pending_confirmations = template_filler.parse_docx_structure(raw)
    except Exception as ee:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"解析模板失败: {ee}", "error_code": "PARSE_FAILED"}), 500

    template_id = str(uuid.uuid4())
    with TEMPLATE_CACHE_LOCK:
        TEMPLATE_CACHE[template_id] = {
            'docx_bytes': raw,
            'structure_map': structure_map,
            'pending_confirmations': pending_confirmations,
            'filename': filename,
            'user_id': user['id'],
            'timestamp': time.time(),
        }

    regions_summary = []
    for region_id, region_info in structure_map.items():
        regions_summary.append({
            'region_id': region_id,
            'type': region_info['type'],
            'region_type': region_info['region_type'],
            'location': region_info.get('location', ''),
            'text_preview': region_info.get('text', '')[:100],
            'is_placeholder': region_info.get('is_placeholder', False),
            'needs_confirmation': region_info.get('needs_confirmation', False),
        })

    return jsonify({
        "ok": True,
        "template_id": template_id,
        "filename": filename,
        "structure_map": regions_summary,
        "pending_confirmations": pending_confirmations,
        "has_pending": len(pending_confirmations) > 0,
    })


@app.route("/api/lesson_plan/template_fill", methods=["POST"])
def api_lesson_plan_template_fill():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    template_id = data.get("template_id")
    if not template_id:
        return jsonify({"error": "缺少 template_id"}), 400

    with TEMPLATE_CACHE_LOCK:
        if template_id not in TEMPLATE_CACHE:
            return jsonify({"error": "模板不存在或已过期"}), 404
        template_data = TEMPLATE_CACHE[template_id]
        if template_data['user_id'] != user['id']:
            return jsonify({"error": "无权访问此模板"}), 403

    content_mapping = data.get("content_mapping", {})
    if not content_mapping:
        return jsonify({"error": "缺少 content_mapping"}), 400

    # 调试：打印 content_mapping 和 structure_map 的键
    print(f"[template_fill] content_mapping keys: {list(content_mapping.keys())[:5]}...")
    print(f"[template_fill] structure_map keys: {list(template_data['structure_map'].keys())[:5]}...")

    try:
        filled_doc = template_filler.fill_lesson_plan_template(
            template_data['docx_bytes'],
            content_mapping,
            template_data['structure_map'],
        )
    except Exception as ee:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"填充模板失败: {ee}", "error_code": "FILL_FAILED"}), 500

    output_filename = f"已填充_{template_data['filename']}"
    return send_file(
        filled_doc,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name=output_filename,
    )


@app.route("/api/lesson_plan/template_confirm", methods=["POST"])
def api_lesson_plan_template_confirm():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    template_id = data.get("template_id")
    if not template_id:
        return jsonify({"error": "缺少 template_id"}), 400

    confirmations = data.get("confirmations", {})
    with TEMPLATE_CACHE_LOCK:
        if template_id not in TEMPLATE_CACHE:
            return jsonify({"error": "模板不存在或已过期"}), 404
        template_data = TEMPLATE_CACHE[template_id]
        if template_data['user_id'] != user['id']:
            return jsonify({"error": "无权访问此模板"}), 403

        # 修复 M10：对 template_data 的 structure_map / pending_confirmations 修改必须在锁内，
        # 否则与 template_fill 并发执行会造成数据竞态
        for region_id, confirmed_type in confirmations.items():
            if region_id in template_data['structure_map']:
                template_data['structure_map'][region_id]['region_type'] = confirmed_type
                template_data['structure_map'][region_id]['needs_confirmation'] = False

        new_pending = [
            p for p in template_data['pending_confirmations']
            if p['region_id'] not in confirmations
        ]
        template_data['pending_confirmations'] = new_pending

    return jsonify({
        "ok": True,
        "has_pending": len(new_pending) > 0,
        "pending_count": len(new_pending),
    })


@app.route("/api/lesson_plan/template_generate", methods=["POST"])
def api_lesson_plan_template_generate():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data:
        return jsonify({"error": "请求数据为空"}), 400

    template_id = data.get("template_id")
    if not template_id:
        return jsonify({"error": "缺少 template_id"}), 400

    with TEMPLATE_CACHE_LOCK:
        if template_id not in TEMPLATE_CACHE:
            return jsonify({"error": "模板不存在或已过期"}), 404
        template_data = TEMPLATE_CACHE[template_id]
        if template_data['user_id'] != user['id']:
            return jsonify({"error": "无权访问此模板"}), 403

    user_prompt = data.get("prompt", "").strip()
    if not user_prompt:
        return jsonify({"error": "缺少生成需求"}), 400

    structure_map = template_data['structure_map']

    regions_info = []
    for region_id, region_info in structure_map.items():
        regions_info.append({
            'region_id': region_id,
            'region_type': region_info.get('region_type', '待确认'),
            'location': region_info.get('location', ''),
            'text_preview': region_info.get('text', '')[:100],
            'is_placeholder': region_info.get('is_placeholder', False),
        })

    ai_prompt = f"""
你是一位专业的教案生成专家。请根据用户的需求，为以下教案模板中的每个可填充区域生成相应的教学内容。

用户需求：{user_prompt}

模板区域列表：
{json.dumps(regions_info, ensure_ascii=False, indent=2)}

请严格按照以下格式输出（必须包裹在 ```json 代码块中）：
```json
{{
    "<region_id>": "<该区域的教学内容>",
    ...
}}
```

注意事项：
1. 每个区域的内容要符合其类型（如"教学目标"区域只生成目标内容）
2. 内容要详细、专业，符合中小学教学规范
3. 如果区域是占位符（is_placeholder: true），请生成完整的内容替换它
4. 如果区域已有文本（text_preview非空），请根据上下文生成补充或替换内容
5. 只输出上述JSON代码块，不要包含其他文字或解释
""".strip()

    try:
        # 修复 M6：补充 timeout，避免教案模板生成请求无限期挂起
        client = OpenAI(api_key=API_KEY, base_url=BASE_URL, timeout=60.0)
        response = client.chat.completions.create(
            model=MODEL,
            messages=[
                {"role": "system", "content": "你是一位专业的教案生成专家，擅长根据模板结构生成高质量的教学内容。"},
                {"role": "user", "content": ai_prompt},
            ],
            temperature=0.7,
            max_tokens=4096,
        )
    except Exception as e:
        print(f"[lesson_plan] AI 请求失败: {type(e).__name__}: {e}")
        return jsonify({"error": _format_ai_error(e)}), 500

    reply = response.choices[0].message.content

    print(f"[AI Reply] length={len(reply)}")
    print(f"[AI Reply] first 500 chars:\n{reply[:500]}")

    # 将AI响应保存到日志文件以便调试
    try:
        with open('ai_reply_debug.log', 'w', encoding='utf-8') as f:
            f.write(f"Timestamp: {datetime.datetime.now()}\n")  # 修复 B5：datetime 已是模块，应调用 datetime.datetime.now()
            f.write(f"Length: {len(reply)}\n")
            f.write(f"Content:\n{reply}\n")
    except Exception as e:
        print(f"[Debug] Failed to write debug log: {e}")

    def _extract_json_robust(text):
        """健壮地提取JSON，处理AI生成的不规范JSON"""
        import re
        # 1. 优先提取 ```json 代码块
        code_block_match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
        if code_block_match:
            json_str = code_block_match.group(1).strip()
        else:
            # 2. 回退：直接找第一个 {...}
            json_match = re.search(r'\{[\s\S]*\}', text)
            if json_match:
                json_str = json_match.group(0)
            else:
                raise ValueError("未找到JSON内容")

        print(f"[JSON Extracted] length={len(json_str)}")

        # 3. 规范化处理：清理不可见控制字符、统一引号
        def _normalize_json(s):
            # 移除不可见控制字符（保留换行和制表符）
            s = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', s)
            # 将中文引号替换为英文引号
            s = s.replace('“', '"').replace('”', '"').replace('‘', "'").replace('’', "'")
            # 移除尾随逗号
            s = re.sub(r',\s*([\]}])', r'\1', s)
            return s

        json_str = _normalize_json(json_str)

        # 4. 尝试使用 json5 解析（支持更宽松的JSON格式）
        try:
            import json5
            result = json5.loads(json_str)
            print(f"[JSON5 Parse] Success! keys={list(result.keys())[:5]}...")
            return result
        except Exception as e:
            print(f"[JSON5 Parse] Failed: {e}")
            # 打印出错位置附近的内容
            try:
                lines = json_str.split('\n')
                if hasattr(e, 'lineno'):
                    line_num = e.lineno
                else:
                    line_num = 74  # 常见错误位置
                start_line = max(0, line_num - 3)
                end_line = min(len(lines), line_num + 3)
                print(f"[JSON Error Context] lines {start_line+1}-{end_line+1}:")
                for i in range(start_line, end_line):
                    prefix = ">>>" if i == line_num - 1 else "   "
                    print(f"{prefix} {i+1}: {lines[i][:150]}")
            except Exception:  # 修复 L3：裸 except 改为 except Exception，避免吞掉 KeyboardInterrupt / SystemExit
                pass

        # 5. 尝试标准 json 解析
        try:
            result = json.loads(json_str)
            print(f"[JSON Parse] Success! keys={list(result.keys())[:5]}...")
            return result
        except json.JSONDecodeError as e:
            print(f"[JSON Parse] Failed: {e}")
            # 打印出错位置附近的内容
            lines = json_str.split('\n')
            start_line = max(0, e.lineno - 3)
            end_line = min(len(lines), e.lineno + 3)
            print(f"[JSON Error Context] lines {start_line+1}-{end_line+1}:")
            for i in range(start_line, end_line):
                prefix = ">>>" if i == e.lineno - 1 else "   "
                print(f"{prefix} {i+1}: {lines[i][:150]}")

        # 6. 修复常见JSON错误：值中未转义的双引号
        # 更健壮的引号处理：逐字段解析
        def _fix_json_with_quotes(s):
            result = {}
            # 按行分割处理
            lines = s.strip().split('\n')
            current_key = None
            current_value = []
            in_value = False
            value_quote_char = None

            for line in lines:
                line = line.strip()
                if not line or line.startswith('{') or line.startswith('}'):
                    continue

                # 检测键值对开始
                if ':' in line and not in_value:
                    # 匹配 "key": "value" 或 "key": "value
                    key_match = re.match(r'"([^"]+)"\s*:\s*["\']?(.*)$', line)
                    if key_match:
                        current_key = key_match.group(1)
                        rest = key_match.group(2).rstrip(',').rstrip()
                        # 检查值是否在本行结束
                        if rest.endswith('"') or rest.endswith("'"):
                            value_quote_char = rest[-1]
                            current_value = [rest[:-1]]
                            in_value = False
                            result[current_key] = ''.join(current_value)
                            current_key = None
                        else:
                            in_value = True
                            value_quote_char = '"'
                            current_value = [rest]
                    continue

                # 在值中
                if in_value and current_key:
                    stripped = line.rstrip(',').rstrip()
                    # 检查值是否结束
                    if stripped.endswith('"') or stripped.endswith("'"):
                        current_value.append(stripped[:-1])
                        in_value = False
                        result[current_key] = '\n'.join(current_value)
                        current_key = None
                    else:
                        current_value.append(stripped)

            return result

        fixed_result = _fix_json_with_quotes(json_str)
        if fixed_result:
            print(f"[Fixed Parse] Success! keys={list(fixed_result.keys())[:5]}...")
            return fixed_result

        # 7. 最后尝试：使用正则逐字段提取
        mapping = {}
        # 匹配 "key": "value" 模式（支持值中包含引号的情况）
        key_value_pattern = r'"([\w_]+)"\s*:\s*"((?:[^"\\]|\\.|[\s\S])*?)"(?=\s*(?:,|\}))'
        for match in re.finditer(key_value_pattern, json_str):
            key = match.group(1)
            value = match.group(2).replace('\\n', '\n').replace('\\"', '"').replace('\\\\', '\\')
            mapping[key] = value

        if mapping:
            print(f"[Regex Parse] Success! keys={list(mapping.keys())[:5]}...")
            return mapping

        raise ValueError("无法解析JSON内容")

    try:
        content_mapping = _extract_json_robust(reply)
    except Exception as e:
        # 解析完全失败时，尝试从文本中逐区域提取内容
        content_mapping = {}
        for region_id in structure_map:
            # 简单启发式：在回复中查找 region_id 后面的内容
            pattern = re.escape(region_id) + r'[\"\']?\s*[:：]\s*[\"\']?([^\"\'\n,}]{10,500})'
            m = re.search(pattern, reply)
            if m:
                content_mapping[region_id] = m.group(1).strip()

        if not content_mapping:
            return jsonify({
                "error": f"解析AI响应失败: {str(e)}",
                "raw_reply": reply[:3000],
            }), 500

    # 兜底：AI 可能没返回所有区域，用默认内容补全
    default_contents = {
        '教学目标': '知识与技能：掌握本课核心知识点\n过程与方法：通过合作探究提升学习能力\n情感态度与价值观：培养学习兴趣与科学素养',
        '教学重难点': '教学重点：本课核心概念与原理\n教学难点：知识的综合应用与迁移',
        '学情分析': '本班学生基础知识较为扎实，具备一定的自主学习能力，但在知识的综合运用方面仍需加强。',
        '教学方法': '讲授法、讨论法、探究法、情境教学法相结合',
        '教学过程': '详见教学过程设计',
        '导入': '通过情境导入，激发学生学习兴趣，引出本课主题。',
        '新授': '教师引导学生自主探究新知，逐步掌握核心内容。',
        '巩固练习': '通过课堂练习巩固所学知识，及时反馈。',
        '小结': '师生共同总结本课重点内容，梳理知识体系。',
        '作业布置': '完成课后配套练习，预习下一课内容。',
        '板书设计': '详见板书设计',
        '教学反思': '课后根据学生学习情况进行反思与改进。',
        '时间分配': '导入5分钟｜新授20分钟｜练习10分钟｜小结5分钟',
        '教具准备': '多媒体课件、教材、黑板',
        '教学媒体': '多媒体课件、投影仪',
        '基本信息': '学科：\t年级：\t课题：',
        '待确认': '（内容由教师补充）',
    }
    for region_id, region_info in structure_map.items():
        if region_id not in content_mapping:
            rtype = region_info.get('region_type', '待确认')
            content_mapping[region_id] = default_contents.get(rtype, '（AI 未生成此区域内容，请手动补充）')

    return jsonify({
        "ok": True,
        "content_mapping": content_mapping,
        "total_regions": len(content_mapping),
    })


MODULES_DATA = {
    "ppt_generate": {
        "name": "PPT生成",
        "emoji": "📊",
        "short_desc": "一键生成教学课件PPT",
        "description": "AI智能生成完整教学课件，支持从课程主题直接输出结构化PPT大纲，并可一键导出为 .pptx 文件。自动包含封面、教学目标、知识点讲解、案例分析、课堂小结、课后作业等标准课件结构。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "生成一份《荷塘月色》的语文课件PPT，包含教学目标、作者介绍、课文赏析、重点语句分析、课堂练习、作业布置，共15-20页",
        "features": [
            "📊 标准课件结构（封面/目标/新知/小结/作业）",
            "📝 每页幻灯片内容完整撰写（标题+要点）",
            "🎨 Markdown格式规范排版，方便二次编辑",
            "⬇️ 一键导出 .pptx 文件，可直接打开使用",
            "🔧 自动适配学科特点（文科赏析/理科例题/工科原理）",
            "⏱ 支持指定页数与各环节时间分配"
        ],
        "examples": [
            {"title": "语文课件：荷塘月色", "desc": "朱自清散文，含作者背景+课文赏析+修辞分析", "prompt": "生成一份《荷塘月色》的语文课件PPT，适合高中，包含教学目标、作者介绍、课文逐段赏析、修辞手法分析、课堂练习、作业布置，共15-20页"},
            {"title": "数学课件：二次函数", "desc": "初中数学，含图像性质+典型例题+课堂练习", "prompt": "生成初中数学《二次函数的图像与性质》课件PPT，包含概念讲解、5个典型例题解析、课堂练习、课后作业，约18页"},
            {"title": "历史课件：辛亥革命", "desc": "高中历史，含背景脉络+重大事件+意义启示", "prompt": "生成高中历史《辛亥革命》课件PPT，包含历史背景、武昌起义经过、中华民国成立、历史意义与局限、思考题，约20页"}
        ],
        "knowledge_points": ["课件结构", "教学目标", "知识点导入", "案例分析", "课堂互动", "例题解析", "归纳小结", "作业布置", "板书设计", "时间分配"]
    },
    "ppt_outline": {
        "name": "PPT讲解纲要",
        "emoji": "🎤",
        "short_desc": "每页PPT的逐页讲解稿",
        "description": "根据PPT大纲为每一页幻灯片生成详细的教师讲解词。包含开场白设计、知识点过渡衔接语、重难点突出提示、提问互动设计、随堂练习引导语等，帮助教师脱稿流畅完成整堂课。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、互动提示、过渡语句、重点强调部分",
        "features": [
            "🎤 逐页讲解词撰写，教师可直接照着读",
            "🔗 页面之间设计自然衔接过渡语句",
            "❓ 穿插课堂提问与互动设计（教师问+预期答）",
            "⭐ 标注重点知识点的语气语调提示（重读、停顿）",
            "⏰ 每页推荐讲解时长，合理分配45分钟",
            "🧑‍🏫 模拟真实课堂的教态与现场感"
        ],
        "examples": [
            {"title": "语文：荷塘月色讲解稿", "desc": "逐页教师用语+情感渲染引导", "prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、课堂提问设计、过渡语句、重点部分标注语气要求"},
            {"title": "数学：二次函数讲解稿", "desc": "概念引入+例题推导步骤+易错提醒", "prompt": "生成初中数学《二次函数图像与性质》每一页PPT的教师讲解纲要，包含概念引入话术、例题逐步推导讲解、学生易错点提醒、课堂练习的处理方式"},
            {"title": "班会：防溺水安全教育稿", "desc": "案例警示+互动问答+总结升华", "prompt": "生成中小学防溺水主题班会PPT的逐页讲解纲要，包含真实案例描述、互动提问、自救技巧讲解、最后的倡议环节，语言符合中小学生认知"}
        ],
        "knowledge_points": ["开场白设计", "过渡衔接", "重难点提示", "互动提问", "案例讲述", "语气语调", "时间把控", "课堂节奏", "情感升华", "结尾总结"]
    },
    "lesson_plan": {
        "name": "教案生成",
        "emoji": "📘",
        "short_desc": "标准格式45分钟完整教案",
        "description": "按中小学教案标准格式自动生成完整45分钟教学设计。包含教学目标（三维/核心素养）、教学重难点、教学方法、教学过程（导入/新授/练习/小结/作业）、板书设计、教学反思模板等全部要素。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "中职", "高职"],
        "default_prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程、板书设计、作业布置、教学反思",
        "features": [
            "🎯 三维教学目标（知识与技能/过程与方法/情感态度）",
            "🎯 核心素养目标适配新课标要求",
            "❗ 教学重点与教学难点明确区分",
            "⏱️ 详细教学过程：5-6个环节精确到分钟",
            "🎨 板书设计框架（主板书+副板书）",
            "📋 配套作业设计（基础+拓展分层）",
            "📝 预留教学反思栏（课堂生成问题记录）"
        ],
        "examples": [
            {"title": "语文：背影教案", "desc": "朱自清经典·情感教育·细节描写", "prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程（各环节含时间分配）、板书设计、分层作业布置、教学反思栏"},
            {"title": "物理：浮力教案", "desc": "实验探究·阿基米德原理·分层练习", "prompt": "生成初中物理《阿基米德原理》45分钟教案，包含实验探究环节设计、演示实验步骤、学生分组活动安排、分层练习设计、板书设计"},
            {"title": "英语：Travel Plan教案", "desc": "听说课型·情境交际·任务型教学", "prompt": "生成初中英语听说课《Travel Plans》完整教案，包含Warm-up、Pre-listening、While-listening、Post-speaking、Summary & Homework环节，配师生对话示例"}
        ],
        "knowledge_points": ["三维目标", "核心素养", "教学重难点", "学情分析", "教学方法", "教学过程", "时间分配", "板书设计", "分层作业", "教学反思"]
    },
    "exercises": {
        "name": "习题生成",
        "emoji": "📝",
        "short_desc": "分层习题+详细答案解析",
        "description": "根据知识点自动生成多难度层次练习题，支持选择题、填空题、判断题、解答题、应用题等多种题型。每道题目均附标准答案与详细解题步骤，可用于课堂练习、课后作业、单元测验。",
        "subject": "全学科",
        "grade_levels": ["小学低年级", "小学高年级", "初中", "高中", "中职", "高职"],
        "default_prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用，所有题目附答案和详细解析",
        "features": [
            "📊 三级难度分层：基础巩固·能力提高·拓展探究",
            "✅ 题型丰富：选择/填空/判断/解答/实验探究/应用题",
            "📝 每题附标准答案 + 详细解题步骤",
            "📌 标注每题考查知识点对应章节",
            "⏱ 推荐完成时长与分值设置",
            "📄 可一键组装为标准试卷（卷头+姓名栏）"
        ],
        "examples": [
            {"title": "数学：勾股定理20题", "desc": "8基础+8提高+4拓展，附详解", "prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用题，所有题目附答案和详细解析步骤"},
            {"title": "英语：一般过去时语法", "desc": "30道选择+填空+句型转换，解析全", "prompt": "生成初中英语语法《一般过去时》专项练习题30道：单项选择15道，用所给动词适当形式填空10道，句型转换5道，附答案和详细解析"},
            {"title": "化学：酸碱盐专题", "desc": "选择+填空+推断+实验探究四大题型", "prompt": "生成初中化学《酸碱盐》单元练习题，包含选择题10道、填空题6道、物质推断题2道、实验探究题2道，附答案与解析"}
        ],
        "knowledge_points": ["基础巩固", "能力提高", "拓展探究", "选择题", "填空题", "解答题", "实验探究", "应用题", "答案解析", "考点分布"]
    },
    "homework_analysis": {
        "name": "作业分析",
        "emoji": "📊",
        "short_desc": "作业批改与学情数据分析",
        "description": "上传作业统计数据或描述作业情况，AI自动生成完整的作业分析报告。包含：正确率统计、高频错题归因、典型错解展示、知识点掌握情况雷达图、学困生针对性辅导建议、下节课教学调整策略等。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "default_prompt": "请为一次高中数学《函数单调性》课后作业生成作业分析报告：全班45人，平均得分率68%，错误集中在：含参数的单调性讨论、复合函数单调性、实际应用最值问题，请给出详细分析与辅导建议",
        "features": [
            "📈 整体成绩统计：平均分/得分率/分数段分布",
            "❌ 高频错题 TOP 5 归因分析（概念/审题/计算/方法）",
            "📝 典型错解案例展示与正解对比",
            "🎯 知识点掌握度雷达图（熟练/一般/薄弱）",
            "👨‍🎓 分层辅导建议：学困生、中等生、优等生",
            "✏️ 下一步教学调整：需补讲的知识点、课堂策略",
            "📑 讲评课时分配建议与讲评例题推荐"
        ],
        "examples": [
            {"title": "数学函数单调性作业分析", "desc": "得分率68%·参数讨论错误高发", "prompt": "为高中数学《函数单调性》课后作业生成详细作业分析报告：全班45人，平均得分率68%，主要错误：①含参数单调性分类讨论（失分率52%）②复合函数同增异减应用（失分率41%）③实际问题求最值忽略定义域（失分率37%），请给出错因分析、辅导建议、教学调整策略"},
            {"title": "英语完形填空作业分析", "desc": "上下文逻辑·固定搭配是薄弱点", "prompt": "分析一次八年级英语完形填空作业情况：全班40人，平均正确率58%；错题中上下文逻辑理解占45%，固定搭配与词组占30%，词汇辨析占15%，语法占10%。请生成完整作业讲评方案"},
            {"title": "物理电路作业错误分析", "desc": "串并联识别+欧姆定律综合应用薄弱", "prompt": "为初三物理《欧姆定律在串并联电路中的应用》作业做分析：全班48人，平均得分62分；常见错：复杂电路等效化简错误、电表测量对象判断错误、比例计算错误。生成作业分析与讲评建议"}
        ],
        "knowledge_points": ["得分率统计", "错题归因", "典型错解", "知识薄弱点", "分层辅导", "学困生帮扶", "教学调整", "讲评课时", "补充练习", "家校沟通"]
    },
    "history": {
        "name": "历史记录",
        "emoji": "🕒",
        "short_desc": "查看所有AI生成记录与历史会话",
        "description": "查看当前账号下所有历史生成记录：历史会话聊天、PPT课件、教案、习题、作业分析等全部生成内容。支持按时间、按模块筛选，快速找到过往生成内容，可再次打开、复用、编辑或重新生成。",
        "subject": "系统功能",
        "grade_levels": ["全部学段"],
        "default_prompt": "（点击下方会话卡片直接查看历史记录）",
        "features": [
            "🕒 按时间倒序展示所有历史会话",
            "🔍 支持按模块类型筛选（PPT/教案/习题/作业等）",
            "📋 可快速查看每条记录的生成时间、标题、内容摘要",
            "♻️ 一键复用：把历史内容重新发送，基于结果再修改",
            "🗂 可导出历史记录为 Word / 文本文件",
            "📊 查看个人使用统计（生成次数、常用模块等）"
        ],
        "examples": [
            {"title": "查看最近7天生成记录", "desc": "快速找到上周生成的PPT并修改", "prompt": "__HISTORY__VIEW__"},
            {"title": "按模块筛选：只看教案", "desc": "查找所有历史生成的教案", "prompt": "__HISTORY_VIEW__ lesson_plan"},
            {"title": "统计本月使用情况", "desc": "查看自己备课模块使用分布", "prompt": "__HISTORY_STATS__"}
        ],
        "knowledge_points": ["聊天记录", "生成文件", "按时间筛选", "按模块筛选", "内容搜索", "记录复用", "批量导出", "使用统计", "收藏夹", "回收站"]
    },
    "visualization": {
        "name": "可视化生成",
        "emoji": "🎬",
        "short_desc": "教学动画与思维导图生成",
        "description": "输入课程知识点或教学目标，AI 自动为您生成：\n✅ HTML5教学演示动画（CSS动态效果、无需安装依赖、可嵌入PPT）\n✅ 一目了然的课程思维导图（交互式、可缩放、关键词精炼）\n支持小学至高职全学段，覆盖数理化生、语文英语等主流学科。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "default_prompt": "请为高中物理《牛顿第一定律》新课教学生成动画和思维导图。",
        "features": [
            "🎬 HTML5教学动画 - 纯CSS动画效果，无需安装Manim/FFmpeg等依赖，点击即生成，支持重播",
            "🧠 交互式思维导图 - markmap渲染，支持缩放、节点展开折叠，自动配色",
            "📐 AI智能布局 - 根据学科内容自动生成动画元素位置、颜色、文字和动画效果",
            "⚡ 极速生成 - 无需等待视频渲染，秒级生成结果",
            "🎨 双栏布局 - 动画和思维导图同屏展示，方便教学使用"
        ],
        "examples": [
            {"title": "牛顿第一定律", "desc": "动画展示物体静止/匀速运动状态；思维导图含定义/条件/推论/应用", "prompt": "牛顿第一定律"},
            {"title": "光合作用", "desc": "动画展示光反应→暗反应过程；导图分光反应/暗反应/影响因素/应用", "prompt": "光合作用"},
            {"title": "勾股定理", "desc": "动画展示直角三角形三边关系；导图含定理/证明/应用/拓展", "prompt": "勾股定理"}
        ],
        "knowledge_points": ["教学动画", "思维导图", "PPT素材", "知识点可视化", "HTML5动画", "Markmap导图", "学科建模", "课堂演示", "教师提效"]
    },
    "exam_paper": {
        "name": "试卷生成",
        "emoji": "📝",
        "short_desc": "AI智能生成试卷，支持多题型、难度分级",
        "description": "AI 智能生成试卷，支持多学科、多题型、难度分级",
        "type": "exam",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "features": [
            "📚 多学科支持 - 语文、数学、英语、物理、化学等全学科覆盖",
            "📋 丰富题型 - 选择题、填空题、判断题、解答题等多种题型",
            "🎯 难度分级 - 简单/中等/困难三档难度自由选择",
            "📊 智能配分 - 自动根据题型和题量分配分值",
            "✅ 答案解析 - 自动生成详细答案和解题思路",
            "📥 一键导出 - 支持导出Word文档和PDF，方便打印使用"
        ],
        "examples": [
            {"title": "一元二次方程", "desc": "初中数学章节测试卷，含选择/填空/解答题", "prompt": "一元二次方程"},
            {"title": "牛顿运动定律", "desc": "高一物理单元测验，含答案解析", "prompt": "牛顿运动定律"},
            {"title": "一般过去时", "desc": "初中英语语法专项试卷", "prompt": "一般过去时"}
        ],
        "knowledge_points": ["试卷生成", "智能出题", "答案解析", "章节测试", "单元测验", "期中期末", "模拟考试", "教师备课", "自动组卷"]
    }
}

# 注册试卷生成模块路由
register_exam_routes(
    app=app,
    get_current_user_func=get_current_user,
    database_mod=database,
    api_key=API_KEY,
    base_url=BASE_URL,
    model_name=MODEL,
    get_user_output_dir_func=get_user_output_dir,
)

# 注册可视化生成模块路由
def _get_openai_client():
    from openai import OpenAI
    return OpenAI(api_key=API_KEY, base_url=BASE_URL)

register_visual_routes(
    app=app,
    get_current_user=get_current_user,
    safe_get_json=safe_get_json,
    database=database,
    openai_client_factory=_get_openai_client,
    model_name=MODEL,
    modules_data=MODULES_DATA,
    qwen_image_api_key=QWEN_IMAGE_API_KEY,
    image_search_provider=IMAGE_SEARCH_PROVIDER,
    image_search_api_key=IMAGE_SEARCH_API_KEY,
)

PUBLIC_MODULES = {"exam_paper"}

@app.route("/module/<module_id>")
def module_detail(module_id):
    user = get_current_user()
    if not user and module_id not in PUBLIC_MODULES:
        print(f"[module_detail] 用户未登录，重定向到登录页")
        return redirect(url_for('login', next=request.path))
    if user:
        print(f"[module_detail] 用户: {user['username']}, module_id: {module_id}")
    else:
        print(f"[module_detail] 公开模块免登录访问: {module_id}")

    if module_id not in MODULES_DATA:
        print(f"[module_detail] 模块不存在: {module_id}")
        return redirect(url_for('index'))
    module = MODULES_DATA[module_id]

    if module_id == "visualization":
        try:
            all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
            print(f"[module_detail] 渲染可视化页面，模块名称: {module['name']}")
            username_display = user['username'] if user else '游客'
            return render_template(
                "module_visual.html",
                module=module,
                module_id=module_id,
                model=MODEL,
                username=username_display,
                all_modules=all_modules
            )
        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"[module_detail] 渲染可视化页面失败: {e}")
            return f"渲染失败", 500

    if module_id == "exam_paper":
        try:
            all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
            print(f"[module_detail] 渲染试卷生成页面，模块名称: {module['name']}")
            username_display = user['username'] if user else '游客'
            return render_template(
                "module_exam.html",
                module=module,
                module_id=module_id,
                model=MODEL,
                username=username_display,
                all_modules=all_modules
            )
        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"[module_detail] 渲染试卷生成页面失败: {e}")
            return f"渲染失败", 500

    all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
    history_sessions = []
    if module_id == "history":
        try:
            MODULE_META = {
                "ppt_generate":         {"name": "PPT生成",     "emoji": "📊", "color": "#667eea", "gradient": "linear-gradient(135deg,#667eea 0%,#764ba2 100%)"},
                "ppt_outline":          {"name": "PPT讲解纲要", "emoji": "📝", "color": "#f093fb", "gradient": "linear-gradient(135deg,#f093fb 0%,#f5576c 100%)"},
                "lesson_plan_generate": {"name": "教案生成",     "emoji": "📘", "color": "#4facfe", "gradient": "linear-gradient(135deg,#4facfe 0%,#00f2fe 100%)"},
                "exercise_generate":    {"name": "习题生成",     "emoji": "📚", "color": "#43e97b", "gradient": "linear-gradient(135deg,#43e97b 0%,#38f9d7 100%)"},
                # 修复 B6：补充与 MODULES_DATA 一致的键名（lesson_plan / exercises），避免历史记录匹配回退到 general
                "lesson_plan":           {"name": "教案生成",     "emoji": "📘", "color": "#4facfe", "gradient": "linear-gradient(135deg,#4facfe 0%,#00f2fe 100%)"},
                "exercises":             {"name": "习题生成",     "emoji": "📚", "color": "#43e97b", "gradient": "linear-gradient(135deg,#43e97b 0%,#38f9d7 100%)"},
                "homework_analysis":    {"name": "作业分析",     "emoji": "📊", "color": "#fa709a", "gradient": "linear-gradient(135deg,#fa709a 0%,#fee140 100%)"},
                "history":              {"name": "历史记录",     "emoji": "📜", "color": "#a8edea", "gradient": "linear-gradient(135deg,#a8edea 0%,#fed6e3 100%)"},
                "visualization":        {"name": "可视化生成",   "emoji": "🎬", "color": "#8b5cf6", "gradient": "linear-gradient(135deg,#8b5cf6 0%,#ec4899 100%)"},
                "general":              {"name": "自由聊天",     "emoji": "💬", "color": "#c3cfe2", "gradient": "linear-gradient(135deg,#c3cfe2 0%,#f5f7fa 100%)"},
            }
            stats = database.get_user_sessions_with_stats(user['id'])
            for s in stats:
                mid = s.get("module_id") or "general"
                if mid not in MODULE_META:
                    mid = "general"
                meta = MODULE_META[mid]
                files = s.get("files") or []
                badge_summary = {}
                for f in files:
                    k = f.get("kind") or "文件"
                    badge_summary.setdefault(k, {"count": 0, "samples": []})
                    if badge_summary[k]["count"] < 3:
                        badge_summary[k]["samples"].append(f.get("name") or "")
                    badge_summary[k]["count"] += 1
                badges = []
                for k, v in badge_summary.items():
                    ext_icon = {"PPT": "📊", "教案": "📘", "练习册": "📚", "作业分析": "📈", "报告": "📄", "文件": "📄"}.get(k, "📄")
                    badges.append({
                        "kind": k,
                        "icon": ext_icon,
                        "count": v["count"],
                        "sample": (v["samples"][0] or "")[:50]
                    })
                history_sessions.append({
                    "id": s["session_id"],
                    "pk": s.get("pk_id") or "",
                    "title": s.get("title") or "(空对话)",
                    "created_at": s.get("created_at") or "",
                    "updated_at": s.get("updated_at") or "",
                    "time": s.get("updated_at") or s.get("created_at") or "",
                    "msg_count": s.get("total_count") or 0,
                    "user_count": s.get("user_count") or 0,
                    "assistant_count": s.get("assistant_count") or 0,
                    "file_count": s.get("file_count") or 0,
                    "module_id": mid,
                    "module_name": meta["name"],
                    "module_emoji": meta["emoji"],
                    "module_color": meta["color"],
                    "module_gradient": meta["gradient"],
                    "preview": s.get("user_preview") or s.get("ai_preview") or "",
                    "user_preview": s.get("user_preview") or "",
                    "ai_preview": s.get("ai_preview") or "",
                    "badges": badges,
                    "files": files[:8],
                })
            history_sessions.sort(key=lambda x: x.get("time") or "", reverse=True)
        except Exception as e:
            import traceback; traceback.print_exc()
            print("[history] 读取会话失败：", e)
    return render_template(
        "module_detail.html",
        module=module,
        module_id=module_id,
        model=MODEL,
        username=user['username'],
        all_modules=all_modules,
        history_sessions=history_sessions
    )



def get_client_ip():
    return _get_client_ip()


def is_admin_logged_in():
    """管理员是否已登录：通过服务端 admin_sessions 表反查会话有效性"""
    token = request.cookies.get('admin_session')
    if not token:
        return False
    session = database.get_admin_session(token)
    return bool(session and session.get('is_active'))


def admin_identity():
    """从服务端会话反查当前管理员登录账号。
    取代旧实现直接信任客户端 admin_user cookie 的做法（已知伪造风险）。
    """
    token = request.cookies.get('admin_session')
    if not token:
        return 'unknown'
    session = database.get_admin_session(token)
    if not session:
        return 'unknown'
    return session.get('username') or 'unknown'


def _current_admin_session():
    """获取当前管理员会话详情（用于路由中读取 admin_id/is_super 等字段）"""
    token = request.cookies.get('admin_session')
    if not token:
        return None
    return database.get_admin_session(token)


def super_admin_required(f):
    """仅超管可访问：用于管理员账号管理接口"""
    @wraps(f)
    def decorated(*args, **kwargs):
        session = _current_admin_session()
        if not session:
            return jsonify({"error": "未授权"}), 401
        if not session.get('is_super'):
            return jsonify({"error": "仅超级管理员可执行此操作"}), 403
        return f(*args, **kwargs)
    return decorated


def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not is_admin_logged_in():
            return jsonify({"error": "未授权"}), 401
        return f(*args, **kwargs)
    return decorated


@app.route("/admin")
def admin_login_page():
    if is_admin_logged_in():
        return redirect(url_for('admin_dashboard'))
    return render_template("admin_login.html")


@app.route("/admin/login", methods=["POST"])
def admin_login():
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    username = (data.get("username") or "").strip()
    password = data.get("password") or ""
    ip = get_client_ip()

    if not username or not password:
        return jsonify({"error": "用户名和密码不能为空"}), 400

    if not _check_login_rate_limit(ip, username):
        return jsonify({"error": "登录尝试过多，请15分钟后再试"}), 429

    result = database.verify_admin_password(username, password)
    if result.get("success"):
        admin = result["admin"]
        _clear_login_failures(ip, username)
        session_token = database.create_admin_session(
            admin["id"], ip_address=ip,
            user_agent=request.headers.get('User-Agent', '')[:200]
        )
        if not session_token:
            return jsonify({"error": "登录失败，会话创建异常"}), 500
        database.record_login_attempt(ip, username, success=True)
        database.log_admin_action(admin["username"], "login",
                                   details=f"管理员登录 display={admin['display_name']}",
                                   ip_address=ip)
        # 顺便清理过期会话（低频清理，避免无界增长）
        try:
            database.cleanup_expired_admin_sessions()
        except Exception:
            pass
        response = jsonify({"status": "success", "username": admin["display_name"]})
        _set_admin_auth_cookies(response, session_token)
        return response
    else:
        _record_login_failure(ip, username)
        database.record_login_attempt(ip, username, success=False)
        # 不暴露具体原因（账号不存在 vs 密码错误），统一返回"用户名或密码错误"
        safe_error = "用户名或密码错误" if result.get("reason") in ("not_found", "wrong_password", "corrupt") else result.get("error", "登录失败")
        return jsonify({"error": safe_error}), 401


@app.route("/admin/dashboard")
def admin_dashboard():
    if not is_admin_logged_in():
        return redirect(url_for('admin_login_page'))
    user = get_current_user()
    return render_template("admin_users.html", user=user)


@app.route("/admin/logout", methods=["GET", "POST"])
def admin_logout():
    admin = admin_identity()
    token = request.cookies.get('admin_session')
    if token:
        database.delete_admin_session(token)
    database.log_admin_action(admin, "logout", details="管理员退出登录", ip_address=get_client_ip())
    if request.method == 'GET':
        response = make_response(redirect('/admin/login'))
    else:
        response = make_response(jsonify({"status": "success"}))
    _clear_auth_cookies(response)
    return response


@app.route("/admin/api/stats", methods=["GET"])
@admin_required
def admin_get_stats():
    try:
        user_stats = database.get_user_stats()
        system_stats = database.get_system_stats()
        activities = database.get_recent_activities(limit=10)
    except Exception as e:
        return jsonify({"success": False, "error": "统计数据获取失败"}), 500

    return jsonify({
        "success": True,
        "user_stats": user_stats,
        "system_stats": system_stats,
        "recent_activities": activities
    })


@app.route("/admin/api/users", methods=["GET"])
@admin_required
def admin_get_users_v2():
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 20, type=int)
    sort_by = request.args.get('sort_by', 'created_at')
    order = request.args.get('order', 'DESC')
    search = request.args.get('search', '').strip()

    if page < 1:
        page = 1
    if per_page < 5:
        per_page = 5
    if per_page > 100:
        per_page = 100

    result = supabase_client.get_all_users()
    if result["success"]:
        users = result["users"]
        # 客户端分页、搜索、排序
        if search:
            kw = search.lower()
            users = [u for u in users if kw in (u.get('username') or '').lower() or kw in (u.get('email') or '').lower()]
        users = sorted(users, key=lambda u: u.get(sort_by) or '', reverse=(order.upper() == 'DESC'))
        total = len(users)
        start = (page - 1) * per_page
        end = start + per_page
        return jsonify({
            "success": True,
            "users": users[start:end],
            "total": total,
            "page": page,
            "per_page": per_page,
            "source": "supabase"
        })
    else:
        local = database.get_all_local_users_paginated(page, per_page, sort_by, order, search)
        return jsonify({
            "success": True,
            "users": local["users"],
            "total": local["total"],
            "page": local["page"],
            "per_page": local["per_page"],
            "source": "local",
            "warning": "未配置 Supabase service_role key，当前显示本地数据库中的用户数据"
        })


@app.route("/admin/api/users/export", methods=["GET"])
@admin_required
def admin_export_users():
    result = supabase_client.get_all_users()
    if result["success"]:
        users = result["users"]
        source = "supabase"
    else:
        users = database.get_all_local_users()
        source = "local"

    # 修复 M12：CSV 注入防护，对以 = + - @ 开头的单元格加前缀 '，避免 Excel 公式注入
    def _csv_escape(value):
        s = str(value) if value is not None else ''
        if s and s[0] in ('=', '+', '-', '@'):
            return "'" + s
        return s

    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["用户ID", "用户名", "邮箱", "角色", "创建时间", "最后登录", "数据源"])
    for u in users:
        writer.writerow([
            _csv_escape(u.get('id', '')),
            _csv_escape(u.get('username', '')),
            _csv_escape(u.get('email', '')),
            _csv_escape('管理员' if u.get('role') == 'admin' else '普通用户'),
            _csv_escape(u.get('created_at', '')),
            _csv_escape(u.get('last_sign_in_at', '')),
            _csv_escape(source)
        ])

    csv_bytes = output.getvalue().encode('utf-8-sig')
    filename = f"users_export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
    return Response(
        csv_bytes,
        mimetype='text/csv; charset=utf-8-sig',
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@app.route("/admin/api/users", methods=["POST"])
@admin_required
def admin_create_user():
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    email = data.get("email")
    password = data.get("password")
    username = data.get("username")
    role = data.get("role", "authenticated")

    if not email or not password:
        return jsonify({"error": "邮箱和密码不能为空"}), 400

    if len(password) < 6:
        return jsonify({"error": "密码长度至少6位"}), 400

    result = supabase_client.admin_create_user(email, password, username, role)
    if result["success"]:
        user = result["user"]
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        try:
            cursor.execute('INSERT OR REPLACE INTO users (user_id, username, email) VALUES (?, ?, ?)',
                          (user['id'], user['username'] or email.split('@')[0], user['email']))
            conn.commit()
        finally:
            conn.close()
        database.log_admin_action(
            admin_identity(), "create_user", target_user_id=user.get('id'),
            details=f"创建用户 {email} 角色 {role}", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "user": user})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500


@app.route("/admin/api/user/<user_id>", methods=["GET"])
@admin_required
def admin_get_user_detail(user_id):
    # 优先尝试 Supabase，失败则使用本地详情
    result = supabase_client.get_user_by_id(user_id)
    if result["success"]:
        return jsonify({"success": True, "user": result["user"], "source": "supabase"})

    local = database.get_user_detail_with_sessions(user_id)
    if local:
        return jsonify({"success": True, "user": local, "source": "local"})

    return jsonify({"success": False, "error": result["error"]}), 404


@app.route("/admin/api/users/<user_id>", methods=["DELETE"])
@admin_required
def admin_delete_user(user_id):
    # 防止管理员误删自己：比对用户名/邮箱
    admin_name = admin_identity()
    user_result = supabase_client.get_user_by_id(user_id)
    if user_result.get("success"):
        target_identity = user_result["user"].get("username") or user_result["user"].get("email") or ''
    else:
        local_user = database.get_user_detail_with_sessions(user_id)
        target_identity = local_user.get("username") or local_user.get("email") or '' if local_user else ''
    if target_identity and target_identity == admin_name:
        return jsonify({"success": False, "error": "不能删除当前登录的管理员账号"}), 400

    result = supabase_client.delete_user(user_id)
    if result["success"]:
        database.delete_local_user(user_id)
        database.log_admin_action(
            admin_identity(), "delete_user", target_user_id=user_id,
            details="删除用户", ip_address=get_client_ip()
        )
        return jsonify({"success": True})
    else:
        local_deleted = database.delete_local_user(user_id)
        if local_deleted:
            database.log_admin_action(
                admin_identity(), "delete_user", target_user_id=user_id,
                details="仅从本地数据库删除用户", ip_address=get_client_ip()
            )
            return jsonify({"success": True, "warning": "仅从本地数据库删除，Supabase中的用户未删除"})
        else:
            return jsonify({"success": False, "error": result["error"]}), 500


def _is_current_admin_user(user_id, admin_name):
    """检查 user_id 是否对应当前管理员账号"""
    if not admin_name:
        return False
    user_result = supabase_client.get_user_by_id(user_id)
    if user_result.get("success"):
        target_identity = user_result["user"].get("username") or user_result["user"].get("email") or ''
    else:
        local_user = database.get_user_detail_with_sessions(user_id)
        target_identity = local_user.get("username") or local_user.get("email") or '' if local_user else ''
    return target_identity and target_identity == admin_name


@app.route("/admin/api/users/batch", methods=["POST"])
@admin_required
def admin_batch_operate():
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    action = data.get("action")
    user_ids = data.get("user_ids", [])
    if not user_ids or not isinstance(user_ids, list):
        return jsonify({"error": "请选择要操作的用户"}), 400

    admin_name = admin_identity()
    # 过滤掉当前管理员自身
    user_ids = [uid for uid in user_ids if not _is_current_admin_user(uid, admin_name)]

    success_count = 0
    failed = []

    if action == "delete":
        for uid in user_ids:
            try:
                res = supabase_client.delete_user(uid)
                if res["success"]:
                    database.delete_local_user(uid)
                    success_count += 1
                else:
                    local_deleted = database.delete_local_user(uid)
                    if local_deleted:
                        success_count += 1
                    else:
                        failed.append({"id": uid, "error": res.get("error", "删除失败")})
            except Exception as e:
                failed.append({"id": uid, "error": str(e)})
        database.log_admin_action(
            admin_identity(), "batch_delete", details=f"批量删除 {success_count} 个用户",
            ip_address=get_client_ip()
        )
        return jsonify({"success": True, "success_count": success_count, "failed": failed})

    elif action == "role":
        role = data.get("role")
        if role not in ["admin", "authenticated"]:
            return jsonify({"error": "无效的角色类型"}), 400
        for uid in user_ids:
            try:
                res = supabase_client.update_user_role(uid, role)
                if res["success"]:
                    success_count += 1
                else:
                    failed.append({"id": uid, "error": res.get("error", "切换失败")})
            except Exception as e:
                failed.append({"id": uid, "error": str(e)})
        database.log_admin_action(
            admin_identity(), "batch_role", details=f"批量切换角色为 {role}，成功 {success_count} 个",
            ip_address=get_client_ip()
        )
        return jsonify({"success": True, "success_count": success_count, "failed": failed})

    return jsonify({"error": "不支持的操作类型"}), 400


@app.route("/admin/api/users/<user_id>", methods=["PUT"])
@admin_required
def admin_edit_user(user_id):
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    updates = {}
    if "email" in data:
        updates["email"] = data["email"]
    if "username" in data:
        updates["username"] = data["username"]
    if "role" in data:
        updates["role"] = data["role"]

    if not updates:
        return jsonify({"error": "没有提供需要更新的字段"}), 400

    # 防止管理员将自己的角色改为普通用户
    admin_name = admin_identity()
    if updates.get("role") and updates.get("role") != "admin":
        user_result = supabase_client.get_user_by_id(user_id)
        if user_result.get("success"):
            target_identity = user_result["user"].get("username") or user_result["user"].get("email") or ''
        else:
            local_user = database.get_user_detail_with_sessions(user_id)
            target_identity = local_user.get("username") or local_user.get("email") or '' if local_user else ''
        if target_identity and target_identity == admin_name:
            return jsonify({"error": "不能将自己的管理员权限取消"}), 400

    if "email" in updates:
        result = supabase_client.update_user_email(user_id, updates["email"])
        if not result["success"]:
            return jsonify({"success": False, "error": result["error"]}), 500

    if "role" in updates:
        result = supabase_client.update_user_role(user_id, updates["role"])
        if not result["success"]:
            return jsonify({"success": False, "error": result["error"]}), 500

    if "username" in updates:
        result = supabase_client.admin_update_user_metadata(user_id, {"username": updates["username"]})
        if not result["success"]:
            return jsonify({"success": False, "error": result["error"]}), 500
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        try:
            cursor.execute('UPDATE users SET username = ? WHERE user_id = ?', (updates["username"], user_id))
            conn.commit()
        finally:
            conn.close()

    database.log_admin_action(
        admin_identity(), "edit_user", target_user_id=user_id,
        details=f"更新字段 {list(updates.keys())}", ip_address=get_client_ip()
    )
    return jsonify({"success": True})


@app.route("/admin/api/users/<user_id>/password", methods=["POST"])
@admin_required
def admin_reset_password(user_id):
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    new_password = data.get("password")
    if not new_password:
        return jsonify({"error": "新密码不能为空"}), 400

    if len(new_password) < 6:
        return jsonify({"error": "密码长度至少6位"}), 400

    result = supabase_client.admin_update_user_password(user_id, new_password)
    if result["success"]:
        database.log_admin_action(
            admin_identity(), "reset_password", target_user_id=user_id,
            details="重置用户密码", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "message": "密码重置成功"})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500


@app.route("/admin/api/users/<user_id>/role", methods=["POST"])
@admin_required
def admin_change_role(user_id):
    data = safe_get_json()
    if not data:
        return jsonify({"error": "无效的请求格式"}), 400

    role = data.get("role")
    if role not in ["admin", "authenticated"]:
        return jsonify({"error": "无效的角色类型"}), 400

    # 防止管理员取消自己的权限
    if role != "admin":
        admin_name = admin_identity()
        user_result = supabase_client.get_user_by_id(user_id)
        if user_result.get("success"):
            target_identity = user_result["user"].get("username") or user_result["user"].get("email") or ''
        else:
            local_user = database.get_user_detail_with_sessions(user_id)
            target_identity = local_user.get("username") or local_user.get("email") or '' if local_user else ''
        if target_identity and target_identity == admin_name:
            return jsonify({"error": "不能将自己的管理员权限取消"}), 400

    result = supabase_client.update_user_role(user_id, role)
    if result["success"]:
        database.log_admin_action(
            admin_identity(), "change_role", target_user_id=user_id,
            details=f"切换角色为 {role}", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "role": role})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500


def generate_random_password(length=10):
    """生成一个易读且足够安全的随机密码"""
    import secrets
    import string
    chars = string.ascii_letters + string.digits
    while True:
        pwd = ''.join(secrets.choice(chars) for _ in range(length))
        # 确保包含大小写和数字
        if (any(c.islower() for c in pwd)
                and any(c.isupper() for c in pwd)
                and any(c.isdigit() for c in pwd)):
            return pwd


@app.route("/admin/api/users/<user_id>/generate_password", methods=["POST"])
@admin_required
def admin_generate_password(user_id):
    """生成一个随机密码并设置给用户，返回明文（仅一次）"""
    new_password = generate_random_password(10)

    result = supabase_client.admin_update_user_password(user_id, new_password)
    if result["success"]:
        database.log_admin_action(
            admin_identity(), "generate_password", target_user_id=user_id,
            details="生成随机密码", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "password": new_password})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500


@app.route("/admin/api/logs", methods=["GET"])
@admin_required
def admin_get_logs():
    limit = request.args.get('limit', 50, type=int)
    offset = request.args.get('offset', 0, type=int)
    action_type = request.args.get('action_type', '')
    target_user = request.args.get('target_user', '')
    start_time = request.args.get('start_time', '')
    end_time = request.args.get('end_time', '')
    if limit > 200:
        limit = 200
    logs = database.get_admin_logs(
        limit=limit, offset=offset, action_type=action_type,
        target_user=target_user, start_time=start_time, end_time=end_time
    )
    return jsonify({"success": True, "logs": logs})


@app.route("/admin/api/admin/change_password", methods=["POST"])
@admin_required
def admin_change_password():
    data = safe_get_json()
    if not data:
        return jsonify({"success": False, "error": "无效的请求格式"}), 400

    old_password = data.get("old_password")
    new_password = data.get("new_password")
    confirm_password = data.get("confirm_password")

    if not old_password or not new_password or not confirm_password:
        return jsonify({"success": False, "error": "请填写所有字段"}), 400

    if new_password != confirm_password:
        return jsonify({"success": False, "error": "两次输入的新密码不一致"}), 400

    if len(new_password) < 6:
        return jsonify({"success": False, "error": "新密码长度至少为6位"}), 400

    # 通过当前会话拿到 admin_id，再验证原密码（避免信任客户端 cookie）
    session = _current_admin_session()
    if not session:
        return jsonify({"success": False, "error": "会话已失效，请重新登录"}), 401
    admin = database.get_admin_by_id(session['admin_id'])
    if not admin:
        return jsonify({"success": False, "error": "账号不存在"}), 404

    verify = database.verify_admin_password(admin['username'], old_password)
    if not verify.get("success"):
        return jsonify({"success": False, "error": "原密码错误"}), 401

    update_result = database.update_admin_password(admin['id'], new_password)
    if not update_result.get("success"):
        return jsonify({"success": False, "error": update_result.get("error", "更新失败")}), 500

    # update_admin_password 已吊销所有会话，前端需重新登录
    database.log_admin_action(
        admin['username'], "change_admin_password", details="修改管理员密码",
        ip_address=get_client_ip()
    )
    response = jsonify({"success": True, "message": "密码修改成功，请重新登录"})
    _clear_auth_cookies(response)
    return response


# ===== 管理员账号管理（仅超管可操作）=====

@app.route("/admin/api/admin/users", methods=["GET"])
@super_admin_required
def admin_list_admin_users():
    """列出所有管理员账号"""
    admins = database.list_admin_users()
    # 附加在线会话数（最近 5 分钟内活跃）
    return jsonify({"success": True, "admins": admins})


@app.route("/admin/api/admin/users", methods=["POST"])
@super_admin_required
def admin_create_admin_user():
    """创建管理员账号（超管或普通管理员）"""
    data = safe_get_json()
    if not data:
        return jsonify({"success": False, "error": "无效的请求格式"}), 400

    username = (data.get("username") or "").strip()
    password = data.get("password") or ""
    display_name = (data.get("display_name") or "").strip() or None
    is_super = bool(data.get("is_super", False))

    if not username or not password:
        return jsonify({"success": False, "error": "用户名和密码不能为空"}), 400
    if len(password) < 6:
        return jsonify({"success": False, "error": "密码长度至少6位"}), 400

    result = database.create_admin_user(username, password, display_name=display_name, is_super=is_super)
    if result.get("success"):
        database.log_admin_action(
            admin_identity(), "create_admin_user", target_user_id=result["admin"]["id"],
            details=f"创建管理员 {username} (super={is_super})", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "admin": result["admin"]})
    return jsonify({"success": False, "error": result.get("error", "创建失败")}), 409


@app.route("/admin/api/admin/users/<admin_id>", methods=["DELETE"])
@super_admin_required
def admin_delete_admin_user(admin_id):
    """删除管理员账号"""
    session = _current_admin_session()
    if session and session['admin_id'] == admin_id:
        return jsonify({"success": False, "error": "不能删除当前登录的管理员账号"}), 400

    target = database.get_admin_by_id(admin_id)
    if not target:
        return jsonify({"success": False, "error": "管理员不存在"}), 404

    result = database.delete_admin_user(admin_id)
    if result.get("success"):
        database.log_admin_action(
            admin_identity(), "delete_admin_user", target_user_id=admin_id,
            details=f"删除管理员 {target['username']}", ip_address=get_client_ip()
        )
        return jsonify({"success": True})
    return jsonify({"success": False, "error": result.get("error", "删除失败")}), 400


@app.route("/admin/api/admin/users/<admin_id>", methods=["PUT"])
@super_admin_required
def admin_update_admin_user(admin_id):
    """修改管理员属性：display_name / is_super / is_active"""
    data = safe_get_json()
    if not data:
        return jsonify({"success": False, "error": "无效的请求格式"}), 400

    target = database.get_admin_by_id(admin_id)
    if not target:
        return jsonify({"success": False, "error": "管理员不存在"}), 404

    changes = []

    # 修改显示名
    if 'display_name' in data:
        new_name = (data['display_name'] or '').strip()
        if not new_name or len(new_name) < 2 or len(new_name) > 20:
            return jsonify({"success": False, "error": "显示名长度需在2-20位之间"}), 400
        if not re.match(r'^[a-zA-Z0-9_\u4e00-\u9fa5]+$', new_name):
            return jsonify({"success": False, "error": "显示名只允许字母、数字、下划线和中文"}), 400
        if database.update_admin_display_name(admin_id, new_name):
            changes.append(f"display_name={new_name}")

    # 修改超管身份
    if 'is_super' in data:
        new_super = bool(data['is_super'])
        # 防止自我降级导致无超管
        session = _current_admin_session()
        if session and session['admin_id'] == admin_id and not new_super:
            return jsonify({"success": False, "error": "不能取消自己的超级管理员身份"}), 400
        # 直接更新 is_super 字段
        import sqlite3 as _sqlite3_for_admin
        conn = _sqlite3_for_admin.connect(database.DB_PATH)
        try:
            cur = conn.cursor()
            cur.execute('UPDATE admin_users SET is_super = ?, updated_at = CURRENT_TIMESTAMP WHERE id = ?',
                        (1 if new_super else 0, admin_id))
            conn.commit()
            changes.append(f"is_super={new_super}")
        finally:
            conn.close()

    # 启用/禁用
    if 'is_active' in data:
        new_active = bool(data['is_active'])
        session = _current_admin_session()
        if session and session['admin_id'] == admin_id and not new_active:
            return jsonify({"success": False, "error": "不能禁用当前登录的管理员账号"}), 400
        result = database.set_admin_active(admin_id, new_active)
        if not result.get("success"):
            return jsonify({"success": False, "error": result.get("error", "操作失败")}), 500
        changes.append(f"is_active={new_active}")

    if changes:
        database.log_admin_action(
            admin_identity(), "update_admin_user", target_user_id=admin_id,
            details=f"修改管理员 {target['username']}：{', '.join(changes)}", ip_address=get_client_ip()
        )

    return jsonify({"success": True, "admin": database.get_admin_by_id(admin_id)})


@app.route("/admin/api/admin/users/<admin_id>/reset_password", methods=["POST"])
@super_admin_required
def admin_reset_admin_password(admin_id):
    """超管重置其他管理员密码（不需要原密码）"""
    data = safe_get_json() or {}
    new_password = data.get("new_password") or ""

    if not new_password:
        # 未提供则生成随机密码
        new_password = generate_random_password(12)

    if len(new_password) < 6:
        return jsonify({"success": False, "error": "密码长度至少6位"}), 400

    target = database.get_admin_by_id(admin_id)
    if not target:
        return jsonify({"success": False, "error": "管理员不存在"}), 404

    result = database.update_admin_password(admin_id, new_password)
    if result.get("success"):
        database.log_admin_action(
            admin_identity(), "reset_admin_password", target_user_id=admin_id,
            details=f"重置管理员密码：{target['username']}", ip_address=get_client_ip()
        )
        # 返回明文密码（仅一次），便于超管转交
        return jsonify({"success": True, "password": new_password})
    return jsonify({"success": False, "error": result.get("error", "重置失败")}), 500


@app.route("/admin/api/system_info", methods=["GET"])
@admin_required
def admin_get_system_info():
    import platform
    import multiprocessing
    import ctypes

    cpu_count = multiprocessing.cpu_count()
    cpu_percent = None
    memory_total = None
    memory_used = None
    memory_percent = None
    disk_total = None
    disk_used = None
    disk_percent = None
    uptime = None

    try:
        is_windows = platform.system() == 'Windows'
        if is_windows:
            try:
                class MEMORYSTATUSEX(ctypes.Structure):
                    _fields_ = [
                        ('dwLength', ctypes.c_ulong),
                        ('dwMemoryLoad', ctypes.c_ulong),
                        ('ullTotalPhys', ctypes.c_ulonglong),
                        ('ullAvailPhys', ctypes.c_ulonglong),
                        ('ullTotalPageFile', ctypes.c_ulonglong),
                        ('ullAvailPageFile', ctypes.c_ulonglong),
                        ('ullTotalVirtual', ctypes.c_ulonglong),
                        ('ullAvailVirtual', ctypes.c_ulonglong),
                        ('ullAvailExtendedVirtual', ctypes.c_ulonglong),
                    ]
                mem = MEMORYSTATUSEX()
                mem.dwLength = ctypes.sizeof(mem)
                ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(mem))
                memory_total = round(mem.ullTotalPhys / (1024**3), 2)
                memory_used = round((mem.ullTotalPhys - mem.ullAvailPhys) / (1024**3), 2)
                memory_percent = round(float(mem.dwMemoryLoad), 1)
            except Exception:
                pass

            try:
                free_bytes = ctypes.c_ulonglong(0)
                total_bytes = ctypes.c_ulonglong(0)
                ctypes.windll.kernel32.GetDiskFreeSpaceExW(
                    ctypes.c_wchar_p('.'),
                    ctypes.byref(free_bytes),
                    ctypes.byref(total_bytes),
                    None
                )
                disk_total = round(total_bytes.value / (1024**3), 2)
                disk_used = round((total_bytes.value - free_bytes.value) / (1024**3), 2)
                disk_percent = round((total_bytes.value - free_bytes.value) / total_bytes.value * 100, 1)
            except Exception:
                pass

            try:
                kernel32 = ctypes.windll.kernel32
                tick_count = kernel32.GetTickCount64()
                uptime = int(tick_count / 1000)
            except Exception:
                try:
                    import ctypes.wintypes
                    class FILETIME(ctypes.Structure):
                        _fields_ = [('dwLowDateTime', ctypes.wintypes.DWORD),
                                    ('dwHighDateTime', ctypes.wintypes.DWORD)]

                    class SYSTEM_INFO(ctypes.Structure):
                        _fields_ = [
                            ('wProcessorArchitecture', ctypes.wintypes.WORD),
                            ('wReserved', ctypes.wintypes.WORD),
                            ('dwPageSize', ctypes.wintypes.DWORD),
                            ('lpMinimumApplicationAddress', ctypes.c_void_p),
                            ('lpMaximumApplicationAddress', ctypes.c_void_p),
                            ('dwActiveProcessorMask', ctypes.c_void_p),
                            ('dwNumberOfProcessors', ctypes.wintypes.DWORD),
                            ('dwProcessorType', ctypes.wintypes.DWORD),
                            ('dwAllocationGranularity', ctypes.wintypes.DWORD),
                            ('wProcessorLevel', ctypes.wintypes.WORD),
                            ('wProcessorRevision', ctypes.wintypes.WORD),
                        ]
                    uptime = int(ctypes.windll.kernel32.GetTickCount() / 1000)
                except Exception:
                    pass
        else:
            try:
                import subprocess
                result = subprocess.run(['df', '-k', '.'], capture_output=True, text=True, timeout=5)
                lines = result.stdout.strip().split('\n')
                if len(lines) >= 2:
                    parts = lines[-1].split()
                    if len(parts) >= 4:
                        disk_total = round(int(parts[1]) / (1024**2), 2)
                        disk_used = round(int(parts[2]) / (1024**2), 2)
                        disk_percent = float(parts[4].replace('%', ''))
            except Exception:
                pass
            try:
                with open('/proc/uptime', 'r') as f:
                    uptime = int(float(f.read().split()[0]))
            except Exception:
                pass
            try:
                with open('/proc/meminfo', 'r') as f:
                    meminfo = {}
                    for line in f:
                        parts = line.split()
                        if len(parts) >= 2:
                            meminfo[parts[0].rstrip(':')] = int(parts[1])
                    if 'MemTotal' in meminfo:
                        memory_total = round(meminfo['MemTotal'] / (1024**2), 2)
                    if 'MemAvailable' in meminfo and 'MemTotal' in meminfo:
                        memory_used = round((meminfo['MemTotal'] - meminfo['MemAvailable']) / (1024**2), 2)
                        memory_percent = round((meminfo['MemTotal'] - meminfo['MemAvailable']) / meminfo['MemTotal'] * 100, 1)
            except Exception:
                pass
    except Exception:
        pass

    return jsonify({
        "success": True,
        "system": {
            "platform": platform.system(),
            "platform_version": platform.version(),
            "python_version": platform.python_version(),
            "hostname": platform.node(),
            "cpu_count": cpu_count,
            "cpu_percent": cpu_percent,
            "memory_total": memory_total,
            "memory_used": memory_used,
            "memory_percent": memory_percent,
            "disk_total": disk_total,
            "disk_used": disk_used,
            "disk_percent": disk_percent,
            "uptime": uptime,
            "current_time": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime()),
            "app_version": "4.0",
            "admin_users": [a['username'] for a in database.list_admin_users()]
        }
    })


@app.route("/admin/api/clean_cache", methods=["POST"])
@admin_required
def admin_clean_cache():
    global TEMPLATE_CACHE
    cache_size = len(TEMPLATE_CACHE)
    TEMPLATE_CACHE.clear()
    database.log_admin_action(
        admin_identity(), "clean_cache", details=f"清理模板缓存，共 {cache_size} 条",
        ip_address=get_client_ip()
    )
    return jsonify({"success": True, "message": f"缓存已清理，共清除 {cache_size} 条记录"})


@app.route("/admin/api/users/<user_id>/disable", methods=["POST"])
@admin_required
def admin_disable_user(user_id):
    admin_name = admin_identity()
    user_result = supabase_client.get_user_by_id(user_id)
    if user_result.get("success"):
        target_identity = user_result["user"].get("username") or user_result["user"].get("email") or ''
        if target_identity and target_identity == admin_name:
            return jsonify({"success": False, "error": "不能禁用当前登录的管理员账号"}), 400
    else:
        # 修复 M8：Supabase 查询失败时不允许跳过自禁检查继续执行禁用操作，避免管理员误禁自己
        return jsonify({"success": False, "error": "用户信息查询失败，禁止执行禁用操作: " + user_result.get("error", "Supabase 不可用")}), 503

    result = supabase_client.admin_update_user_by_id(user_id, {"disabled": True})
    if result["success"]:
        database.log_admin_action(
            admin_name, "disable_user", target_user_id=user_id,
            details="禁用用户账号", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "message": "用户已禁用"})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500


@app.route("/admin/api/users/<user_id>/enable", methods=["POST"])
@admin_required
def admin_enable_user(user_id):
    result = supabase_client.admin_update_user_by_id(user_id, {"disabled": False})
    if result["success"]:
        database.log_admin_action(
            admin_identity(), "enable_user", target_user_id=user_id,
            details="启用用户账号", ip_address=get_client_ip()
        )
        return jsonify({"success": True, "message": "用户已启用"})
    else:
        return jsonify({"success": False, "error": result["error"]}), 500

if __name__ == "__main__":
    print("""
============================================
       BeiKe Assistant (Web) - 登录版 5.3
       http://localhost:5000
       Ctrl+C to exit
============================================
""")
    import webbrowser
    threading.Timer(1.5, lambda: webbrowser.open("http://127.0.0.1:5000/")).start()
    app.run(host="0.0.0.0", port=5000, debug=False, threaded=True)
