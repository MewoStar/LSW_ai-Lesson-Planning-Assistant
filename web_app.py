# ============================================
# 备课助手 Web 版 - 用户登录版
# 启动: python web_app.py
# 访问: http://localhost:5024
# ============================================

import re
import os
import sys
import yaml
import random
import sqlite3
import threading
import queue
from pathlib import Path
from flask import Flask, render_template, request, jsonify, send_file, send_from_directory, Response, redirect, url_for, make_response
from openai import OpenAI
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import time
import json

# 导入数据库模块
import database

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys._MEIPASS)
    return Path(__file__).parent

def get_data_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys.executable).parent
    return Path(__file__).parent

BASE_DIR = get_base_dir()
DATA_DIR = get_data_dir()

# 读配置
CONFIG_PATH = BASE_DIR / "config.yaml"
cfg = {}
if CONFIG_PATH.exists():
    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

SECRET_KEY = cfg.get("secret_key") or os.environ.get("SECRET_KEY")
if not SECRET_KEY:
    raise ValueError("必须配置 secret_key（在 config.yaml 中设置或设置环境变量 SECRET_KEY）")

template_dir = str(BASE_DIR / "templates")
app = Flask(__name__, template_folder=template_dir)
app.secret_key = SECRET_KEY
app.config['TEMPLATES_AUTO_RELOAD'] = True

API_KEY = os.environ.get("DEEPSEEK_API_KEY") or os.environ.get("OPENAI_API_KEY") or cfg.get("api_key") or ""
BASE_URL = cfg.get("base_url") or os.environ.get("DEEPSEEK_BASE_URL") or os.environ.get("OPENAI_BASE_URL") or "https://api.deepseek.com/v1"
MODEL = cfg.get("model", "deepseek-chat")
TEMP = cfg.get("temperature", 0.7)
MAX_TOKENS = cfg.get("max_tokens", 3072)
MAX_HISTORY = cfg.get("max_history", 6)
OUTPUT_DIR = cfg.get("output_dir", ".")
if not os.path.isabs(OUTPUT_DIR):
    OUTPUT_DIR = os.path.abspath(os.path.join(str(DATA_DIR), OUTPUT_DIR))
else:
    OUTPUT_DIR = os.path.abspath(OUTPUT_DIR)

os.makedirs(OUTPUT_DIR, exist_ok=True)

HISTORY_DIR = os.path.join(OUTPUT_DIR, "历史记录")
os.makedirs(HISTORY_DIR, exist_ok=True)

VISUAL_OUTPUT_DIR = DATA_DIR / "output"
os.makedirs(str(VISUAL_OUTPUT_DIR), exist_ok=True)

QWEN_IMAGE_API_KEY = cfg.get("qwen_image_api_key") or os.environ.get("QWEN_IMAGE_API_KEY") or ""
QWEN_IMAGE_MODEL = cfg.get("qwen_image_model", "qwen-image-plus")

IMAGE_SEARCH_PROVIDER = cfg.get("image_search_provider", "unsplash")
IMAGE_SEARCH_API_KEY = cfg.get("image_search_api_key") or os.environ.get("IMAGE_SEARCH_API_KEY") or ""

ALLOWED_UPLOAD_EXT = {'.md', '.txt', '.docx', '.xlsx', '.csv', '.pdf'}
MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # 20MB
MAX_TEXT_CHARS = 30000  # 单文件最多送3万字符给AI


def safe_get_json():
    try:
        data = request.get_json(silent=True)
        if data is None:
            data = request.form if request.form else None
        return data
    except Exception:
        return None


def get_user_output_dir(user_id=None):
    if user_id:
        user_dir = os.path.join(HISTORY_DIR, f"user_{user_id}")
    else:
        user_dir = os.path.join(HISTORY_DIR, "public")
    os.makedirs(user_dir, exist_ok=True)
    return user_dir

SYSTEM_PROMPT = """你是专业的「AI 备课助手」，服务对象是中小学 / 职业院校的一线教师。你的任务是根据教师的需求，按规范格式生成高质量的教学资源，并**严格用指定标签包裹可保存的文件内容**，方便系统自动解析并下载。

---

## 🎯 核心能力（6 大模块）
1. **PPT 课件生成**：输出结构化 PPT 大纲，系统会自动解析为 `.pptx` 课件
2. **PPT 讲解纲要**：为每页幻灯片生成详细的教师讲解稿（逐页讲稿 + 互动设计）
3. **教案生成**：45 分钟标准格式教案（三维目标 / 教学过程 / 板书 / 反思）
4. **习题生成**：基础 / 提高 / 拓展 三级难度分层习题，附答案解析
5. **作业分析**：作业批改数据分析 + 错因归因 + 辅导建议 + 教学调整策略
6. **复习提纲**：期末 / 单元复习资料，知识脉络 + 核心考点 + 典型例题

---

## 🚫 最重要的输出规则（不遵守会导致文件无法下载）
**必须用 `[文件:xxx.md]  …内容…  [/文件]` 或 `[PPT:xxx.pptx]  …大纲…  [/PPT]` 标签把生成结果包裹起来**，系统会自动保存为可下载的文件。
- 如果用户要的是「教案 / 习题 / 作业分析 / 讲解纲要 / 复习提纲」 → 用 `[文件:xxx.md]`
- 如果用户要的是「PPT 课件 / 幻灯片」 → 用 `[PPT:xxx.pptx]`
- 若用户同时要多份内容（如「同时生成教案和PPT」），可同时出现多个标签块
- 纯文字的简短说明 / 寒暄可以写在标签外，但**核心生成内容必须全部放在标签块内**

✅ 正确示例：
```
好的，为您生成《背影》的教案如下：
[文件:初中语文《背影》教案.md]
# 初中语文《背影》教案

## 一、三维教学目标
- **知识与技能**：……
- **过程与方法**：……
- **情感态度与价值观**：……

## 二、教学重难点
……（以下内容全部写在标签内）……

## 六、教学反思栏
……
[/文件]
```

❌ 错误示例（严禁！会导致用户无法下载文件）：
- 直接输出 Markdown，不写任何标签包裹
- 标签名写错，如 `[教案:xxx]`、`[下载:xxx]`
- 文件内容一半在标签里，一半在标签外

---

## 📖 模块 1：教案生成 标准格式（`[文件:xxx教案.md]`）
结构必须完整，标题层级严格按下面顺序：

```
# 《课题名》教案（学段 + 学科）

## 一、基本信息
- 学科：xxx
- 学段：xxx（小学/初中/高中/中职）
- 年级：xxx
- 课时：1课时（45分钟）
- 课型：新授课 / 复习课 / 讲评课 / 实验课 / 活动课

## 二、三维教学目标 / 核心素养目标
- 知识与技能：……
- 过程与方法：……
- 情感态度与价值观：……

## 三、教学重难点
- **教学重点**：……
- **教学难点**：……

## 四、学情分析与教学方法
- 学情分析：……
- 教学方法：讲授法 / 讨论法 / 实验法 / 情境教学法 / 任务驱动法

## 五、教学准备
- 教具准备：多媒体课件、实验器材、学案……
- 学生准备：……

## 六、详细教学过程（精确到分钟）
### 1. 导入新课（约 5 分钟）
……教师活动……
……学生活动……

### 2. 新知讲授（约 15 分钟）
……分步骤讲解……

### 3. 课堂练习 / 小组讨论（约 12 分钟）
……

### 4. 课堂小结（约 5 分钟）
……由学生总结 + 教师补充……

### 5. 作业布置（约 3 分钟）
- 基础作业：……
- 拓展作业：……

### 6. 板书设计（约 5 分钟同步板书）
主板书 | 副板书
-------|-------
……    | ……

## 七、板书设计
- 主板书：……（结构化呈现）
- 副板书：……

## 八、分层作业设计
- ✅ **基础层（必做）**：……
- 📈 **提高层（选做）**：……
- 🚀 **拓展层（挑战）**：……

## 九、教学反思
- 本节课亮点：……
- 待改进点：……
- 课堂生成问题记录（留白教师填写）：
```

---

## 📝 模块 2：习题生成 标准格式（`[文件:xxx习题.md]`）
```
# 《课题 / 知识点》分层练习题

## 一、基本信息
- 学段学科：xxx
- 考查知识点：xxx
- 题目总数：xx 道
- 建议用时：xx 分钟
- 难度分布：基础 60% / 提高 30% / 拓展 10%

---

## 二、基础巩固题（约 60%）
### 一、选择题（每题 3 分，共 xx 分）
1. 题目内容……
   A. xxx   B. xxx   C. xxx   D. xxx
2. ……

### 二、填空题（每题 2 分，共 xx 分）
1. _________________
2. _________________

---

## 三、能力提高题（约 30%）
### 三、解答题 / 计算题（每题 xx 分，共 xx 分）
1. 题目内容……

---

## 四、拓展探究题（约 10%）
### 四、综合应用题 / 探究题
1. 题目内容……

---

## 五、参考答案与详细解析
1. 选择题 1：答案 B
   - 解析：……
   - 考点：……
   - 易错分析：……
2. ……
```

---

## 📊 模块 3：作业分析报告（`[文件:xxx作业分析.md]`）
```
# 《作业名称》学情分析报告

## 一、整体概况
- 学科：xxx
- 班级：xxx（共 xx 人）
- 应交：xx 份，实交：xx 份，上交率：xx%
- 平均分：xx 分 ｜ 最高分：xx ｜ 最低分：xx ｜ 及格率：xx% ｜ 优秀率：xx%

### 分数段分布
| 分数段    | 人数 | 占比  |
|-----------|------|-------|
| 90-100 分 | x 人 | xx%  |
| 80-89 分  | x 人 | xx%  |
| 70-79 分  | x 人 | xx%  |
| 60-69 分  | x 人 | xx%  |
| 60 分以下 | x 人 | xx%  |

## 二、知识点掌握情况雷达
| 知识点         | 平均得分率 | 掌握等级 |
|----------------|------------|----------|
| 知识点 A       | xx%        | 熟练 / 一般 / 薄弱 |
| 知识点 B       | xx%        | ……     |

## 三、高频错题 TOP 5 + 归因分析
| 排名 | 题号 | 知识点 | 得分率 | 主要错解 | 错因归类 |
|------|------|--------|--------|----------|----------|
| 1    | 第x题 | xxx | xx% | …… | ⚠️ 概念不清 / 审题错误 / 计算失误 / 方法未掌握 |
| 2    | …… | …… | …… | …… | …… |

### 典型错解展示（第 x 题）
> 错误解法示例：……
> ✅ 正确解法：……
> 🔍 错因剖析：……

## 四、分层辅导建议
### 🔴 学困生（xx 人，60 分以下）
- 薄弱点：……
- 辅导策略：① …… ② …… ③ ……
- 补充练习题：……

### 🟡 中等生（xx 人，60-84 分）
- 提升点：……
- 辅导策略：……
- 强化练习：……

### 🟢 优等生（xx 人，85 分以上）
- 拓展方向：……
- 挑战性题目：……

## 五、下节课教学调整建议
- 需要补讲的知识点：……（建议用时 xx 分钟）
- 课堂讲评顺序建议：先讲第 x/x/x 题（得分率最低）
- 教学方法调整：……
- 与家长沟通要点（家校共育）：……

## 六、讲评课时分配建议
| 环节            | 时间  | 内容说明 |
|-----------------|-------|----------|
| 整体情况通报    | 3 分钟 | 成绩分布 + 表扬优秀进步学生 |
| 高频错题精讲    | 20 分钟 | TOP 5 错因分析 + 变式训练 |
| 小组互助订正    | 10 分钟 | 学困生结对，同伴讲解 |
| 针对性补充练习  | 8 分钟  | 相似题型，当堂巩固 |
| 总结 + 二次过关 | 4 分钟  | 关键方法总结，小测过关 |
```

---

## 🎤 模块 4：PPT 讲解纲要（逐页讲稿，`[文件:xxx讲解纲要.md]`）
```
# 《课件名》PPT 逐页讲解纲要

## 使用说明
- 总页数：xx 页
- 建议总时长：45 分钟
- 适用对象：xxx 年级学生

---

## 第 1 页：封面（约 1 分钟）
🎤 **教师讲解词**：
同学们好，今天我们一起来学习《……》。在上课之前请大家看屏幕上的这张图片 / 这个问题，有没有同学能说一说……
❓ **互动提问**（可选）：……
⭐ **语气提示**：语速稍慢，吸引注意力

---

## 第 2 页：教学目标（约 2 分钟）
🎤 **教师讲解词**：
本节课我们要达成三个学习目标：第一，……；第二，……；第三，……
⭐ **强调**：第 2 条是重点，请用红色笔在学案上画出来
⏰ 本页用时：2 分钟

---

## 第 3 页：xxx（按 PPT 每页依次写）
🎤 **教师讲解词**：……
❓ **互动提问**：……（预留学生作答的留白，可写学生可能的回答）
⚠️ **易错提醒**：……
⏰ 本页用时：xx 分钟

---
（每页 PPT 都按上面格式写一节）
---

## 最后一页：课堂小结 + 作业布置（约 3 分钟）
🎤 **教师讲解词**：
好，本节课我们学习了……，主要内容可以用三句话记住：①……②……③……
下课后请大家完成：……（作业内容）
下课，同学们再见！
```

---

## 🤖 模块 5：复习提纲（`[文件:xxx复习提纲.md]`）
```
# 《课程名》期末 / 单元复习提纲

## 📚 一、复习内容概览
- 章节范围：第 x 章 — 第 x 章
- 建议复习用时：xx 课时
- 重要程度：★★★★★（5 星为必考）

## 🧠 二、核心知识体系（思维导图式）
### 模块一：xxx
- 核心概念 1：……
  - 定义：……
  - 关键词：……
- 核心概念 2：……

### 模块二：xxx
……

## 📐 三、重点公式 / 定理 / 结论速记
| 名称 | 公式 / 结论 | 适用条件 | 考频 |
|------|-------------|----------|------|
| xxx  | ……         | ……       | ⭐⭐⭐ |

## 📝 四、典型例题精讲
### 例题 1（★★ 基础题）
> 题目：……
> 解题步骤：
> ① ……
> ② ……
> 💡 点拨：……

### 例题 2（★★★★ 高频考点）
……

## ⚠️ 五、易错点警示 TOP 10
1. ❌ 错误理解：…… → ✅ 正确：……
2. ❌ ……

## 🎯 六、考点预测 + 分值分布
| 考点         | 预测分值 | 题型       | 难度 |
|--------------|----------|------------|------|
| xxx          | 8-12 分  | 选择 + 解答 | ★★★ |

## ✅ 七、自我检测（附答案）
……（10 道精选小题，覆盖全部考点）
```

---

## 📊 模块 6：PPT 课件生成（必须用 `[PPT:xxx.pptx]` … `[/PPT]` 包裹）
**结构严格遵守以下格式，系统会自动解析成 `.pptx` 文件：**

```
[PPT:荷塘月色_语文课件.pptx]
# 《荷塘月色》语文课件
朱自清散文 · 高中语文必修上册

## 第1页：情境导入
- 展示荷花池夜色图片
- 提问：同学们记忆中关于"月"和"荷"的诗句有哪些？
- 引出作者：朱自清
- 板书课题

## 第2页：学习目标
- 📖 知识与技能：掌握重点词语，理解关键语句含义
- 🧠 过程与方法：通过朗读体会文章的语言美和意境美
- ❤️ 情感态度：体会作者情感，感受文学作品的感染力

## 第3页：作者介绍
- 朱自清（1898-1948），原名自华，字佩弦
- 现代著名散文家、诗人、学者
- 代表作品：《背影》《春》《匆匆》《荷塘月色》
- 散文风格：语言洗练，文笔清丽，情感真挚

## 第4页：写作背景
……

## 第N页：课堂小结 + 作业布置
……
[/PPT]
```

**PPT 格式要求（务必遵守）：**
1. 第一行 `# 主标题` 是封面主标题，第二行（无特殊符号）是副标题
2. 每页幻灯片用 `## 第N页：页面标题` 开头
3. 每页下面的内容用 `- 要点` 列（最多 7 条，每条简短）
4. 页数控制在 8-20 页之间，结构完整
5. **绝对不要**把多页内容塞进一个 `##` 里

---

## 🎨 其它输出规范
- 全文中文表达，专业简洁
- Markdown 层级清晰：# 一级 → ## 二级 → ### 三级
- 重点用 **粗体**、表格整理数据、列表梳理条目
- 编号必须连续，避免乱码字符

## 📎 上传文件智能识别规则（最重要！）
用户消息中如果出现 `【用户上传文件：文件名（类型）】` 标记，说明用户上传了该文件的全文内容，你必须按以下逻辑处理：

### 情况 A：文件类型 = "教案/教学设计模板"（文件名或内容含「教案」「教学设计」「学案」「模板」「____」「[填空]」「__」等占位符）
- 你的任务 = **按模板中所有的空白 / 占位符 / 待填写项，自动填充完整的教学内容**
- 规则：
  1. 严格**保留模板原有的章节结构、编号、表格结构、填空占位符的位置**
  2. 模板里写着「____」「___」「...」「待填写」「填写」「留白」「XXX」「[     ]」的位置，全部填上具体、详实、符合学科逻辑的内容
  3. 若用户在上传文件之外还写了需求（如"请按人教版高二语文"），需优先按该需求填充内容
  4. 没有的学科/年级信息时，按模板里能推断的学段和学科默认选一个最通用的
  5. **生成结果必须用 `[文件:xxx_已填充.md]` 完整包裹起来**，把整个模板+填充后的全部内容放进去，方便用户下载

### 情况 B：文件类型 = "学生作业/成绩/批改记录"（文件名或内容含「作业」「成绩」「得分」「分数」「错题」「批改」「考勤」「班级」「学生」「姓名」等）
- 你的任务 = **按「模块 5 作业分析」标准格式生成一份完整的学情分析报告**
- 报告内容必须包含 6 部分：
  1. **整体概况**（班级人数、均分、及格率、优秀率、最高分/最低分、分数段分布）
  2. **错题归因分析**（哪题错最多？错因归类：概念不清 / 计算失误 / 审题错误 / 知识点没掌握 / 书写规范）
  3. **知识点掌握雷达**（按知识点列出 掌握率 = 做对人数/总人数）
  4. **学生分层画像**（优生/中等/待进步各占多少 %，每一层的典型问题）
  5. **针对性辅导建议**（全班教学调整 + 分层布置作业 + 个别辅导名单）
  6. **后续教学调整策略**（下一节课怎么改、是否加小测、是否安排讲评课）
- **生成结果必须用 `[文件:xxx_作业分析报告.md]` 完整包裹**

### 情况 C：其他类型文件（如参考资料、课文原文、讲义等）
- 先明确告诉用户："已读取到您上传的《xxx》（共 x 字 / x 行）"
- 再结合用户的具体问题给出对应的回答；如果用户没有明确问题，自动判断：内容像课文原文 → 建议生成《xxx》教案/PPT/习题；内容像复习资料 → 建议生成复习提纲

最后再次强调：**所有生成内容必须放在对应的标签块中，否则用户点击"下载"时将没有任何文件！**"""

# 存会话（内存缓存 + 数据库持久化）
# 用 user_id:sesssion_id 作为 key，确保用户隔离
sessions: dict[str, list[dict]] = {}

def _session_key(user_id, session_id):
    return f"{user_id}:{session_id}"

def get_current_user():
    token = request.cookies.get('session_token')
    if token:
        return database.get_user_by_token(token)
    return None

def trim_history(messages):
    if not messages or len(messages) <= MAX_HISTORY + 1:
        return messages
    system_msg = messages[0]
    recent = messages[-(MAX_HISTORY):]
    return [system_msg] + recent

def markdown_to_word(markdown_text: str, title: str = "教案") -> BytesIO:
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)

    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title)
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(31, 78, 121)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    lines = markdown_text.split('\n')
    in_code_block = False
    code_content = []

    for line in lines:
        if line.strip().startswith('```'):
            if in_code_block:
                if code_content:
                    code_para = doc.add_paragraph()
                    code_run = code_para.add_run('\n'.join(code_content))
                    code_run.font.name = 'Courier New'
                    code_run.font.size = Pt(9)
                    code_run.font.color.rgb = RGBColor(128, 128, 128)
                    code_para.paragraph_format.left_indent = Pt(20)
                code_content = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_content.append(line)
            continue

        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('#### '):
            doc.add_heading(line[5:], level=4)
        elif '**' in line:
            parts = line.split('**')
            para = doc.add_paragraph()
            for i, part in enumerate(parts):
                if i % 2 == 1:
                    run = para.add_run(part)
                    run.bold = True
                else:
                    para.add_run(part)
        elif line.strip().startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.strip().startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif re.match(r'^\d+\.\s', line.strip()):
            doc.add_paragraph(re.sub(r'^\d+\.\s', '', line.strip()), style='List Number')
        elif '|' in line and line.strip():
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if len(cells) >= 2:
                if not all(re.match(r'^-+$', cell) for cell in cells):
                    para = doc.add_paragraph()
                    para.add_run(line)
        elif line.strip().startswith('>'):
            para = doc.add_paragraph(line[1:].strip())
            para.paragraph_format.left_indent = Pt(20)
            para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
        elif line.strip() == '---':
            doc.add_paragraph('_' * 50)
        elif line.strip():
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def _strip_md(text):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'^#{1,6}\s*', '', text)
    text = text.replace('•', '').replace('◦', '').replace('▪', '')
    text = re.sub(r'^\s*[-*]\s+', '', text)
    text = re.sub(r'^\s*\d+[.、)]\s*', '', text)
    return text.strip()


def _detect_header_level(s):
    m = re.match(r'^(#{1,6})\s+(.+)$', s)
    if m:
        return len(m.group(1)), m.group(2).strip()
    return 0, s


def _is_bullet(s):
    return bool(re.match(r'^\s*[-*•◦▪]\s+', s))


def _is_numbered(s):
    return bool(re.match(r'^\s*\d+[.、)]\s*', s))


def _indent_level(raw):
    stripped = raw.lstrip(' \t')
    n = len(raw) - len(stripped)
    if n == 0:
        return 0
    if n <= 3:
        return 1
    if n <= 6:
        return 2
    return 3


def parse_ppt_outline(markdown_text):
    lines = markdown_text.split('\n')
    title_main = ""
    title_sub = ""
    raw_sections = []

    header_levels = set()
    for raw in lines:
        s = raw.strip()
        if not s:
            continue
        lvl, _ = _detect_header_level(s)
        if lvl > 0:
            header_levels.add(lvl)

    if header_levels:
        if 1 in header_levels:
            cover_level = 1
            slide_levels = {l for l in header_levels if l > 1}
        else:
            cover_level = 0
            slide_levels = header_levels
    else:
        cover_level = 0
        slide_levels = set()

    min_slide = min(slide_levels) if slide_levels else 0

    cur_title = None
    cur_items = []

    def flush():
        nonlocal cur_title, cur_items
        if cur_title is not None:
            raw_sections.append((cur_title, cur_items))
        cur_title = None
        cur_items = []

    for raw in lines:
        if not raw.strip():
            continue
        s = raw.strip()
        lvl, content = _detect_header_level(s)

        if lvl == cover_level and cover_level > 0:
            flush()
            title_main = content
        elif lvl in slide_levels:
            flush()
            t = re.sub(r'^第\s*\d+\s*页\s*[：:]\s*', '', content)
            cur_title = t
        elif lvl > 0 and min_slide > 0:
            rel = max(0, lvl - min_slide)
            cur_items.append((min(rel, 3), content))
        elif _is_bullet(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*[-*•◦▪]\s+', '', s)
            cur_items.append((ind, content))
        elif _is_numbered(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*\d+[.、)]\s*', '', s)
            cur_items.append((ind, content))
        else:
            if title_main and not title_sub and not raw_sections and cur_title is None:
                title_sub = s
            else:
                cur_items.append((_indent_level(raw), s))
    flush()

    merged = []
    for stitle, items in raw_sections:
        if not items and merged:
            prev_t, prev_items = merged[-1]
            prev_items.append((0, stitle))
            merged[-1] = (prev_t, prev_items)
        else:
            merged.append((stitle, items))
    raw_sections = merged

    if not raw_sections and title_main:
        all_items = []
        for raw in lines:
            s = raw.strip()
            if not s or _detect_header_level(s)[0] == cover_level:
                continue
            if _is_bullet(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*[-*•◦▪]\s+', '', s)))
            elif _is_numbered(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*\d+[.、)]\s*', '', s)))
            elif _detect_header_level(s)[0] > 0:
                _, content = _detect_header_level(s)
                all_items.append((0, content))
            else:
                all_items.append((_indent_level(raw), s))
        if all_items:
            raw_sections = [("内容概览", all_items)]

    return title_main, title_sub, raw_sections


def _split_overflow(items, max_items=7):
    if len(items) <= max_items:
        return [items]
    result = []
    cur = []
    for lvl, text in items:
        if len(cur) >= max_items and lvl == 0:
            result.append(cur)
            cur = []
        cur.append((lvl, text))
    if cur:
        result.append(cur)
    return result


class PptTheme:
    PRIMARY = PptRGBColor(0x1E, 0x40, 0xAF)
    DARK = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT_SEC = PptRGBColor(0x4B, 0x55, 0x63)
    MUTED = PptRGBColor(0x9C, 0xA3, 0xAF)
    WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = PptRGBColor(0x3B, 0x82, 0xF6)
    FONT = 'Microsoft YaHei'


class PptBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.sw = self.prs.slide_width
        self.sh = self.prs.slide_height
        self.m = Inches(0.8)

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _tbox(self, slide, left, top, w, h, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def _run(self, para, text, size, color, bold=False):
        run = para.add_run()
        run.text = text
        run.font.size = PptPt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = PptTheme.FONT
        return run

    def _rect(self, slide, left, top, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _page_num(self, slide, num, total):
        tf = self._tbox(slide, self.sw - Inches(1.6), self.sh - Inches(0.45),
                        Inches(1.4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        self._run(p, f"{num} / {total}", 11, PptTheme.MUTED)

    def _footer(self, slide):
        tf = self._tbox(slide, self.m, self.sh - Inches(0.45),
                        Inches(4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        self._run(p, "AI 备课助手", 10, PptTheme.MUTED)

    def _font_size(self, items):
        n = len(items)
        chars = sum(len(_strip_md(t)) for _, t in items)
        max_lvl = max((lvl for lvl, _ in items), default=0)
        if n <= 4 and chars <= 120:
            base = 30
        elif n <= 6 and chars <= 250:
            base = 26
        elif n <= 8 and chars <= 400:
            base = 22
        elif n <= 12 and chars <= 600:
            base = 18
        else:
            base = 15
        if max_lvl >= 2 and base > 22:
            base = 22
        return base

    def cover(self, title, subtitle, speaker='', date_str=''):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(2.8), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.1), self.sw - self.m * 2,
                        Inches(2), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        t = _strip_md(title) or "课件"
        size = 46 if len(t) <= 18 else (36 if len(t) <= 36 else 28)
        self._run(p, t, size, PptTheme.WHITE, bold=True)

        if subtitle:
            tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                            Inches(0.8), MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            self._run(p, _strip_md(subtitle), 22, PptTheme.MUTED)

        footer = '  |  '.join(filter(None, [speaker, date_str]))
        if footer:
            tf = self._tbox(slide, self.m, self.sh - Inches(0.8),
                            self.sw - self.m * 2, Inches(0.4), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            self._run(p, footer, 14, PptTheme.MUTED)

    def content(self, title, items, page_num, total, dark=False):
        slide = self._blank()

        bg = PptTheme.DARK if dark else PptTheme.WHITE
        tc = PptTheme.WHITE if dark else PptTheme.TEXT
        sc = PptTheme.MUTED if dark else PptTheme.TEXT_SEC
        bc = PptTheme.ACCENT if dark else PptTheme.PRIMARY

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        tf = self._tbox(slide, self.m, Inches(0.5), self.sw - self.m * 2,
                        Inches(1), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        clean_t = _strip_md(title)
        ts = 36 if len(clean_t) <= 20 else (28 if len(clean_t) <= 40 else 22)
        self._run(p, clean_t, ts, tc, bold=True)

        self._rect(slide, self.m, Inches(1.5), Inches(2), Inches(0.06), bc)

        base = self._font_size(items)

        BULLETS = ["●", "○", "▪", "·"]
        INDENT_EMU = [0, 457200, 914400, 1371600]
        SIZE_DEC = [0, 3, 5, 6]

        valid_items = [(min(lvl, 3), _strip_md(text)) for lvl, text in items
                       if _strip_md(text)]
        n = len(valid_items)
        est_line_h = (base * 1.3 + 8) / 72.0
        est_height = Inches(min(n * est_line_h + 0.4, 4.6))

        avail_top = Inches(1.85)
        avail_bottom = self.sh - Inches(0.7)
        avail_h = avail_bottom - avail_top
        if est_height < avail_h:
            content_top = avail_top + (avail_h - est_height) / 2
            content_h = est_height
            anchor = MSO_ANCHOR.MIDDLE
        else:
            content_top = avail_top
            content_h = avail_h
            anchor = MSO_ANCHOR.TOP

        tf = self._tbox(slide, self.m, content_top, self.sw - self.m * 2,
                        content_h, anchor)

        first = True
        for lvl, clean in valid_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PptPt(6)
            p.line_spacing = 1.3

            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(INDENT_EMU[lvl]))
            pPr.set('indent', str(-228600))

            bchar = BULLETS[lvl]
            bsize = base - SIZE_DEC[lvl]
            tsize = base - SIZE_DEC[lvl]
            tcolor = tc if lvl == 0 else sc
            self._run(p, f"{bchar}  ", bsize, bc, bold=(lvl == 0))
            self._run(p, clean, tsize, tcolor, bold=False)

        self._page_num(slide, page_num, total)
        if not dark:
            self._footer(slide)

    def ending(self):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(3.0), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.3), self.sw - self.m * 2,
                        Inches(1.5), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "谢谢观看", 48, PptTheme.WHITE, bold=True)

        tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                        Inches(0.6), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "AI 备课助手 · 自动生成", 16, PptTheme.MUTED)

    def build(self, title, subtitle, content_slides, speaker='', date_str=''):
        self.cover(title, subtitle, speaker, date_str)

        expanded = []
        for stitle, items in content_slides:
            chunks = _split_overflow(items, max_items=7)
            for ci, chunk in enumerate(chunks):
                ct = stitle if ci == 0 else f"{stitle}（续{ci}）"
                expanded.append((ct, chunk))

        total = len(expanded) + 1
        for i, (stitle, items) in enumerate(expanded):
            is_summary = any(k in stitle for k in
                             ['总结', '小结', '结语', '谢谢', '感谢', 'CTA', '行动'])
            self.content(stitle, items, i + 2, total, dark=is_summary)

        if not expanded:
            self.content("内容概览", [(0, "暂无具体内容")], 2, 2)

        self.ending()

        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf


def _parse_kv(items):
    kv = {}
    for lvl, text in items:
        clean = _strip_md(text).strip()
        if not clean:
            continue
        for sep in ['：', ':', '-', '—']:
            if sep in clean:
                parts = clean.split(sep, 1)
                if len(parts) == 2:
                    key = parts[0].strip()
                    val = parts[1].strip()
                    if key and val:
                        kv[key] = val
                    break
    return kv


def markdown_to_ppt(markdown_text, title="课件"):
    title_main, title_sub, slides = parse_ppt_outline(markdown_text)
    if not title_main:
        title_main = title
    if not slides:
        slides = [(title_main or title,
                   [(0, t) for t in markdown_text.split('\n') if t.strip()])]

    cover_data = {}
    content_slides = []
    for stitle, items in slides:
        if stitle == '封面' or stitle.startswith('封面'):
            cover_data = _parse_kv(items)
        else:
            content_slides.append((stitle, items))

    if cover_data:
        c_title = _strip_md(cover_data.get('主标题', title_main))
        c_sub = _strip_md(cover_data.get('副标题', title_sub))
        speaker = _strip_md(cover_data.get('演讲人', ''))
        date_str = _strip_md(cover_data.get('日期', time.strftime('%Y年%m月%d日')))
    else:
        c_title = _strip_md(title_main) if title_main else title
        c_sub = _strip_md(title_sub) if title_sub else ''
        speaker = ''
        date_str = time.strftime('%Y年%m月%d日')

    builder = PptBuilder()
    return builder.build(c_title, c_sub, content_slides, speaker, date_str)


def generate_ppt_files(reply_text, user_id=None):
    base_dir = get_user_output_dir(user_id)
    ppt_blocks = re.findall(r"\[PPT:\s*([^\]]+)\](.*?)\[/PPT\]", reply_text, re.DOTALL)
    saved_ppts = []
    for filename, outline in ppt_blocks:
        filename = filename.strip()
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"
        safe_name = re.sub(r'[<>:"/\\|?*]', '', filename)
        try:
            title = os.path.splitext(safe_name)[0]
            ppt_buffer = markdown_to_ppt(outline, title=title)
            filepath = os.path.join(base_dir, safe_name)
            with open(filepath, "wb") as f:
                f.write(ppt_buffer.getvalue())
            saved_ppts.append(safe_name)
        except Exception as e:
            print(f"[PPT] 生成失败 {safe_name}: {e}")
    display = reply_text
    if ppt_blocks:
        display = re.sub(r"\[PPT:\s*[^\]]+\]", "", display)
        display = re.sub(r"\[/PPT\]", "", display)
    return saved_ppts, display


def auto_wrap_and_save_fallback(user_message: str, ai_reply: str, user_id):
    """兜底：AI 忘记写 [文件:] / [PPT:] 标签时，自动检测内容类型并保存文件，保证一定能下载。
    返回 (extra_files, extra_ppts, display)"""
    import os as _os, re as _re, time as _time
    if not ai_reply or not ai_reply.strip():
        return [], [], ai_reply
    has_file = bool(_re.search(r"\[文件:\s*[^\]]+\]", ai_reply))
    has_ppt = bool(_re.search(r"\[PPT:\s*[^\]]+\]", ai_reply))
    if has_file or has_ppt:
        return [], [], ai_reply

    text = (user_message or "") + " " + (ai_reply or "")
    lower = text.lower()

    safe_title = _re.sub(r'[<>:"/\\|?*]', '', (user_message or "AI生成内容").strip())[:20] or "AI生成内容"
    ts = _time.strftime("%Y%m%d_%H%M%S")

    if ("ppt" in lower) or ("课件" in text) or ("幻灯片" in text):
        filename = f"课件_{safe_title}_{ts}.pptx"
        wrapped = f"[PPT:{filename}]\n{ai_reply.strip()}\n[/PPT]"
        saved_ppts, display = generate_ppt_files(wrapped, user_id)
        return [], saved_ppts, display

    if "作业" in text and ("分析" in text or "批改" in text or "错因" in text or "讲评" in text):
        fn = f"作业分析_{safe_title}_{ts}.md"
    elif "教案" in text or "教学设计" in text or "说课稿" in text or "教学过程" in text:
        fn = f"教案_{safe_title}_{ts}.md"
    elif "习题" in text or "练习" in text or "试卷" in text or "测试题" in text or ("题" in text and "答案" in text):
        fn = f"习题_{safe_title}_{ts}.md"
    elif ("讲解" in text and ("纲要" in text or "讲稿" in text or "逐页" in text)) or "说课稿" in text:
        fn = f"讲解纲要_{safe_title}_{ts}.md"
    elif "复习" in text or "提纲" in text or "知识点总结" in text or "知识体系" in text:
        fn = f"复习提纲_{safe_title}_{ts}.md"
    else:
        fn = f"备课资料_{safe_title}_{ts}.md"

    base_dir = get_user_output_dir(user_id)
    filepath = _os.path.join(base_dir, fn)
    content = ai_reply.strip()
    if content.startswith("```"):
        content = _re.sub(r"^```\w*\s*\n", "", content)
        content = _re.sub(r"\n```\s*$", "", content)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    return [fn], [], ai_reply


def extract_and_save_all(user_message, reply_text, user_id):
    """统一处理 AI 回复：解析文件块 + 生成 PPT + 兜底保存，返回(saved_files, saved_ppts, display_for_chat)"""
    user_output_dir = get_user_output_dir(user_id)

    saved_files = []
    files = re.findall(r"\[文件:\s*([^\]]+)\](.*?)\[/文件\]", reply_text, re.DOTALL)
    for filename, content in files:
        filepath = os.path.join(user_output_dir, filename)
        clean = content.strip()
        if clean.startswith("```"):
            clean = re.sub(r"^```\w*\s*\n", "", clean)
            clean = re.sub(r"\n```\s*$", "", clean)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(clean)
        saved_files.append(filename)

    display = reply_text
    if files:
        display = re.sub(r"\[文件:.*?\[/文件\]", "", reply_text, flags=re.DOTALL).strip()

    saved_ppts, display = generate_ppt_files(display, user_id)

    if not saved_files and not saved_ppts:
        extra_fs, extra_ppts, display = auto_wrap_and_save_fallback(
            user_message, reply_text, user_id
        )
        saved_files.extend(extra_fs)
        saved_ppts.extend(extra_ppts)

    return saved_files, saved_ppts, display


@app.route("/health")
def health_check():
    return {"status": "ok"}

@app.route("/")
def index():
    user = get_current_user()
    if not user:
        return redirect(url_for('login'))
    return render_template("index.html", model=MODEL, username=user['username'])

@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        data = safe_get_json()
        if not data:
            return jsonify({"error": "无效的请求格式"}), 400
        username = data.get("username")
        password = data.get("password")
        
        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400
        
        user = database.authenticate_user(username, password)
        if user:
            token = database.generate_session_token()
            database.set_session_token(user['id'], token)
            response = jsonify({"status": "success", "username": user['username']})
            response.set_cookie('session_token', token, httponly=True, secure=False, max_age=86400*7)
            return response
        else:
            return jsonify({"error": "用户名或密码错误"}), 401
    
    return render_template("login.html")

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        data = safe_get_json()
        if not data:
            return jsonify({"error": "无效的请求格式"}), 400
        username = data.get("username")
        password = data.get("password")
        email = data.get("email", "")
        
        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400
        
        if len(password) < 6:
            return jsonify({"error": "密码长度至少6位"}), 400
        
        success = database.register_user(username, password, email)
        if success:
            user = database.authenticate_user(username, password)
            if user:
                token = database.generate_session_token()
                database.set_session_token(user['id'], token)
                response = jsonify({"status": "success", "username": user['username']})
                response.set_cookie('session_token', token, httponly=True, secure=False, max_age=86400*7)
                return response
        else:
            return jsonify({"error": "用户名已存在"}), 409
    
    return render_template("register.html")

@app.route("/logout")
def logout():
    token = request.cookies.get('session_token')
    if token:
        database.clear_session_token(token)
    response = make_response(redirect(url_for('login')))
    response.set_cookie('session_token', '', expires=0)
    return response

@app.route("/api/chat", methods=["POST"])
def chat():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    session_id = str(data.get("session_id", "default"))[:64]
    user_message = data.get("message", "").strip()
    attachments = data.get("attachments") or []

    if isinstance(attachments, list):
        blocks = []
        for att in attachments:
            if not isinstance(att, dict):
                continue
            name = str(att.get("name", "未命名文件"))[:120]
            kind = str(att.get("kind", "资料"))[:30]
            text = str(att.get("text", ""))[:MAX_TEXT_CHARS]
            blocks.append(
                f"【用户上传文件：{name}（{kind}）】\n"
                f"文件内容预览（节选，共约 {len(text)} 字）：\n"
                f"```\n{text}\n```\n"
                f"【上传文件结束】"
            )
        if blocks:
            if user_message:
                user_message = user_message + "\n\n---\n\n" + "\n\n".join(blocks)
            else:
                user_message = "\n\n".join(blocks)

    if not user_message:
        return jsonify({"error": "Message cannot be empty"}), 400
    
    if len(user_message) > 24000:
        user_message = user_message[:24000] + "\n\n[内容过长已截断]"

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return jsonify({"error": "Invalid temperature or max_tokens value"}), 400

    key = _session_key(user['id'], session_id)
    if key not in sessions:
        sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]

    sessions[key].append({"role": "user", "content": user_message})

    database.save_search_history(user['id'], user_message)

    try:
        client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
        msgs = trim_history(sessions[key])
        response = client.chat.completions.create(
            model=_model,
            messages=msgs,
            temperature=_temp,
            max_tokens=_max_tokens,
        )
    except Exception as e:
        return jsonify({"error": f"API request failed: {str(e)}"}), 500

    reply = response.choices[0].message.content
    usage = response.usage

    sessions[key].append({"role": "assistant", "content": reply})

    first_user_msg = next((m["content"] for m in sessions[key] if m["role"] == "user"), "")
    title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
    database.save_user_chat_session(user['id'], session_id, title)

    saved_files, saved_ppts, display = extract_and_save_all(user_message, reply, user['id'])

    return jsonify({
        "reply": display,
        "raw_reply": reply,
        "model": response.model,
        "usage": {
            "prompt_tokens": usage.prompt_tokens,
            "completion_tokens": usage.completion_tokens,
        },
        "saved_files": saved_files,
        "saved_ppts": saved_ppts,
    })

def extract_keywords(message):
    keywords = []
    subjects = ["物理", "数学", "语文", "英语", "化学", "生物", "历史", "地理", "政治", "信息技术", "通用技术", "体育", "音乐", "美术", "科学", "道德与法治", "思想品德"]
    for subj in subjects:
        if subj in message:
            keywords.append(subj)
            break
    levels = ["小学", "初中", "高中", "大学", "职业院校", "中职", "高职", "一年级", "二年级", "三年级", "四年级", "五年级", "六年级", "初一", "初二", "初三", "高一", "高二", "高三"]
    for level in levels:
        if level in message:
            keywords.append(level)
            break
    types = ["教案", "课件", "习题", "试卷", "复习", "提纲", "实验", "实训", "说课稿", "导学案", "任务单", "教学设计"]
    for t in types:
        if t in message:
            keywords.append(t)
            break
    if not keywords:
        keywords = ["教学", "备课"]
    return keywords

knowledge_base = {
    "物理": ["牛顿运动定律", "能量守恒", "电磁感应", "光学", "热力学", "力学", "相对论", "量子物理", "波动", "磁场"],
    "数学": ["函数", "方程", "几何", "概率统计", "数列", "三角函数", "向量", "导数", "积分", "不等式"],
    "语文": ["文言文", "现代文", "诗词鉴赏", "写作", "阅读", "修辞手法", "表现手法", "作文", "记叙文", "议论文"],
    "英语": ["语法", "词汇", "阅读", "写作", "听力", "口语", "时态", "句型", "翻译", "完形填空"],
    "化学": ["元素周期表", "化学反应", "有机化学", "无机化学", "化学平衡", "电化学", "化学实验", "化学键"],
    "生物": ["细胞", "遗传", "生态", "光合作用", "呼吸作用", "进化论", "微生物", "人体生理"],
    "历史": ["中国古代史", "中国近代史", "世界史", "历史事件", "历史人物", "朝代", "战争", "改革"],
    "地理": ["自然地理", "人文地理", "气候", "地形", "地图", "环境保护", "区域地理", "人口城市"],
    "政治": ["经济生活", "政治生活", "文化生活", "哲学生活", "法律", "道德", "国情"],
    "default": ["教学目标", "教学重难点", "教学过程", "教学方法", "板书设计", "课后作业", "教学反思", "学情分析"]
}

def generate_thinking_phases(user_message):
    keywords = extract_keywords(user_message)
    subject = keywords[0] if keywords else "default"
    kbs = knowledge_base.get(subject, knowledge_base["default"])
    selected_kb = random.sample(kbs, min(4, len(kbs)))

    phases = [
        {"title": "📋 需求分析阶段", "items": [
            f"解析用户需求：{user_message[:25]}{'...' if len(user_message) > 25 else ''}",
            f"识别学科与学段：{'、'.join(keywords)}",
            "明确输出类型与格式要求",
        ]},
        {"title": "🔍 知识检索阶段", "items": [
            "在知识库中检索相关教学资源...",
            f"匹配知识点：{selected_kb[0]}、{selected_kb[1]}",
            "查找课程标准与教学大纲要求",
            f"筛选参考资料：{selected_kb[2]}相关案例",
        ]},
        {"title": "🧠 内容规划阶段", "items": [
            "梳理知识体系与逻辑结构",
            "设计三维教学目标",
            "确定教学重点与难点",
            "规划教学环节与时间分配",
            "选择教学方法与策略",
        ]},
        {"title": "✍️ 内容生成阶段", "items": [
            "组织语言，生成详细内容...",
            "优化表述，确保专业准确",
            "调整结构，保证逻辑清晰",
            "检查内容完整性与合理性",
        ]},
    ]
    return phases

def thinking_to_text(phases, phase_idx, item_idx):
    lines = []
    for i, phase in enumerate(phases):
        if i > phase_idx:
            break
        lines.append(f"【{phase['title']}】")
        for j, item in enumerate(phase['items']):
            if i == phase_idx and j > item_idx:
                break
            if i < phase_idx or (i == phase_idx and j < item_idx):
                icon = "✓"
            elif i == phase_idx and j == item_idx:
                icon = "⏳"
            else:
                icon = "○"
            lines.append(f"  {icon} {item}")
        lines.append("")
    return "\n".join(lines)

@app.route("/api/chat/stream", methods=["POST"])
def chat_stream():
    user = get_current_user()
    if not user:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '请先登录'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Invalid request body'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    session_id = str(data.get("session_id", "default"))[:64]
    user_message = data.get("message", "").strip()
    attachments = data.get("attachments") or []
    
    if isinstance(attachments, list):
        blocks = []
        for att in attachments:
            if not isinstance(att, dict):
                continue
            name = str(att.get("name", "未命名文件"))[:120]
            kind = str(att.get("kind", "资料"))[:30]
            text = str(att.get("text", ""))[:MAX_TEXT_CHARS]
            blocks.append(
                f"【用户上传文件：{name}（{kind}）】\n"
                f"文件内容预览（节选，共约 {len(text)} 字）：\n"
                f"```\n{text}\n```\n"
                f"【上传文件结束】"
            )
        if blocks:
            if user_message:
                user_message = user_message + "\n\n---\n\n" + "\n\n".join(blocks)
            else:
                user_message = "\n\n".join(blocks)

    if not user_message:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Message cannot be empty'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )
    
    if len(user_message) > 24000:
        user_message = user_message[:24000] + "\n\n[内容过长已截断]"

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Invalid temperature or max_tokens value'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    key = _session_key(user['id'], session_id)
    if key not in sessions:
        sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]

    sessions[key].append({"role": "user", "content": user_message})
    database.save_search_history(user['id'], user_message)

    def generate():
        thinking_phases = generate_thinking_phases(user_message)
        thinking_done = False
        ai_done = False
        ai_error = None
        full_reply = ""
        ai_queue = queue.Queue()
        first_content_received = False

        def ai_worker():
            nonlocal full_reply, ai_done, ai_error
            try:
                client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
                msgs = trim_history(sessions[key])
                stream = client.chat.completions.create(
                    model=_model,
                    messages=msgs,
                    temperature=_temp,
                    max_tokens=_max_tokens,
                    stream=True,
                )
                for chunk in stream:
                    if chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        full_reply += content
                        ai_queue.put(("content", content))
                ai_queue.put(("done", None))
                ai_done = True
            except Exception as e:
                ai_error = str(e)
                ai_queue.put(("error", str(e)))
                ai_done = True

        ai_thread = threading.Thread(target=ai_worker)
        ai_thread.start()

        total_thinking_time = 0
        min_thinking_time = 1.2
        max_thinking_time = 3.0

        for phase_idx, phase in enumerate(thinking_phases):
            if thinking_done:
                break
            for item_idx, item in enumerate(phase["items"]):
                if thinking_done:
                    break

                thinking_text = thinking_to_text(thinking_phases, phase_idx, item_idx)
                yield f"data: {json.dumps({'type': 'thinking', 'content': thinking_text}, ensure_ascii=False)}\n\n"

                sleep_time = random.uniform(0.25, 0.4)
                if total_thinking_time + sleep_time > max_thinking_time and phase_idx >= 2:
                    sleep_time = 0.1
                total_thinking_time += sleep_time
                time.sleep(sleep_time)

                if total_thinking_time >= min_thinking_time:
                    while not ai_queue.empty():
                        msg_type, msg_data = ai_queue.get()
                        if msg_type == "content":
                            first_content_received = True
                            break
                        elif msg_type == "done":
                            first_content_received = True
                            ai_done = True
                            break
                        elif msg_type == "error":
                            first_content_received = True
                            ai_done = True
                            break

                if first_content_received and total_thinking_time >= min_thinking_time:
                    thinking_done = True
                    final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
                    yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
                    break

        if not thinking_done:
            final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
            yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
            thinking_done = True

        if ai_error:
            yield f"data: {json.dumps({'type': 'error', 'content': ai_error}, ensure_ascii=False)}\n\n"
            return

        while not ai_done or not ai_queue.empty():
            try:
                msg_type, msg_data = ai_queue.get(timeout=0.1)
                if msg_type == "content":
                    yield f"data: {json.dumps({'type': 'content', 'content': msg_data}, ensure_ascii=False)}\n\n"
                elif msg_type == "done":
                    ai_done = True
            except queue.Empty:
                if not ai_done:
                    time.sleep(0.05)

        sessions[key].append({"role": "assistant", "content": full_reply})

        first_user_msg = next((m["content"] for m in sessions[key] if m["role"] == "user"), "")
        title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
        database.save_user_chat_session(user['id'], session_id, title)

        saved_files, saved_ppts, display = extract_and_save_all(user_message, full_reply, user['id'])

        final_thinking_text = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
        yield f"data: {json.dumps({'type': 'done', 'content': display, 'raw_content': full_reply, 'saved_files': saved_files, 'saved_ppts': saved_ppts, 'thinking': final_thinking_text}, ensure_ascii=False)}\n\n"

    return Response(generate(), mimetype="text/event-stream")

@app.route("/api/clear", methods=["POST"])
def clear():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json() or {}
    session_id = data.get("session_id", "default")
    key = _session_key(user['id'], session_id)
    sessions[key] = [{"role": "system", "content": SYSTEM_PROMPT}]
    return jsonify({"status": "ok"})

@app.route("/api/sessions", methods=["GET"])
def list_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    db_sessions = database.get_user_sessions(user['id'])
    result = []
    for s in db_sessions:
        result.append({
            "id": s['session_id'],
            "title": s['title'],
            "message_count": 0,
            "created_at": s['created_at']
        })
    return jsonify({"sessions": result})

@app.route("/api/session/<session_id>", methods=["GET"])
def get_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    key = _session_key(user['id'], session_id)
    if key in sessions:
        msgs = sessions[key]
        user_msgs = [m for m in msgs if m["role"] == "user"]
        title = "新对话"
        if user_msgs:
            first_user = user_msgs[0]["content"]
            title = first_user[:20] + "..." if len(first_user) > 20 else first_user
        return jsonify({
            "id": session_id,
            "title": title,
            "messages": msgs
        })
    
    db_sessions = database.get_user_sessions(user['id'])
    db_session = next((s for s in db_sessions if s['session_id'] == session_id), None)
    if db_session:
        return jsonify({
            "id": session_id,
            "title": db_session['title'],
            "messages": []
        })
    
    return jsonify({"error": "会话不存在"}), 404

@app.route("/api/session/<session_id>", methods=["DELETE"])
def delete_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    database.delete_user_session(user['id'], session_id)
    key = _session_key(user['id'], session_id)
    if key in sessions:
        del sessions[key]
    return jsonify({"status": "ok"})

@app.route("/api/sessions/batch", methods=["DELETE"])
def delete_sessions_batch():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json() or {}
    ids = data.get("ids", [])
    deleted = 0
    for sid in ids:
        database.delete_user_session(user['id'], sid)
        key = _session_key(user['id'], sid)
        if key in sessions:
            del sessions[key]
        deleted += 1
    return jsonify({"status": "ok", "deleted": deleted})

@app.route("/api/sessions/all", methods=["DELETE"])
def delete_all_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    prefix = f"{user['id']}:"
    to_delete = [key for key in sessions if key.startswith(prefix)]
    for key in to_delete:
        del sessions[key]
    db_deleted = database.delete_all_user_sessions(user['id'])
    return jsonify({"status": "ok", "memory_deleted": len(to_delete), "db_deleted": db_deleted})

@app.route("/api/search_history", methods=["GET"])
def get_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    history = database.get_search_history(user['id'], limit=50)
    return jsonify({"history": history})

@app.route("/api/search_history", methods=["DELETE"])
def clear_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    deleted = database.clear_search_history(user['id'])
    return jsonify({"status": "ok", "deleted": deleted})

@app.route("/api/config", methods=["GET"])
def get_config():
    return jsonify({
        "model": MODEL,
        "temperature": TEMP,
        "max_tokens": MAX_TOKENS,
    })

MIME_MAP = {
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.md': 'text/markdown; charset=utf-8',
    '.txt': 'text/plain; charset=utf-8',
    '.csv': 'text/csv; charset=utf-8',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.pdf': 'application/pdf',
}


def _read_text_bytes(b: bytes) -> str:
    for enc in ('utf-8-sig', 'utf-8', 'gbk', 'gb18030', 'utf-16'):
        try:
            return b.decode(enc)
        except UnicodeDecodeError:
            continue
    return b.decode('utf-8', errors='ignore')


def _guess_file_kind(name: str, text: str) -> str:
    low = (name + "\n" + text[:1500]).lower()
    template_keys = ['教案', '教学设计', '学案', '模板', '____', '___', '待填写', '留白', '[    ]', '[___]', '填写说明', '教学目标', '教学重难点']
    homework_keys = ['作业', '成绩', '得分', '分数', '错题', '批改', '班级', '学生', '姓名', '学号', '排名', '及格', '均分', '考勤', '总分']
    t_hit = sum(1 for k in template_keys if k in low)
    h_hit = sum(1 for k in homework_keys if k in low)
    if t_hit >= 2 and t_hit >= h_hit:
        return "教案/教学设计模板"
    if h_hit >= 2:
        return "学生作业/成绩/批改记录"
    if '.xlsx' in name.lower() or '.csv' in name.lower():
        return "数据表"
    if '.pdf' in name.lower():
        return "PDF文档"
    return "参考资料"


def _table_to_markdown(rows):
    rows = list(rows)
    if not rows:
        return ""
    rows = [[("" if v is None else str(v)).replace('|', '\\|').replace('\n', ' ') for v in r] for r in rows]
    max_cols = max(len(r) for r in rows)
    rows = [r + [''] * (max_cols - len(r)) for r in rows]
    header = rows[0]
    sep = ['---'] * max_cols
    body = rows[1:]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join(sep) + " |"]
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


@app.route("/api/upload", methods=["POST"])
def api_upload():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    if 'file' not in request.files:
        return jsonify({"error": "未收到上传文件字段 file"}), 400

    f = request.files['file']
    if not f or not f.filename:
        return jsonify({"error": "文件名为空"}), 400

    filename = os.path.basename(f.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXT:
        return jsonify({"error": f"不支持的格式 {ext}，支持：md / txt / docx / xlsx / csv / pdf"}), 400

    raw = f.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return jsonify({"error": f"文件过大 {len(raw)} bytes，上限 {MAX_UPLOAD_BYTES} bytes（20MB）"}), 400

    text = ""
    try:
        if ext in ('.txt', '.md'):
            text = _read_text_bytes(raw)
        elif ext == '.docx':
            bio = BytesIO(raw)
            doc = Document(bio)
            parts = []
            for p in doc.paragraphs:
                if p.text:
                    parts.append(p.text)
            for t in doc.tables:
                rows = [[c.text for c in row.cells] for row in t.rows]
                md = _table_to_markdown(rows)
                if md:
                    parts.append("\n[表格]\n" + md + "\n")
            text = "\n".join(parts)
        elif ext == '.csv':
            s = _read_text_bytes(raw)
            import csv as _csv
            reader = _csv.reader(s.splitlines())
            rows = list(reader)[:200]
            text = _table_to_markdown(rows) or (s[:MAX_TEXT_CHARS])
        elif ext == '.xlsx':
            import openpyxl
            bio = BytesIO(raw)
            wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
            sheets_text = []
            for ws in wb.worksheets:
                rows = []
                for idx, row in enumerate(ws.iter_rows(values_only=True)):
                    if idx > 60:
                        break
                    rows.append(list(row)[:16])
                md = _table_to_markdown(rows)
                if md:
                    sheets_text.append(f"# 工作表：{ws.title}\n{md}")
            text = "\n\n".join(sheets_text)
        elif ext == '.pdf':
            try:
                import PyPDF2
                bio = BytesIO(raw)
                reader = PyPDF2.PdfReader(bio)
                pages = []
                for i, page in enumerate(reader.pages[:30]):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    if t:
                        pages.append(f"\n---第 {i+1} 页---\n{t}")
                text = "\n".join(pages)
            except ImportError:
                return jsonify({"error": "PDF 需要先安装 PyPDF2：pip install PyPDF2"}), 400
    except Exception as ee:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"解析文件失败: {ee}"}), 500

    if not text or not text.strip():
        return jsonify({"error": "文件内容为空或无法识别文本"}), 400

    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n[原文过长已截断，仅保留前部分内容用于分析]"

    kind = _guess_file_kind(filename, text)
    preview = text[:600]

    return jsonify({
        "ok": True,
        "name": filename,
        "kind": kind,
        "chars": len(text),
        "preview": preview,
        "full_text": text,
    })


def _is_safe_filename(filename):
    if not filename or not isinstance(filename, str):
        return False
    base = os.path.basename(filename)
    if base != filename or base in ('.', '..') or not base:
        return False
    if any(c in filename for c in ('\\', '/', '..')):
        return False
    if filename.startswith('.'):
        return False
    return True


def _find_user_file(filename):
    if not _is_safe_filename(filename):
        return None
    user = get_current_user()
    if user:
        user_dir = os.path.realpath(get_user_output_dir(user['id']))
        fp = os.path.realpath(os.path.join(user_dir, filename))
        if os.path.isfile(fp) and fp.startswith(user_dir + os.sep):
            return fp
    if os.path.isdir(HISTORY_DIR):
        hist_real = os.path.realpath(HISTORY_DIR)
        for sub in os.listdir(HISTORY_DIR):
            sub_dir = os.path.realpath(os.path.join(hist_real, sub))
            if not os.path.isdir(sub_dir):
                continue
            fp = os.path.realpath(os.path.join(sub_dir, filename))
            if os.path.isfile(fp) and fp.startswith(sub_dir + os.sep):
                return fp
    out_real = os.path.realpath(OUTPUT_DIR)
    fp = os.path.realpath(os.path.join(out_real, filename))
    if os.path.isfile(fp) and fp.startswith(out_real + os.sep):
        return fp
    return None


@app.route("/api/download/<path:filename>")
def download(filename):
    from urllib.parse import unquote, quote as _qd
    filename = unquote(filename)
    filepath = _find_user_file(filename)
    if not filepath:
        return jsonify({"error": "文件不存在"}), 404

    ext = os.path.splitext(filename)[1].lower()
    mimetype = MIME_MAP.get(ext, 'application/octet-stream')
    try:
        safe_ascii = filename.encode('ascii', errors='ignore').decode('ascii') or ('download' + ext)
    except Exception:
        safe_ascii = 'download' + ext
    try:
        resp = send_file(filepath, as_attachment=True, download_name=safe_ascii, mimetype=mimetype)
        resp.headers['X-Filename-Encoded'] = _qd(filename)
        resp.headers['X-Filename'] = safe_ascii
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as ee:
        print(f"[download] 发送文件异常: {ee}")
        import traceback; traceback.print_exc()
        return jsonify({"error": f"文件下载失败: {ee}"}), 500


@app.route("/api/export/word", methods=["POST"])
def export_to_word():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    markdown_content = data.get("content", "").strip()
    title = str(data.get("title", "教案")).strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])

    try:
        word_buffer = markdown_to_word(markdown_content, title)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:20] or "untitled"
        filename = f"教案_{safe_title}_{timestamp}.docx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(word_buffer.getvalue())

        from urllib.parse import quote as _q
        from unicodedata import normalize as _norm
        try:
            _ascii_fn = filename.encode('ascii', errors='ignore').decode('ascii') or 'document.docx'
        except Exception:
            _ascii_fn = 'document.docx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn,
                         mimetype=MIME_MAP['.docx'])
        resp.headers['X-Filename-Encoded'] = _q(filename)
        resp.headers['X-Filename'] = _ascii_fn
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"导出失败: {str(e)}"}), 500

@app.route("/api/export/ppt", methods=["POST"])
def export_to_ppt():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = safe_get_json()
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    markdown_content = data.get("content", "").strip()
    title = str(data.get("title", "课件")).strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])

    try:
        ppt_buffer = markdown_to_ppt(markdown_content, title)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:20] or "untitled"
        filename = f"课件_{safe_title}_{timestamp}.pptx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(ppt_buffer.getvalue())

        from urllib.parse import quote as _q2
        try:
            _ascii_fn2 = filename.encode('ascii', errors='ignore').decode('ascii') or 'slides.pptx'
        except Exception:
            _ascii_fn2 = 'slides.pptx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn2,
                         mimetype=MIME_MAP['.pptx'])
        resp.headers['X-Filename-Encoded'] = _q2(filename)
        resp.headers['X-Filename'] = _ascii_fn2
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"导出 PPT 失败: {str(e)}"}), 500

MODULES_DATA = {
    "ppt_generate": {
        "name": "PPT生成",
        "emoji": "📊",
        "short_desc": "一键生成教学课件PPT",
        "description": "AI智能生成完整教学课件，支持从课程主题直接输出结构化PPT大纲，并可一键导出为 .pptx 文件。自动包含封面、教学目标、知识点讲解、案例分析、课堂小结、课后作业等标准课件结构。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "生成一份《荷塘月色》的语文课件PPT，包含教学目标、作者介绍、课文赏析、重点语句分析、课堂练习、作业布置，共15-20页",
        "features": [
            "📊 标准课件结构（封面/目标/新知/小结/作业）",
            "📝 每页幻灯片内容完整撰写（标题+要点）",
            "🎨 Markdown格式规范排版，方便二次编辑",
            "⬇️ 一键导出 .pptx 文件，可直接打开使用",
            "🔧 自动适配学科特点（文科赏析/理科例题/工科原理）",
            "⏱ 支持指定页数与各环节时间分配"
        ],
        "examples": [
            {"title": "语文课件：荷塘月色", "desc": "朱自清散文，含作者背景+课文赏析+修辞分析", "prompt": "生成一份《荷塘月色》的语文课件PPT，适合高中，包含教学目标、作者介绍、课文逐段赏析、修辞手法分析、课堂练习、作业布置，共15-20页"},
            {"title": "数学课件：二次函数", "desc": "初中数学，含图像性质+典型例题+课堂练习", "prompt": "生成初中数学《二次函数的图像与性质》课件PPT，包含概念讲解、5个典型例题解析、课堂练习、课后作业，约18页"},
            {"title": "历史课件：辛亥革命", "desc": "高中历史，含背景脉络+重大事件+意义启示", "prompt": "生成高中历史《辛亥革命》课件PPT，包含历史背景、武昌起义经过、中华民国成立、历史意义与局限、思考题，约20页"}
        ],
        "knowledge_points": ["课件结构", "教学目标", "知识点导入", "案例分析", "课堂互动", "例题解析", "归纳小结", "作业布置", "板书设计", "时间分配"]
    },
    "ppt_outline": {
        "name": "PPT讲解纲要",
        "emoji": "🎤",
        "short_desc": "每页PPT的逐页讲解稿",
        "description": "根据PPT大纲为每一页幻灯片生成详细的教师讲解词。包含开场白设计、知识点过渡衔接语、重难点突出提示、提问互动设计、随堂练习引导语等，帮助教师脱稿流畅完成整堂课。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、互动提示、过渡语句、重点强调部分",
        "features": [
            "🎤 逐页讲解词撰写，教师可直接照着读",
            "🔗 页面之间设计自然衔接过渡语句",
            "❓ 穿插课堂提问与互动设计（教师问+预期答）",
            "⭐ 标注重点知识点的语气语调提示（重读、停顿）",
            "⏰ 每页推荐讲解时长，合理分配45分钟",
            "🧑‍🏫 模拟真实课堂的教态与现场感"
        ],
        "examples": [
            {"title": "语文：荷塘月色讲解稿", "desc": "逐页教师用语+情感渲染引导", "prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、课堂提问设计、过渡语句、重点部分标注语气要求"},
            {"title": "数学：二次函数讲解稿", "desc": "概念引入+例题推导步骤+易错提醒", "prompt": "生成初中数学《二次函数图像与性质》每一页PPT的教师讲解纲要，包含概念引入话术、例题逐步推导讲解、学生易错点提醒、课堂练习的处理方式"},
            {"title": "班会：防溺水安全教育稿", "desc": "案例警示+互动问答+总结升华", "prompt": "生成中小学防溺水主题班会PPT的逐页讲解纲要，包含真实案例描述、互动提问、自救技巧讲解、最后的倡议环节，语言符合中小学生认知"}
        ],
        "knowledge_points": ["开场白设计", "过渡衔接", "重难点提示", "互动提问", "案例讲述", "语气语调", "时间把控", "课堂节奏", "情感升华", "结尾总结"]
    },
    "lesson_plan": {
        "name": "教案生成",
        "emoji": "📘",
        "short_desc": "标准格式45分钟完整教案",
        "description": "按中小学教案标准格式自动生成完整45分钟教学设计。包含教学目标（三维/核心素养）、教学重难点、教学方法、教学过程（导入/新授/练习/小结/作业）、板书设计、教学反思模板等全部要素。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "中职", "高职"],
        "default_prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程、板书设计、作业布置、教学反思",
        "features": [
            "🎯 三维教学目标（知识与技能/过程与方法/情感态度）",
            "🎯 核心素养目标适配新课标要求",
            "❗ 教学重点与教学难点明确区分",
            "⏱️ 详细教学过程：5-6个环节精确到分钟",
            "🎨 板书设计框架（主板书+副板书）",
            "📋 配套作业设计（基础+拓展分层）",
            "📝 预留教学反思栏（课堂生成问题记录）"
        ],
        "examples": [
            {"title": "语文：背影教案", "desc": "朱自清经典·情感教育·细节描写", "prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程（各环节含时间分配）、板书设计、分层作业布置、教学反思栏"},
            {"title": "物理：浮力教案", "desc": "实验探究·阿基米德原理·分层练习", "prompt": "生成初中物理《阿基米德原理》45分钟教案，包含实验探究环节设计、演示实验步骤、学生分组活动安排、分层练习设计、板书设计"},
            {"title": "英语：Travel Plan教案", "desc": "听说课型·情境交际·任务型教学", "prompt": "生成初中英语听说课《Travel Plans》完整教案，包含Warm-up、Pre-listening、While-listening、Post-speaking、Summary & Homework环节，配师生对话示例"}
        ],
        "knowledge_points": ["三维目标", "核心素养", "教学重难点", "学情分析", "教学方法", "教学过程", "时间分配", "板书设计", "分层作业", "教学反思"]
    },
    "exercises": {
        "name": "习题生成",
        "emoji": "📝",
        "short_desc": "分层习题+详细答案解析",
        "description": "根据知识点自动生成多难度层次练习题，支持选择题、填空题、判断题、解答题、应用题等多种题型。每道题目均附标准答案与详细解题步骤，可用于课堂练习、课后作业、单元测验。",
        "subject": "全学科",
        "grade_levels": ["小学低年级", "小学高年级", "初中", "高中", "中职", "高职"],
        "default_prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用，所有题目附答案和详细解析",
        "features": [
            "📊 三级难度分层：基础巩固·能力提高·拓展探究",
            "✅ 题型丰富：选择/填空/判断/解答/实验探究/应用题",
            "📝 每题附标准答案 + 详细解题步骤",
            "📌 标注每题考查知识点对应章节",
            "⏱ 推荐完成时长与分值设置",
            "📄 可一键组装为标准试卷（卷头+姓名栏）"
        ],
        "examples": [
            {"title": "数学：勾股定理20题", "desc": "8基础+8提高+4拓展，附详解", "prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用题，所有题目附答案和详细解析步骤"},
            {"title": "英语：一般过去时语法", "desc": "30道选择+填空+句型转换，解析全", "prompt": "生成初中英语语法《一般过去时》专项练习题30道：单项选择15道，用所给动词适当形式填空10道，句型转换5道，附答案和详细解析"},
            {"title": "化学：酸碱盐专题", "desc": "选择+填空+推断+实验探究四大题型", "prompt": "生成初中化学《酸碱盐》单元练习题，包含选择题10道、填空题6道、物质推断题2道、实验探究题2道，附答案与解析"}
        ],
        "knowledge_points": ["基础巩固", "能力提高", "拓展探究", "选择题", "填空题", "解答题", "实验探究", "应用题", "答案解析", "考点分布"]
    },
    "homework_analysis": {
        "name": "作业分析",
        "emoji": "📊",
        "short_desc": "作业批改与学情数据分析",
        "description": "上传作业统计数据或描述作业情况，AI自动生成完整的作业分析报告。包含：正确率统计、高频错题归因、典型错解展示、知识点掌握情况雷达图、学困生针对性辅导建议、下节课教学调整策略等。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "default_prompt": "请为一次高中数学《函数单调性》课后作业生成作业分析报告：全班45人，平均得分率68%，错误集中在：含参数的单调性讨论、复合函数单调性、实际应用最值问题，请给出详细分析与辅导建议",
        "features": [
            "📈 整体成绩统计：平均分/得分率/分数段分布",
            "❌ 高频错题 TOP 5 归因分析（概念/审题/计算/方法）",
            "📝 典型错解案例展示与正解对比",
            "🎯 知识点掌握度雷达图（熟练/一般/薄弱）",
            "👨‍🎓 分层辅导建议：学困生、中等生、优等生",
            "✏️ 下一步教学调整：需补讲的知识点、课堂策略",
            "📑 讲评课时分配建议与讲评例题推荐"
        ],
        "examples": [
            {"title": "数学函数单调性作业分析", "desc": "得分率68%·参数讨论错误高发", "prompt": "为高中数学《函数单调性》课后作业生成详细作业分析报告：全班45人，平均得分率68%，主要错误：①含参数单调性分类讨论（失分率52%）②复合函数同增异减应用（失分率41%）③实际问题求最值忽略定义域（失分率37%），请给出错因分析、辅导建议、教学调整策略"},
            {"title": "英语完形填空作业分析", "desc": "上下文逻辑·固定搭配是薄弱点", "prompt": "分析一次八年级英语完形填空作业情况：全班40人，平均正确率58%；错题中上下文逻辑理解占45%，固定搭配与词组占30%，词汇辨析占15%，语法占10%。请生成完整作业讲评方案"},
            {"title": "物理电路作业错误分析", "desc": "串并联识别+欧姆定律综合应用薄弱", "prompt": "为初三物理《欧姆定律在串并联电路中的应用》作业做分析：全班48人，平均得分62分；常见错：复杂电路等效化简错误、电表测量对象判断错误、比例计算错误。生成作业分析与讲评建议"}
        ],
        "knowledge_points": ["得分率统计", "错题归因", "典型错解", "知识薄弱点", "分层辅导", "学困生帮扶", "教学调整", "讲评课时", "补充练习", "家校沟通"]
    },
    "history": {
        "name": "历史记录",
        "emoji": "🕒",
        "short_desc": "查看所有AI生成记录与历史会话",
        "description": "查看当前账号下所有历史生成记录：历史会话聊天、PPT课件、教案、习题、作业分析等全部生成内容。支持按时间、按模块筛选，快速找到过往生成内容，可再次打开、复用、编辑或重新生成。",
        "subject": "系统功能",
        "grade_levels": ["全部学段"],
        "default_prompt": "（点击下方会话卡片直接查看历史记录）",
        "features": [
            "🕒 按时间倒序展示所有历史会话",
            "🔍 支持按模块类型筛选（PPT/教案/习题/作业等）",
            "📋 可快速查看每条记录的生成时间、标题、内容摘要",
            "♻️ 一键复用：把历史内容重新发送，基于结果再修改",
            "🗂 可导出历史记录为 Word / 文本文件",
            "📊 查看个人使用统计（生成次数、常用模块等）"
        ],
        "examples": [
            {"title": "查看最近7天生成记录", "desc": "快速找到上周生成的PPT并修改", "prompt": "__HISTORY__VIEW__"},
            {"title": "按模块筛选：只看教案", "desc": "查找所有历史生成的教案", "prompt": "__HISTORY_VIEW__ lesson_plan"},
            {"title": "统计本月使用情况", "desc": "查看自己备课模块使用分布", "prompt": "__HISTORY_STATS__"}
        ],
        "knowledge_points": ["聊天记录", "生成文件", "按时间筛选", "按模块筛选", "内容搜索", "记录复用", "批量导出", "使用统计", "收藏夹", "回收站"]
    },
    "visualization": {
        "name": "可视化生成",
        "emoji": "🎬",
        "short_desc": "教学动画与思维导图生成",
        "description": "输入课程知识点或教学目标，AI 自动为您生成：\n✅ HTML5教学演示动画（CSS动态效果、无需安装依赖、可嵌入PPT）\n✅ 一目了然的课程思维导图（交互式、可缩放、关键词精炼）\n支持小学至高职全学段，覆盖数理化生、语文英语等主流学科。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "default_prompt": "请为高中物理《牛顿第一定律》新课教学生成动画和思维导图。",
        "features": [
            "🎬 HTML5教学动画 - 纯CSS动画效果，无需安装Manim/FFmpeg等依赖，点击即生成，支持重播",
            "🧠 交互式思维导图 - markmap渲染，支持缩放、节点展开折叠，自动配色",
            "📐 AI智能布局 - 根据学科内容自动生成动画元素位置、颜色、文字和动画效果",
            "⚡ 极速生成 - 无需等待视频渲染，秒级生成结果",
            "🎨 双栏布局 - 动画和思维导图同屏展示，方便教学使用"
        ],
        "examples": [
            {"title": "牛顿第一定律", "desc": "动画展示物体静止/匀速运动状态；思维导图含定义/条件/推论/应用", "prompt": "牛顿第一定律"},
            {"title": "光合作用", "desc": "动画展示光反应→暗反应过程；导图分光反应/暗反应/影响因素/应用", "prompt": "光合作用"},
            {"title": "勾股定理", "desc": "动画展示直角三角形三边关系；导图含定理/证明/应用/拓展", "prompt": "勾股定理"}
        ],
        "knowledge_points": ["教学动画", "思维导图", "PPT素材", "知识点可视化", "HTML5动画", "Markmap导图", "学科建模", "课堂演示", "教师提效"]
    }
}

@app.route("/visual")
def visual_redirect():
    return redirect(url_for('module_detail', module_id='visualization'))

@app.route("/module/<module_id>")
def module_detail(module_id):
    user = get_current_user()
    if not user:
        print(f"[module_detail] 用户未登录，重定向到登录页")
        return redirect(url_for('login'))
    print(f"[module_detail] 用户: {user['username']}, module_id: {module_id}")
    
    if module_id not in MODULES_DATA:
        print(f"[module_detail] 模块不存在: {module_id}")
        return redirect(url_for('index'))
    module = MODULES_DATA[module_id]
    
    if module_id == "visualization":
        try:
            all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
            print(f"[module_detail] 渲染可视化页面，模块名称: {module['name']}")
            return render_template(
                "module_visual.html",
                module=module,
                module_id=module_id,
                model=MODEL,
                username=user['username'],
                all_modules=all_modules
            )
        except Exception as e:
            import traceback
            traceback.print_exc()
            print(f"[module_detail] 渲染可视化页面失败: {e}")
            return f"渲染失败: {str(e)}", 500
    
    all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
    history_sessions = []
    if module_id == "history":
        try:
            MODULE_META = {
                "ppt_generate":         {"name": "PPT生成",     "emoji": "📊", "color": "#667eea", "gradient": "linear-gradient(135deg,#667eea 0%,#764ba2 100%)"},
                "ppt_outline":          {"name": "PPT讲解纲要", "emoji": "📝", "color": "#f093fb", "gradient": "linear-gradient(135deg,#f093fb 0%,#f5576c 100%)"},
                "lesson_plan_generate": {"name": "教案生成",     "emoji": "📘", "color": "#4facfe", "gradient": "linear-gradient(135deg,#4facfe 0%,#00f2fe 100%)"},
                "exercise_generate":    {"name": "习题生成",     "emoji": "📚", "color": "#43e97b", "gradient": "linear-gradient(135deg,#43e97b 0%,#38f9d7 100%)"},
                "homework_analysis":    {"name": "作业分析",     "emoji": "📊", "color": "#fa709a", "gradient": "linear-gradient(135deg,#fa709a 0%,#fee140 100%)"},
                "history":              {"name": "历史记录",     "emoji": "📜", "color": "#a8edea", "gradient": "linear-gradient(135deg,#a8edea 0%,#fed6e3 100%)"},
                "visualization":        {"name": "可视化生成",   "emoji": "🎬", "color": "#8b5cf6", "gradient": "linear-gradient(135deg,#8b5cf6 0%,#ec4899 100%)"},
                "general":              {"name": "自由聊天",     "emoji": "💬", "color": "#c3cfe2", "gradient": "linear-gradient(135deg,#c3cfe2 0%,#f5f7fa 100%)"},
            }
            stats = database.get_user_sessions_with_stats(user['id'])
            for s in stats:
                mid = s.get("module_id") or "general"
                if mid not in MODULE_META:
                    mid = "general"
                meta = MODULE_META[mid]
                files = s.get("files") or []
                badge_summary = {}
                for f in files:
                    k = f.get("kind") or "文件"
                    badge_summary.setdefault(k, {"count": 0, "samples": []})
                    if badge_summary[k]["count"] < 3:
                        badge_summary[k]["samples"].append(f.get("name") or "")
                    badge_summary[k]["count"] += 1
                badges = []
                for k, v in badge_summary.items():
                    ext_icon = {"PPT": "📊", "教案": "📘", "练习册": "📚", "作业分析": "📈", "报告": "📄", "文件": "📄"}.get(k, "📄")
                    badges.append({
                        "kind": k,
                        "icon": ext_icon,
                        "count": v["count"],
                        "sample": (v["samples"][0] or "")[:50]
                    })
                history_sessions.append({
                    "id": s["session_id"],
                    "pk": s.get("pk_id") or "",
                    "title": s.get("title") or "(空对话)",
                    "created_at": s.get("created_at") or "",
                    "updated_at": s.get("updated_at") or "",
                    "time": s.get("updated_at") or s.get("created_at") or "",
                    "msg_count": s.get("total_count") or 0,
                    "user_count": s.get("user_count") or 0,
                    "assistant_count": s.get("assistant_count") or 0,
                    "file_count": s.get("file_count") or 0,
                    "module_id": mid,
                    "module_name": meta["name"],
                    "module_emoji": meta["emoji"],
                    "module_color": meta["color"],
                    "module_gradient": meta["gradient"],
                    "preview": s.get("user_preview") or s.get("ai_preview") or "",
                    "user_preview": s.get("user_preview") or "",
                    "ai_preview": s.get("ai_preview") or "",
                    "badges": badges,
                    "files": files[:8],
                })
            history_sessions.sort(key=lambda x: x.get("time") or "", reverse=True)
        except Exception as e:
            import traceback; traceback.print_exc()
            print("[history] 读取会话失败：", e)
    return render_template(
        "module_detail.html",
        module=module,
        module_id=module_id,
        model=MODEL,
        username=user['username'],
        all_modules=all_modules,
        history_sessions=history_sessions
    )

@app.route("/api/visual/generate", methods=["POST"])
def api_visual_generate():
    try:
        print(f"[api_visual_generate] 请求收到")
        user = get_current_user()
        print(f"[api_visual_generate] 当前用户: {user}")
        if not user:
            return jsonify({"success": False, "error": "请先登录"}), 401
    
        data = safe_get_json()
        print(f"[api_visual_generate] 请求数据: {data}")
        if not data:
            return jsonify({"success": False, "error": "无效请求"}), 400
        
        topic = data.get("topic", "").strip()
        gen_type = data.get("type", "mindmap")
        layout = data.get("layout", "logic")
        
        if not topic:
            return jsonify({"success": False, "error": "请输入课程主题"}), 400
        
        client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
        
        if gen_type == "mindmap":
            # 根据布局类型生成不同的prompt
            if layout == "tree":
                layout_desc = "组织结构图/树状图"
                format_hint = """# 角色定义
你是一位精通组织管理与信息分类的可视化专家。你必须严格区分并遵循「组织结构图」与「树状图」的专属规范，禁止将其与流程图、概念图或发散式思维导图混淆。

# 核心定义
- 组织结构图（Org Chart）：以垂直自上而下为主，展示职位/部门间的汇报、隶属与管理跨度关系，强调权责层级。
- 树状图（Tree Diagram）：可垂直或水平展开，展示事物的分类、组成或分解关系，强调MECE（相互独立、完全穷尽）原则。

# 格式规范（必须100%遵守）
1. 使用 Mermaid flowchart TD（自上而下）格式
2. 节点形状：仅允许矩形方框 []，根节点/高管节点可用加粗边框，禁止圆形、椭圆、菱形、云朵形、圆角矩形
3. 连线样式：仅允许直线 --> 表示父子关系，禁止曲线、波浪线、虚线
4. 布局方向：强制自上而下（TD），同一层级节点水平对齐，上级居中于下级上方
5. 层级表达：通过分支位置体现父子级，最多3-4级深度，跨级连接需标注说明
6. 文字规范：节点内文字≤8字，名词/职位名称优先，禁止动词长句、解释性描述
7. 至少3-5个一级分支，每个分支下有2-4个子节点
8. 同级节点必须严格对齐，禁止斜线连接"""
                output_format = "mermaid"
            elif layout == "rect":
                layout_desc = "矩形思维导图"
                format_hint = """# 角色定义
你是一位精通结构化思维的可视化专家。当用户提到“矩形思维导图”时，你必须将其理解为「以矩形为节点、直线为连接的层级拆解图」，禁止生成传统曲线发散式思维导图。

# 核心定义
矩形思维导图 = 思维导图的层级逻辑 + 逻辑图的视觉规范。
它用于对主题进行MECE拆解、分类归纳或流程梳理，强调结构的严谨性与信息的可读性，而非创意发散。

# 格式规范（必须100%遵守）
1. 使用 Mermaid flowchart LR（左→右）格式，内容较多时可用TD（上→下）
2. 节点形状：所有节点（含中心主题）均为矩形方框 []，禁止圆形、椭圆、云朵形、圆角矩形
3. 连线样式：仅允许直线 --> 表示父子关系，禁止任何曲线、波浪线、手绘线条
4. 布局方向：优先左→右展开（中心主题在最左侧），内容较多时可用上→下；禁止中心辐射状布局
5. 对齐规则：同级节点严格垂直/水平对齐，父子节点间距一致，整体呈现网格化秩序感
6. 文字规范：节点内文字≤8字，关键词/短语优先，禁止长句、解释性描述、emoji
7. 至少3-5个一级分支，每个分支下有2-4个子节点
8. 同级节点必须严格对齐，禁止斜线连接"""
                output_format = "mermaid"
            else:
                layout_desc = "逻辑图/逻辑结构图，以矩形方框为节点、直线/直角折线为连接、从左到右或从上到下呈现层级递进关系的结构化图表，强调因果、流程、分类或组成关系"
                format_hint = """格式要求（必须100%遵守）：
1. 使用 Mermaid flowchart LR（从左到右）格式
2. 节点形状：仅允许矩形方框 []，禁止圆形、椭圆、云朵形、圆角矩形、手绘边框
3. 连线样式：仅允许直线 --> 或直角折线 ---，禁止曲线、波浪线；备选路径用虚线 -.-> 表示
4. 布局方向：仅允许左→右（LR）或上→下（TD），禁止中心辐射、环形、自由散点布局
5. 层级表达：通过分支位置体现父子级，同一层级节点水平/垂直对齐
6. 文字规范：节点内文字≤10字，动宾结构优先，禁止长句、解释性描述"""
                output_format = "mermaid"

            if output_format == "mermaid":
                if layout == "tree":
                    prompt = f"""请为课程主题"{topic}"生成一个{layout_desc}的Mermaid流程图。

{format_hint}

# 正例格式锚点（必须模仿此结构）

【组织结构图示例】
flowchart TD
    A[总经理] --> B[市场部]
    A --> C[技术部]
    A --> D[人事部]
    B --> E[品牌组]
    B --> F[推广组]
    C --> G[后端组]
    C --> H[前端组]

【树状图示例】
flowchart TD
    A[电子产品] --> B[手机]
    A --> C[电脑]
    B --> D[智能手机]
    B --> E[功能手机]
    C --> F[笔记本]
    C --> G[台式机]

# 反例（禁止生成）
- 用箭头表示"市场部→技术部"的协作关系（× 这是流程图）
- 中心写"公司"，四周发散出各部门曲线分支（× 传统思维导图）
- 同级节点未对齐、连线为斜线或曲线（× 非标准树状结构）
- 节点内写"负责产品推广与品牌建设"等长句（× 违反文字规范）
- 使用()、{{}}、>等非矩形节点形状（× 违反节点形状规范）

# 执行指令
1. 先判断场景：涉及职位/汇报→组织结构图；涉及分类/拆解→树状图
2. 输出时必须严格按上述正例格式锚点的结构呈现
3. 若课程内容不满足MECE或层级混乱，主动调整使分类相互独立、完全穷尽
4. 强制采用自上而下布局（TD），永远不要使用LR/RL/BT
5. 永远不要添加"如图所示""参见下图"等无法渲染的描述，仅输出结构化文本
6. 节点文字必须≤8字，名词优先，禁止动词长句

直接输出Mermaid代码，不要其他解释，不要代码块包裹

节点命名规则：使用A、B、C等字母作为节点ID，方括号内写中文内容（≤8字）。
连接线规则：仅使用-->表示父子隶属关系，禁止曲线和虚线。
至少3-5个核心节点，形成完整的层级结构。
"""
                elif layout == "rect":
                    prompt = f"""请为课程主题"{topic}"生成一个{layout_desc}的Mermaid流程图。

{format_hint}

# 正例格式锚点（必须模仿此结构）

【左→右矩形思维导图示例】
flowchart LR
    A[新媒体运营] --> B[内容生产]
    A --> C[渠道分发]
    A --> D[数据复盘]
    B --> E[选题策划]
    B --> F[脚本撰写]
    B --> G[视觉设计]
    G --> H[封面制作]
    G --> I[排版美化]
    C --> J[微信公众号]
    C --> K[视频号]
    C --> L[小红书]
    D --> M[阅读量分析]
    D --> N[转化率优化]

【上→下矩形思维导图示例（适用于宽层级）】
flowchart TD
    A[年度营销计划] --> B[Q1拉新]
    A --> C[Q2留存]
    A --> D[Q3变现]
    B --> E[社媒投放]
    B --> F[KOL合作]
    C --> G[会员体系]
    D --> H[直播带货]
    D --> I[私域转化]

# 反例（禁止生成）
- 中心写"运营"，四周用曲线发散出"内容""渠道"等分支（× 传统博赞式思维导图）
- 节点为圆角矩形或带阴影立体效果（× 非标准矩形）
- 同级节点未对齐、连线为斜线或自由曲线（× 失去网格秩序）
- 节点内写"负责公众号推文撰写与排版"等长句（× 违反文字规范）
- 使用()、{{}}、>等非矩形节点形状（× 违反节点形状规范）

# 执行指令
1. 自动将需求转换为左→右或上→下的矩形层级结构
2. 输出时必须严格按格式锚点的符号与对齐方式呈现
3. 若用户提供的内容存在层级交叉或非MECE问题，主动指出并建议调整后再输出
4. 优先采用左→右布局（LR），分支较多时可用上→下（TD）
5. 永远不要添加"如图所示""见下图"等无法渲染的描述，仅输出纯结构化文本
6. 节点文字必须≤8字，关键词/短语优先，禁止长句、解释性描述、emoji

直接输出Mermaid代码，不要其他解释，不要代码块包裹

节点命名规则：使用A、B、C等字母作为节点ID，方括号内写中文内容（≤8字）。
连接线规则：仅使用-->表示父子关系，禁止曲线和虚线。
至少3-5个核心节点，形成完整的层级结构。
"""
                else:
                    prompt = f"""请为课程主题"{topic}"生成一个{layout_desc}的Mermaid流程图。

{format_hint}

正例（逻辑图）：
flowchart LR
    A[用户注册] --> B[手机号验证] --> C[设置密码] --> D[完成注册]
    B -.-> E[邮箱验证]

反例（禁止生成）：
- 中心写"注册"，四周发散出"手机""邮箱""密码"等曲线分支（× 传统思维导图）
- 用圆角矩形+箭头表示"验证成功/失败"的判断菱形（× 流程图）
- 节点为手绘体圆圈，连线为彩色曲线（× 博赞式思维导图）

执行指令：
1. 先确认内容是否适合逻辑图（若为纯创意发散，主动建议改用其他形式）
2. 输出时必须严格按上述格式规范组织文本结构
3. 默认采用左→右布局（LR）
4. 永远不要添加"如图所示""参见下图"等无法渲染的描述

直接输出Mermaid代码，不要其他解释，不要代码块包裹

节点命名规则：使用A、B、C等字母，或简短中文（不超过4字）。
连接线规则：使用-->表示主要路径，-.->表示备选路径。
至少3-5个核心节点，形成完整的逻辑链。
"""
            else:
                prompt = f"""请为课程主题"{topic}"生成一个{layout_desc}的Markdown思维导图，{format_hint}

直接输出Markdown，不要其他解释，不要代码块包裹

示例格式：
# {topic}
- 核心概念1
  - 细节1
  - 细节2
- 核心概念2
  - 细节1
  - 细节2
"""
            response = client.chat.completions.create(
                model=MODEL,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=2000
            )
            content = response.choices[0].message.content.strip()
            
            content = re.sub(r'^```\w*\s*', '', content)
            content = re.sub(r'\s*```$', '', content)
            
            if output_format == "markdown" and not content.startswith('#'):
                content = f"# {topic}\n" + content
            
            database.save_search_history(user['id'], f"可视化生成-思维导图({layout})：{topic}")
            
            return jsonify({
                "success": True,
                "type": "mindmap",
                "layout": layout,
                "format": output_format,
                "content": content
            })
            
        elif gen_type == "animation":
            if QWEN_IMAGE_API_KEY:
                try:
                    import dashscope
                    from dashscope import ImageSynthesis

                    dashscope.api_key = QWEN_IMAGE_API_KEY

                    anim_prompt = f"""教学主题"{topic}"的教学演示插图，用于课件展示，风格：教育类插画，清晰简洁，适合课堂教学使用"""

                    result = ImageSynthesis.call(
                        model=ImageSynthesis.Models.wanx_v1,
                        prompt=anim_prompt,
                        n=1,
                        size='1280*720'
                    )

                    if result.status_code == 200 and result.output and result.output.task_status == 'SUCCEEDED' and result.output.results:
                        image_url = result.output.results[0].url

                        anim_html = f'''<div class="classroom-animation" style="width:500px;height:320px;background:transparent;display:flex;align-items:center;justify-content:center;">
                            <img src="{image_url}" alt="{topic}" style="max-width:100%;max-height:100%;border-radius:8px;" />
                        </div>'''

                        database.save_search_history(user['id'], f"可视化生成-Qwen图像：{topic}")

                        return jsonify({
                            "success": True,
                            "type": "animation",
                            "format": "html",
                            "html": anim_html,
                            "image_url": image_url
                        })
                    else:
                        print(f"[api_visual_generate] Qwen-Image生成失败: {result}")
                except Exception as e:
                    print(f"[api_visual_generate] Qwen-Image调用异常: {e}")
            
            search_images = []
            if IMAGE_SEARCH_PROVIDER:
                try:
                    import requests
                    
                    search_query = topic
                    if IMAGE_SEARCH_PROVIDER.lower() == "unsplash":
                        url = f"https://api.unsplash.com/search/photos?query={requests.utils.quote(search_query)}&per_page=3"
                        headers = {"Authorization": f"Client-ID {IMAGE_SEARCH_API_KEY}"} if IMAGE_SEARCH_API_KEY else {}
                        response = requests.get(url, headers=headers, timeout=10)
                        if response.status_code == 200:
                            data = response.json()
                            if data.get("results"):
                                search_images = [img["urls"]["regular"] for img in data["results"][:3]]
                    elif IMAGE_SEARCH_PROVIDER.lower() == "pexels":
                        url = f"https://api.pexels.com/v1/search?query={requests.utils.quote(search_query)}&per_page=3"
                        headers = {"Authorization": IMAGE_SEARCH_API_KEY} if IMAGE_SEARCH_API_KEY else {}
                        response = requests.get(url, headers=headers, timeout=10)
                        if response.status_code == 200:
                            data = response.json()
                            if data.get("photos"):
                                search_images = [img["src"]["medium"] for img in data["photos"][:3]]
                    
                    if search_images:
                        print(f"[api_visual_generate] 搜索到{len(search_images)}张图片: {search_images}")
                except Exception as e:
                    print(f"[api_visual_generate] 图片搜索失败: {e}")
            
            image_prompt = ""
            if search_images:
                image_prompt = f"""

# 可用图片素材（在动画中合理使用这些图片）
以下是搜索到的相关图片URL，请在动画中使用<img>标签引用：
{chr(10).join([f"- 图片{i+1}: {url}" for i, url in enumerate(search_images)])}

使用示例：<img src="图片URL" style="width:100px;height:auto;" />
"""

            if "勾股定理" in topic:
                anim_prompt = f"""请为数学课程"{topic}"生成一个科学准确的几何演示动画HTML代码。

# 勾股定理可视化（a² + b² = c²）

## 核心几何原理
直角三角形两直角边分别为a和b，斜边为c，则 a² + b² = c²。
动画必须清晰展示：两个直角边上的正方形面积之和等于斜边上的正方形面积。

## 构图要求

### 1. 直角三角形（底部中央）
- 放置在坐标系第一象限
- 直角顶点在原点(0,0)附近
- 底边a沿x轴方向（蓝色）
- 高边b沿y轴方向（红色）
- 斜边c（黑色粗线）
- 直角标记（小正方形或直角符号）

### 2. 三个正方形（必须准确构造）
- **a边上的正方形**：沿x轴，边长=a，蓝色填充，面积标注a²
- **b边上的正方形**：沿y轴，边长=b，红色填充，面积标注b²
- **c边上的正方形**：以斜边c为一条边的正方形，绿色填充，面积标注c²

### 3. c边上正方形的几何构造（非常重要！必须按向量法构造）
设斜边端点为A(x1,y1)和C(x2,y2)：
- 向量 v = (x2-x1, y2-y1)
- 边长 c = sqrt((x2-x1)² + (y2-y1)²)
- 垂直向量 n = (-v.y, v.x) = (y1-y2, x2-x1)
- 正方形四个顶点：A, C, C+n, A+n
- 使用SVG <polygon>或<path>绘制
- **禁止**使用"以中点为中心旋转的轴对齐正方形"的错误方法

## 动画效果（CSS动画，无限循环）
1. **面积闪烁**：三个正方形依次高亮闪烁（a²→b²→c²→a²），配合文字标注
2. **面积累加**：a²和b²的面积块依次"移动"到c²位置，演示a²+b²=c²
3. **勾股公式显示**：底部公式"a² + b² = c²"逐字显示或闪烁高亮
4. **边长标注**：a、b、c三边长依次标注，箭头指示

## 尺寸与布局
- viewBox="0 0 500 320"
- 直角三角形底边a=120px，高b=90px（3:4:5比例），斜边c=150px
- 直角顶点坐标(80, 200)
- 底边终点(200, 200)
- 高边终点(80, 110)
- 斜边端点(80, 110)到(200, 200)
- 文字标注清晰，字体大小12-14px

## 色彩规范
- 底边a及正方形：蓝色 #3498DB
- 高边b及正方形：红色 #E74C3C
- 斜边c及正方形：绿色 #27AE60
- 文字：深灰 #333
- 背景：透明

## 格式规范
1. 输出纯HTML代码，包含在<div class="classroom-animation">中
2. CSS放在<style>标签内，SVG直接内联
3. 动画自动循环（animation-iteration-count: infinite）
4. 背景透明
5. 禁止使用JavaScript
6. SVG必须使用精确的几何计算，不能有视觉误差

## 几何验证检查
1. 三个正方形的面积是否满足a²+b²=c²？
2. c边上的正方形是否以斜边为一条边？（不是旋转的轴对齐正方形）
3. 直角标记是否清晰？
4. 公式"a² + b² = c²"是否准确标注？

直接输出HTML代码，不要其他解释。
""" + image_prompt
            else:
                anim_prompt = f"""请为理科课程主题"{topic}"生成一个具有教学逻辑性的课堂演示动画HTML代码。

# 角色定义
你是一位擅长用SVG+CSS动画演示理科教学逻辑的专业课件动画师。你必须使用内联SVG绘制精细的教学图形，配合CSS动画展示知识点的逻辑关系、过程演变或因果关系。本动画主要服务于理科教学（物理、化学、数学、生物）。

# 核心原则（必须遵守）
- **必须使用SVG内联绘图**来绘制教学图形（分子结构、电路图、几何图形、实验装置、生物细胞等）
- SVG图形必须精细、美观、专业，像教科书插图一样清晰
- **动画必须具有教学逻辑性**，展示因果关系、演变过程、步骤流程或对比关系
- 动画是理科教学演示工具，必须能帮助学生理解"{topic}"的核心概念
- 逻辑清晰：从起点→过程→结果，或从问题→分析→结论的完整逻辑链条
- 背景必须透明，便于嵌入PPT

# SVG绘图要求（非常重要！）
- 所有的教学图形、实验装置、分子结构、几何图形等必须用内联<svg>标签绘制
- SVG必须设置viewBox属性，推荐viewBox="0 0 500 320"
- SVG元素使用stroke和fill设置颜色，stroke-width设置线宽
- 用<circle>画圆、<rect>画矩形、<line>画直线、<path>画曲线、<polygon>/<polyline>画多边形
- 用<text>添加标注文字，设置font-size和fill颜色
- 用<g>标签分组，配合CSS动画让整个组动起来
- 用<defs>和<marker>定义箭头等可复用元素
- SVG图形要专业精细，线条流畅，颜色协调，不要画简陋的火柴人

# 理科SVG绘图示例（主要服务以下学科）

**物理**：用SVG画弹簧+方块（弹簧用<path>的正弦曲线）、斜面+方块、电路图（<rect>电阻+<line>导线+<circle>电池）、光的折射（<line>入射光+折射光+<rect>界面）、磁场线（<path>曲线+箭头）、波动图（<path>正弦波）

**化学**：用SVG画分子结构（<circle>原子+<line>化学键）、烧杯试管（<rect>+<path>液面）、电子云（<circle>不同透明度的圆叠加）、反应方程式（<text>+上下标）、原子结构（<circle>原子核+<ellipse>电子轨道）

**数学**：用SVG画坐标系（<line>坐标轴+<text>刻度）、函数曲线（<path>贝塞尔曲线）、几何图形（<polygon>+标注）、数列变化（<rect>柱状图+动画）、导数/积分示意（<path>曲线+<rect>面积）

**生物**：用SVG画细胞结构（<circle>细胞膜+<circle>细胞核+<rect>线粒体）、DNA双螺旋（<path>两条螺旋线+<line>碱基对）、光合作用（<circle>叶绿体+<path>光能+<text>产物）、细胞分裂（<circle>动态分裂过程）

# 理科教学逻辑动画类型

【类型1：过程演示型】适用于物理实验、化学反应、生物生长等
示例结构：SVG绘制起点状态 → CSS动画过渡 → SVG绘制终态（循环）
例：化学反应H₂+O₂→H₂O、细胞分裂过程、电路通电过程

【类型2：因果关系型】适用于物理原理、化学规律、生物机制等
示例结构：SVG绘制原因/条件 → 动画展示作用机制 → SVG绘制结果/现象
例：力→加速度→速度变化、温度升高→分子运动加快→状态改变

【类型3：对比演示型】适用于概念辨析、正确vs错误、变量对比等
示例结构：SVG左侧画A + SVG右侧画B → 动态高亮差异
例：光合作用vs呼吸作用、串联vs并联电路、酸性vs碱性

【类型4：循环系统型】适用于物质循环、能量循环、生物循环等
示例结构：SVG画各环节 → 依次高亮激活 → 形成循环回路
例：碳循环、水循环、能量流动、细胞呼吸链

【类型5：层级展开型】适用于知识结构、分类体系、公式推导等
示例结构：SVG从中心向外逐层展开
例：生物分类树、数学公式推导步骤、物质分类体系

# 格式规范（必须100%遵守）
1. 输出纯HTML代码，不要markdown代码块包裹
2. 代码必须包含在 <div class="classroom-animation"> 中
3. CSS样式放在 <style> 标签内，SVG图形直接内联
4. 动画必须是自动循环播放的（animation-iteration-count: infinite）
5. 背景色必须设为透明（background: transparent）
6. 整体尺寸限制在 500px 宽 × 320px 高以内
7. 文字使用中文，简短明确（关键词≤6字），字体大小12-16px
8. 可以使用提供的图片素材（通过<img>标签引用），禁止使用其他外部字体、JS库资源
9. 禁止使用JavaScript，仅用CSS animation/transition实现动画
10. SVG图形必须标注教学含义（如受力标注F、加速度标注a等）
11. 逻辑链条必须清晰可见，用动画顺序或位置关系体现因果/过程

# 正例格式锚点

【示例：牛顿第二定律因果链】
<div class="classroom-animation" style="width:500px;height:320px;...">
  <style>
    @keyframes pushForce {{ from {{ transform: translateX(0); }} to {{ transform: translateX(30px); }} }}
    @keyframes moveBlock {{ from {{ transform: translateX(0); }} to {{ transform: translateX(80px); }} }}
    .force-arrow {{ animation: pushForce 2s ease-in-out infinite alternate; }}
    .block {{ animation: moveBlock 2s ease-in-out infinite alternate; }}
  </style>
  <svg viewBox="0 0 500 320" width="500" height="320">
    <!-- 地面 -->
    <line x1="20" y1="220" x2="480" y2="220" stroke="#888" stroke-width="2"/>
    <!-- 方块 -->
    <g class="block">
      <rect x="120" y="170" width="60" height="50" fill="#4A90D9" rx="4"/>
      <text x="150" y="200" text-anchor="middle" fill="white" font-size="14">m</text>
    </g>
    <!-- 力的箭头 -->
    <g class="force-arrow">
      <line x1="60" y1="195" x2="115" y2="195" stroke="#E74C3C" stroke-width="3" marker-end="url(#arrowhead)"/>
      <text x="85" y="185" text-anchor="middle" fill="#E74C3C" font-size="14" font-weight="bold">F</text>
    </g>
    <!-- 箭头定义 -->
    <defs><marker id="arrowhead" markerWidth="10" markerHeight="7" refX="10" refY="3.5" orient="auto"><polygon points="0 0, 10 3.5, 0 7" fill="#E74C3C"/></marker></defs>
    <!-- 公式标注 -->
    <text x="250" y="280" text-anchor="middle" fill="#333" font-size="16" font-weight="bold">F = ma</text>
  </svg>
</div>

【示例：化学反应过程】
<div class="classroom-animation" style="...">
  <style>
    @keyframes bond {{ from {{ opacity:1; }} to {{ opacity:0; }} }}
    @keyframes newBond {{ from {{ opacity:0; }} to {{ opacity:1; }} }}
    .old-bond {{ animation: bond 3s ease-in-out infinite; }}
    .new-bond {{ animation: newBond 3s ease-in-out 1.5s infinite; }}
  </style>
  <svg viewBox="0 0 500 320" width="500" height="320">
    <g class="old-bond">
      <circle cx="150" cy="160" r="25" fill="#E74C3C"/><text x="150" y="165" text-anchor="middle" fill="white" font-size="14">H₂</text>
      <circle cx="350" cy="160" r="25" fill="#3498DB"/><text x="350" y="165" text-anchor="middle" fill="white" font-size="14">O₂</text>
    </g>
    <text x="250" y="165" text-anchor="middle" fill="#333" font-size="20">→</text>
    <g class="new-bond">
      <circle cx="250" cy="160" r="30" fill="#2ECC71"/><text x="250" y="165" text-anchor="middle" fill="white" font-size="14">H₂O</text>
    </g>
    <text x="250" y="270" text-anchor="middle" fill="#333" font-size="14">2H₂+O₂→2H₂O</text>
  </svg>
</div>

# 反例（禁止生成）
- 用CSS div方块画教学图形（× 必须用SVG画精细图形）
- 纯装饰性动画：花瓣飘落、星星闪烁（× 无教学意义）
- 简陋火柴人/粗糙贴图式图形（× SVG要精细专业）
- 纯文字静态展示（× 必须有动画效果）
- 包含JavaScript代码（× 纯CSS动画）
- 背景不透明（× 必须透明）
- 使用未提供的外部图片URL（× 只能使用提供的图片素材）
- 文科类动画（× 本模块专注理科：物理/化学/数学/生物）

# 动画设计检查清单
1. 是否使用了SVG绘制精细教学图形？
2. 这个动画展示了"{topic}"的什么理科核心逻辑？
3. 学生看完能理解哪个知识点？
4. SVG图形是否专业清晰，像教科书插图？
5. 逻辑链条是否清晰（起点→过程→结果）？

直接输出HTML代码，不要其他解释。
""" + image_prompt
            response = client.chat.completions.create(
                model=MODEL,
                messages=[{"role": "user", "content": anim_prompt}],
                temperature=0.7,
                max_tokens=4096
            )
            anim_html = response.choices[0].message.content.strip()
            
            anim_html = re.sub(r'^```\w*\s*', '', anim_html)
            anim_html = re.sub(r'\s*```$', '', anim_html)
            
            if 'classroom-animation' not in anim_html:
                anim_html = f'<div class="classroom-animation" style="width:500px;height:320px;background:transparent;display:flex;align-items:center;justify-content:center;color:#333;font-size:18px;border:2px dashed #ccc;border-radius:12px;">{topic}</div>'
            
            database.save_search_history(user['id'], f"可视化生成-教学动画：{topic}")
            
            return jsonify({
                "success": True,
                "type": "animation",
                "format": "html",
                "html": anim_html
            })
        else:
            return jsonify({"success": False, "error": "未知的生成类型"}), 400
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"[api_visual_generate] 异常: {e}")
        return jsonify({"success": False, "error": f"生成失败: {str(e)}"}), 500

if __name__ == "__main__":
    print("""
============================================
       BeiKe Assistant (Web) - 登录版 2.0
       http://localhost:5024
       Ctrl+C to exit
============================================
""")
    import webbrowser
    threading.Timer(1.5, lambda: webbrowser.open("http://127.0.0.1:5024/")).start()
    app.run(host="0.0.0.0", port=5024, debug=False, threaded=True)
